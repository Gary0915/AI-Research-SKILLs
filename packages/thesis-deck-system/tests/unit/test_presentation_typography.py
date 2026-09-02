"""Closed semantic typography authority tests."""

from __future__ import annotations

from pathlib import Path

import pytest


ROOT = Path(__file__).resolve().parents[4]


def test_typography_profile_has_all_semantic_roles_and_truthful_font_evidence():
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.presentation_typography import build_presentation_typography_profile

    profile = build_presentation_typography_profile(ROOT)
    roles = {item["role"] for item in profile["roles"]}

    assert {"deck_title", "section_title", "slide_title", "body", "caption", "citation", "metric_primary", "table_header", "formula", "footer", "page_number"} <= roles
    assert profile["font_family_fidelity"] == "insufficient_evidence"
    assert profile["theme_binding_valid"] is True
    assert profile["native_render_verified"] is False
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", schema_names=("presentation-typography-profile",))
    assert registry.errors("presentation-typography-profile", profile) == []


def test_typography_governor_resolves_semantic_role_and_rejects_uncontrolled_overrides():
    from thesis_deck_system.presentation_typography import TypographyGovernorError, build_presentation_typography_profile, resolve_typography

    profile = build_presentation_typography_profile(ROOT)
    resolved = resolve_typography(profile, "caption")

    assert resolved["role"] == "caption"
    assert resolved["font_size_pt"] < resolve_typography(profile, "body")["font_size_pt"]
    assert resolve_typography(profile, "slide_title")["font_size_pt"] > resolve_typography(profile, "body")["font_size_pt"]
    with pytest.raises(TypographyGovernorError):
        resolve_typography(profile, "caption", override={"font_family": "uncontrolled"})


def test_typography_governor_preserves_mixed_cjk_latin_and_scientific_text():
    from thesis_deck_system.presentation_typography import build_presentation_typography_profile, validate_editable_text

    profile = build_presentation_typography_profile(ROOT)
    text = "水凝膠 / Hydrogel Δ — 量測結果 / Result (μm)"

    assert validate_editable_text(profile, "caption", text) == text


def test_typography_governor_applies_a_role_to_an_editable_pptx_textbox(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.presentation_typography import apply_typography_to_shape, build_presentation_typography_profile

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 1000000, 300000)
    shape.text = "量測結果 / Result"
    applied = apply_typography_to_shape(shape, build_presentation_typography_profile(ROOT), "caption")

    assert applied["role"] == "caption"
    assert shape.text_frame.paragraphs[0].runs[0].font.size.pt == 12
