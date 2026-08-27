from __future__ import annotations

import copy
import json
from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation

from thesis_deck_system.ledger import Ledger


ROOT = Path(__file__).resolve().parents[4]


def _event(event_type: str, payload: dict):
    ledger = Ledger()
    ledger.append(event_type, payload)
    return ledger


def test_transition_rejects_downstream_experiment_result_as_precursor_even_when_card_is_early():
    from thesis_deck_system.hypothesis import validate_evidence_causal_roles

    ledger = Ledger()
    ledger.append("claim_created", {"claim_id": "C101"})
    ledger.append("claim_created", {"claim_id": "C201"})
    # The card is deliberately early, but its declared origin is still the
    # later discriminating H02 experiment and therefore cannot be historical
    # transition evidence.
    ledger.append("evidence_linked", {
        "evidence_id": "E201", "causal_role": "experiment_result",
        "origin": {"experiment_stage_ref": "ST-EXP201", "layer_ref": "H002"},
    })
    ledger.append("stage_revised", {"stage_id": "ST-RES101", "stage_type": "result", "block_ref": {"block_id": "B101"}})
    ledger.append("decision_recorded", {"decision_id": "D101"})
    ledger.append("hypothesis_transition_recorded", {
        "transition_id": "TR-H001-H002", "previous_hypothesis_claim_ref": "C101",
        "new_hypothesis_claim_ref": "C201", "key_result_refs": ["RES101"],
        "decision_refs": ["D101"], "observation_or_uncertainty_refs": ["E201"],
    })
    ledger.append("stage_revised", {"stage_id": "ST-EXP201", "stage_type": "experiment", "block_ref": {"block_id": "B201"}})

    assert "P2-CAUSAL-TRANSITION-DOWNSTREAM-EVIDENCE" in {
        finding.rule_id for finding in validate_evidence_causal_roles(ledger)
    }


def test_split_resolution_rejects_self_approval_future_evidence_and_unresolved_plan():
    from thesis_deck_system.layout import validate_split_resolution

    plan = {"slide_id": "S-H001-PROBLEM-01", "split_recommendation": True}
    self_approved = {"resolution_type": "external_review_override", "approved_by": "Phase 2 synthetic acceptance review", "approval_artifact": "none"}
    future = {"resolution_type": "automated_fit_exception", "measurement_artifact": "render/future.json", "measurement_cursor": 9, "available_cursor": 8}
    assert validate_split_resolution(plan, self_approved)
    assert validate_split_resolution(plan, future)
    assert validate_split_resolution(plan, None)
    assert validate_split_resolution(plan, {"resolution_type": "split", "continuation_slide_ids": ["S-H001-EXP-01", "S-H001-EXP-02"]}) == []
    assert validate_split_resolution(plan, {"resolution_type": "automated_fit_exception", "measurement_artifact": "render/fit.json", "measurements_pass": True, "measurement_cursor": 8, "available_cursor": 8}) == []


def test_render_pixel_evidence_changes_when_the_same_spec_render_is_blank_or_cropped(tmp_path: Path):
    from thesis_deck_system.qa2 import run_visual_qa_v2

    spec = {
        "slide_id": "S-H001-PROBLEM-01", "semantic_role": "problem_definition",
        "title": {"text": "Problem"}, "content": {"slots": {"previous_finding": "a", "unresolved_conflict": "b", "research_question": "c"}},
        "placement_plan": [
            {"slot": "previous_finding", "left": .8, "top": 1.4, "width": 3, "height": 2, "font_size_pt": 18, "element_role": "evidence"},
            {"slot": "unresolved_conflict", "left": 4, "top": 1.4, "width": 3, "height": 2, "font_size_pt": 18, "element_role": "problem"},
            {"slot": "research_question", "left": 7.2, "top": 1.4, "width": 3, "height": 2, "font_size_pt": 20, "element_role": "question"},
        ],
    }
    good = tmp_path / "good.png"
    image = Image.new("RGB", (1280, 720), "white")
    ImageDraw.Draw(image).rectangle((70, 120, 1200, 650), fill="#4b5563")
    image.save(good)
    blank = tmp_path / "blank.png"
    Image.new("RGB", (1280, 720), "white").save(blank)
    good_report = run_visual_qa_v2([spec], {spec["slide_id"]: good}, expected_size=(1280, 720))
    blank_report = run_visual_qa_v2([spec], {spec["slide_id"]: blank}, expected_size=(1280, 720))
    assert good_report["render_pixel_qa"]["slides"][0]["render_sha256"] != blank_report["render_pixel_qa"]["slides"][0]["render_sha256"]
    assert blank_report["status"] == "fail"


def test_professor_qa_discovers_h003_transition_without_literal_h001_h002_contracts():
    from thesis_deck_system.qa2 import run_professor_qa_v2
    import inspect

    def layer(layer_id: str, previous: str | None, revision: int):
        number = layer_id[-3:]
        return {
            "hypothesis_layer_id": layer_id, "hypothesis_claim_ref": f"C{number}", "problem_ref": f"P{number}",
            "research_block_refs": [f"B{number}"], "fishbone_snapshot_ref": {"fishbone_id": "FB001", "revision": revision},
            "fishbone_focus_refs": [f"FB-{number}"], "result_refs": [f"RES{number}"],
            "layer_discussion_ref": f"DISC-{layer_id}", "layer_summary_ref": f"SUM-{layer_id}",
            "derived_from": None if previous is None else {"previous_layer_ref": previous},
        }

    layers = [layer("H001", None, 1), layer("H002", "H001", 2), layer("H003", "H002", 3)]
    state = {
        "hypothesis_layers": {entry["hypothesis_layer_id"]: entry for entry in layers},
        "claims": {f"C{n}": {} for n in ("001", "002", "003")},
        "stages": {f"ST-RES{n}": {} for n in ("001", "002", "003")},
        "decisions": {"D002": {}}, "evidence": {"E003": {"causal_role": "transition_precursor"}},
        "hypothesis_transitions": {"TR-H002-H003": {
            "from_layer_ref": "H002", "to_layer_ref": "H003", "previous_hypothesis_claim_ref": "C002",
            "new_hypothesis_claim_ref": "C003", "key_result_refs": ["RES002"], "decision_refs": ["D002"],
            "observation_or_uncertainty_refs": ["E003"],
        }},
    }
    slides = []
    for item in layers:
        layer_id = item["hypothesis_layer_id"]
        slides.extend([
            {"semantic_role": "hypothesis_title", "hypothesis_layer_ref": layer_id},
            {"semantic_role": "problem_definition", "hypothesis_layer_ref": layer_id},
            {"semantic_role": "fishbone_locator", "hypothesis_layer_ref": layer_id, "fishbone_snapshot_ref": item["fishbone_snapshot_ref"], "fishbone_focus_refs": item["fishbone_focus_refs"]},
            {"semantic_role": "result_single", "hypothesis_layer_ref": layer_id},
            {"semantic_role": "layer_integrated_discussion", "hypothesis_layer_ref": layer_id},
            {"semantic_role": "layer_summary_decision", "hypothesis_layer_ref": layer_id},
        ])
    slides.append({"semantic_role": "hypothesis_transition", "hypothesis_layer_ref": "H002", "object_ref": "TR-H002-H003"})
    report = run_professor_qa_v2({"profile_id": "PROF-SYNTH-001", "version": "1.0.0"}, {"layers": layers, "slides": slides, "state": state, "previous_commitments": [{"owner": "Gary", "target_window": {"due": "2026-09-10"}, "dependency_refs": [], "parallelizable": True, "status": "planned"}]})
    assert "PROF-TRANSITION-PROVENANCE" in report["executed_checks"]
    assert "H003" in report["evidence"]["PROF-HISTORY-REACHABLE"]["detail"]
    source = inspect.getsource(run_professor_qa_v2)
    assert '"H001"' not in source and '"H002"' not in source and '"TR-H001-H002"' not in source


def test_persisted_ledger_rebuilds_layout_without_seed_fixture_reads(tmp_path: Path, monkeypatch):
    import shutil
    import yaml
    import thesis_deck_system.phase2_build as phase2_build
    from thesis_deck_system.phase2_build import build_phase2, rebuild_specs_and_layouts_from_ledger

    build_phase2(output_root=tmp_path)
    ledger = Ledger.load(tmp_path / "ledger-events.json")
    canonical_specs = json.loads((tmp_path / "slide-specs.json").read_text(encoding="utf-8"))
    canonical_plans = json.loads((tmp_path / "layout-plans.json").read_text(encoding="utf-8"))
    profile = json.loads((tmp_path / "template-profile.json").read_text(encoding="utf-8"))
    mutated_root = tmp_path / "mutated-seed"
    shutil.copytree(ROOT / "thesis-deck-system/examples/synthetic-project/phase2", mutated_root)
    fixture_path = mutated_root / "fixture.yaml"
    fixture = yaml.safe_load(fixture_path.read_text(encoding="utf-8"))
    fixture["hypothesis_layers"][0]["title"] = "MUTATED UNCOMMITTED SCIENTIFIC STORY"
    fixture["hypothesis_layers"][1]["experiment_refs"] = ["EXP999"]
    fixture["hypothesis_transition"]["rationale"] = "MUTATED FUTURE LEAK"
    fixture_path.write_text(yaml.safe_dump(fixture, allow_unicode=True, sort_keys=False), encoding="utf-8")
    monkeypatch.setattr(phase2_build, "FIXTURE_ROOT", mutated_root)
    specs, plans = rebuild_specs_and_layouts_from_ledger(ledger, tmp_path, profile)
    assert specs == canonical_specs
    assert plans == canonical_plans


def test_structural_audit_rejects_three_slot_problem_when_only_one_physical_slot_survives(tmp_path: Path):
    from thesis_deck_system.phase2_build import build_phase2
    from thesis_deck_system.pptx import audit_pptx

    build_phase2(output_root=tmp_path)
    specs = json.loads((tmp_path / "slide-specs.json").read_text(encoding="utf-8"))
    target_index = next(index for index, spec in enumerate(specs) if spec["semantic_role"] == "problem_definition")
    deck = tmp_path / "acceptance-deck.pptx"
    presentation = Presentation(deck)
    target_slide = presentation.slides[-len(specs) + target_index]
    for shape in list(target_slide.shapes):
        if shape.name.startswith("tds-slot:") and shape.name != "tds-slot:research_question":
            target_slide.shapes._spTree.remove(shape._element)
    broken = tmp_path / "problem-one-slot.pptx"
    presentation.save(broken)
    profile = json.loads((tmp_path / "template-profile.json").read_text(encoding="utf-8"))
    audit = audit_pptx(broken, tmp_path / "synthetic-template.pptx", profile, specs)
    target = next(item for item in audit["generated_slides"] if item["slide_spec_id"] == specs[target_index]["slide_id"])
    assert target["governed_geometry_match"] is False
    assert target["governed_slot_matches"]["previous_finding"] is False
    assert target["governed_slot_matches"]["unresolved_conflict"] is False
