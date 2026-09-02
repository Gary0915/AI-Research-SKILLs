"""Production application tests for the bounded Planner v1 layer."""

from pathlib import Path

import pytest


ROOT = Path(__file__).resolve().parents[4]


def test_application_cases_select_semantic_recipes_and_reject_single_visual_for_dense_validation():
    from thesis_deck_system.presentation_planner_application import build_planner_application

    application = build_planner_application(ROOT)
    by_case = {item["case_id"]: item for item in application["cases"]}
    assert all("eligible_body_families" not in item for item in application["scenario_inputs"])
    assert "BCF-HARDWARE-DESIGN-PROCEDURE" in by_case["CASE-A-EXPERIMENT"]["eligible_body_families"]
    assert "BCF-PHYSICAL-VALIDATION-MATRIX" in by_case["CASE-B-PHYSICAL"]["eligible_body_families"]
    assert "BCF-REAL-RESULT-VALIDATION" in by_case["CASE-C-RESULT"]["eligible_body_families"]
    assert "BCF-LITERATURE-VISUAL-MATRIX" in by_case["CASE-E-LITERATURE"]["eligible_body_families"]
    assert "BCF-TECHNOLOGY-COMPARISON" in by_case["CASE-F-COMPARISON"]["eligible_body_families"]
    assert "BCF-PRINCIPLE-EQUIPMENT-SPLIT" in by_case["CASE-G-PRINCIPLE"]["eligible_body_families"]
    assert "single_visual_capacity" in by_case["CASE-B-PHYSICAL"]["rejection_reasons"]


def test_candidates_are_structurally_distinct_and_tie_break_is_deterministic():
    from thesis_deck_system.presentation_planner_application import build_planner_application

    first = build_planner_application(ROOT)
    second = build_planner_application(ROOT)
    first_j = next(item for item in first["cases"] if item["case_id"] == "CASE-J-DIVERSITY")
    second_j = next(item for item in second["cases"] if item["case_id"] == "CASE-J-DIVERSITY")
    assert len(first_j["candidates"]) == 2
    assert len({item["structure_fingerprint"] for item in first_j["candidates"]}) == 2
    assert first_j["selected_decision"]["selected_candidate_id"] == second_j["selected_decision"]["selected_candidate_id"]
    assert first["candidate_difference_audit"]["fake_candidate_variant_count"] == 0


def test_candidate_breadth_exposes_at_least_four_genuinely_valid_structural_alternatives():
    from thesis_deck_system.presentation_planner_application import build_planner_application

    application = build_planner_application(ROOT)
    multi = [case for case in application["cases"] if len(case["candidates"]) >= 2]

    assert len(multi) >= 4
    assert all(len({candidate["structure_fingerprint"] for candidate in case["candidates"]}) == len(case["candidates"]) for case in multi)


def test_all_closed_body_families_have_distinct_normalized_physical_recipes():
    from thesis_deck_system.presentation_planner_application import build_body_composition_recipe_registry

    recipes = build_body_composition_recipe_registry()
    assert len(recipes) == 10
    assert all(recipe["normalized_content_bounds"] == {"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0} for recipe in recipes)
    assert len({recipe["geometry_hash"] for recipe in recipes}) == 10
    by_family = {recipe["body_family_id"]: recipe for recipe in recipes}
    assert {region["presentation_role"] for region in by_family["BCF-HARDWARE-DESIGN-PROCEDURE"]["regions"]} >= {"primary_visual", "procedure", "go_criterion"}
    assert sum(region["presentation_role"] == "citation_strip" for region in by_family["BCF-LITERATURE-VISUAL-MATRIX"]["regions"]) >= 2


def test_candidate_physical_plans_bind_each_candidate_to_its_recipe_geometry():
    from thesis_deck_system.presentation_planner_application import build_physical_composition_plans, build_planner_application

    application = build_planner_application(ROOT)
    plans = build_physical_composition_plans(application)
    assert len(plans) == application["metrics"]["candidate_count"]
    assert all(plan["required_role_coverage_status"] == "pass" for plan in plans)
    by_candidate = {plan["candidate_id"]: plan for plan in plans}
    for record in application["candidate_difference_audit"]["records"]:
        assert by_candidate[record["candidate_a"]]["geometry_hash"] != by_candidate[record["candidate_b"]]["geometry_hash"]


def test_each_physical_plan_binds_a_closed_body_style_recipe_and_spacing_scale():
    from thesis_deck_system.presentation_planner_application import build_physical_composition_plans, build_planner_application

    plans = build_physical_composition_plans(build_planner_application(ROOT))

    assert all(plan["body_style_recipe_id"].startswith("BSR-") for plan in plans)
    assert all(plan["spacing_scale_id"] == "PSS-001" for plan in plans)
    assert all(plan["style_authority_status"] == "body_only_no_shell_override" for plan in plans)


def test_physical_plans_bind_the_professor_shell_profile_without_claiming_body_measurement():
    from thesis_deck_system.presentation_planner_application import build_physical_composition_plans, build_planner_application
    from thesis_deck_system.professor_shell import build_professor_shell_profile

    shell = build_professor_shell_profile(ROOT)
    plans = build_physical_composition_plans(build_planner_application(ROOT), shell_profile=shell)

    assert all(plan["shell_profile_id"] == shell["shell_profile_id"] for plan in plans)
    assert all(plan["content_bounds"] == shell["body_content_safe_region"]["geometry_inches"] for plan in plans)
    assert shell["body_content_safe_region"]["evidence_level"] == "synthetic_system_owned"


def test_reverse_physical_audit_recovers_deterministic_region_identity_from_pptx(tmp_path: Path):
    from thesis_deck_system.presentation_planner_application import reverse_audit_physical_composition, write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    audit = reverse_audit_physical_composition(outputs["review_pptx"])
    assert audit["missing_required_region_count"] == 0
    assert audit["out_of_content_bounds_count"] == 0
    assert audit["hard_overlap_violation_count"] == 0
    assert all(item["planner_shape_count"] > 0 for item in audit["slides"])


def test_review_deck_uses_governed_typography_and_native_synthetic_visuals_not_wireframe_labels(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    presentation = Presentation(outputs["review_pptx"])
    texts = [shape.text for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame]
    effect_shapes = [shape for slide in presentation.slides for shape in slide.shapes if shape.name.startswith("PPAFX::")]
    titles = [shape for slide in presentation.slides for shape in slide.shapes if shape.name.startswith("tds-title:")]

    assert not any(text.startswith("SYNTHETIC_NON_EVIDENCE") for text in texts if text)
    assert effect_shapes
    assert titles and all(shape.text_frame.paragraphs[0].runs[0].font.size.pt == 26 for shape in titles)


def test_review_deck_persists_shell_typography_style_and_occupancy_audits(tmp_path: Path):
    import json

    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    shell = json.loads(outputs["shell_reverse_audit"].read_text(encoding="utf-8"))
    typography = json.loads(outputs["typography_reverse_audit"].read_text(encoding="utf-8"))
    style = json.loads(outputs["style_reverse_audit"].read_text(encoding="utf-8"))
    occupancy = json.loads(outputs["occupancy_audit"].read_text(encoding="utf-8"))

    assert shell["aggregate_status"] == "pass"
    assert shell["title_safe_region_violation_count"] == 0
    assert typography["typography_role_mismatch_count"] == 0
    assert style["uncontrolled_style_override_count"] == 0
    assert occupancy["recipe_registry_family_coverage"] == "10/10"
    assert occupancy["review_deck_materialized_family_coverage"] == "10/10"
    assert occupancy["fully_occupied_golden_family_coverage"] == "10/10"


def test_reverse_audit_compares_each_review_slide_to_the_full_physical_recipe(tmp_path: Path):
    from thesis_deck_system.presentation_planner_application import (
        reverse_audit_physical_composition,
        write_planner_application_artifacts,
    )

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    audit = reverse_audit_physical_composition(outputs["review_pptx"], outputs["physical_plans"])

    assert audit["missing_required_region_count"] == 0
    assert audit["physical_recipe_identity_mismatch"] == 0
    assert audit["selected_candidate_materialization_mismatch"] == 0
    assert audit["review_slide_mapping_failure_count"] == 0
    assert all("region_id" in region for slide in audit["slides"] for region in slide["regions"])


def test_incremental_physical_application_reuses_historical_artifacts_and_inserts_new_children():
    from thesis_deck_system.presentation_planner_application import build_incremental_physical_application_audit

    audit = build_incremental_physical_application_audit(ROOT)

    assert audit["historical_reused"] == 20
    assert audit["new_planned_slides"] == 2
    assert audit["new_physical_slides"] == 2
    assert audit["historical_visual_migration_count"] == 0
    assert audit["semantic_insertion_status"] == "pass"


def test_historical_lock_and_stale_reviewer_selection_never_migrate_or_change_science():
    from thesis_deck_system.presentation_planner_application import PlannerApplicationError, build_planner_application, apply_reviewer_selection

    application = build_planner_application(ROOT)
    historical = next(item for item in application["cases"] if item["case_id"] == "CASE-H-HISTORICAL")
    assert historical["selected_decision"]["selection_mode"] == "historical_reuse"
    candidate = historical["candidates"][0]
    review = {"slide_id": historical["slide_id"], "candidate_id": candidate["candidate_id"], "dependency_hash": "0" * 64, "selection_origin": "reviewer_selection", "layout_locked": True}
    with pytest.raises(PlannerApplicationError):
        apply_reviewer_selection(historical, review)
    assert "scientific_fields" not in review


def test_review_overlay_applies_only_same_dependency_candidate_and_bounded_presentation_adjustment():
    from thesis_deck_system.presentation_planner_application import (
        PlannerApplicationError,
        apply_presentation_review_overlay,
        build_physical_composition_plans,
        build_planner_application,
    )

    application = build_planner_application(ROOT)
    physical_plans = build_physical_composition_plans(application)
    case = next(item for item in application["cases"] if item["case_id"] == "CASE-A-EXPERIMENT")
    selected = case["selected_decision"]["selected_candidate_id"]
    plan = next(item for item in physical_plans if item["candidate_id"] == selected)
    overlay = {
        "overlay_id": "PRO-001",
        "slide_id": case["slide_id"],
        "dependency_hash": case["dependency_hash"],
        "selected_candidate_id": selected,
        "layout_locked": True,
        "meeting_visibility": "visible",
        "bounded_region_adjustments": [{"region_id": plan["content_item_assignments"][0]["region_id"], "delta_x": 0.01, "delta_y": 0.0}],
        "review_note": "synthetic presentation adjustment",
        "review_origin": "reviewer_selection",
    }

    applied = apply_presentation_review_overlay(case, plan, overlay)

    assert applied["status"] == "applied"
    assert applied["selection_mode"] == "reviewer_selection"
    assert applied["layout_locked"] is True
    assert applied["adjusted_physical_plan_hash"] != plan["physical_composition_hash"]
    with pytest.raises(PlannerApplicationError):
        apply_presentation_review_overlay(case, plan, overlay | {"scientific_value": "forbidden"})
    stale = apply_presentation_review_overlay(case, plan, overlay | {"dependency_hash": "0" * 64})
    assert stale["status"] == "stale"
    assert stale["selection_applied"] is False


def test_review_overlay_rejects_a_tampered_persisted_overlay_hash():
    from thesis_deck_system.presentation_planner_application import (
        PlannerApplicationError,
        apply_presentation_review_overlay,
        build_physical_composition_plans,
        build_planner_application,
    )

    application = build_planner_application(ROOT)
    case = next(item for item in application["cases"] if item["case_id"] == "CASE-A-EXPERIMENT")
    plan = next(
        item for item in build_physical_composition_plans(application)
        if item["candidate_id"] == case["selected_decision"]["selected_candidate_id"]
    )
    overlay = {
        "overlay_id": "PRO-002", "slide_id": case["slide_id"], "dependency_hash": case["dependency_hash"],
        "selected_candidate_id": plan["candidate_id"], "layout_locked": True, "meeting_visibility": "visible",
        "bounded_region_adjustments": [], "review_note": "synthetic review", "review_origin": "reviewer_selection",
        "overlay_sha256": "0" * 64,
    }

    with pytest.raises(PlannerApplicationError, match="hash"):
        apply_presentation_review_overlay(case, plan, overlay)


def test_physical_planner_contract_schemas_are_closed_and_registered():
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.presentation_planner_application import (
        build_body_composition_recipe_registry,
        build_physical_composition_plans,
        build_planner_application,
    )

    registry = SchemaRegistry(
        ROOT / "thesis-deck-system" / "schemas",
        schema_names=(
            "body-composition-recipe-registry",
            "physical-composition-plans",
            "presentation-review-overlay",
            "planner-physical-reverse-audit",
            "incremental-planner-physical-application-audit",
        ),
    )
    application = build_planner_application(ROOT)
    registry.validate("body-composition-recipe-registry", {"schema_version": "2.0.0", "registry_id": "BCR-REG-001", "recipes": build_body_composition_recipe_registry()})
    registry.validate("physical-composition-plans", {"schema_version": "2.0.0", "planner_version": "2.0.0", "records": build_physical_composition_plans(application)})


def test_new_physical_planner_schemas_do_not_leave_untyped_object_array_items():
    import json

    names = (
        "body-composition-recipe-registry",
        "physical-composition-plans",
        "presentation-review-overlay",
        "planner-physical-reverse-audit",
        "incremental-planner-physical-application-audit",
    )
    for name in names:
        payload = json.loads((ROOT / "thesis-deck-system" / "schemas" / f"{name}.schema.json").read_text(encoding="utf-8"))
        assert '"items":{"type":"object"}' not in json.dumps(payload, separators=(",", ":"))


def test_application_materializes_review_only_pptx_and_incremental_scenario(tmp_path: Path):
    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts
    from thesis_deck_system.contracts import SchemaRegistry

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    assert outputs["review_pptx"].is_file()
    assert outputs["acceptance"].is_file()
    assert outputs["physical_plans"].is_file()
    assert outputs["reverse_audit"].is_file()
    assert outputs["review_overlays"].is_file()
    import json
    acceptance = json.loads(outputs["acceptance"].read_text(encoding="utf-8"))
    physical_plans = json.loads(outputs["physical_plans"].read_text(encoding="utf-8"))
    assert len(physical_plans["records"]) > json.loads(outputs["review_json"].read_text(encoding="utf-8"))["metrics"]["candidate_count"]
    from pptx import Presentation
    review_application = json.loads(outputs["review_json"].read_text(encoding="utf-8"))
    assert len(Presentation(outputs["review_pptx"]).slides) == len(physical_plans["records"])
    assert acceptance["structural_audit"]["missing_required_region_count"] == 0
    assert acceptance["structural_audit"]["hard_capacity_violation_count"] == 0
    assert acceptance["structural_audit"]["selected_candidate_materialization_mismatch"] == 0
    assert acceptance["incremental_scenario"]["historical_migrations"] == 0
    assert acceptance["incremental_scenario"]["new_physical_slides"] == 2
    reverse_audit = json.loads(outputs["reverse_audit"].read_text(encoding="utf-8"))
    assert reverse_audit["aggregate_status"] == "pass"
    assert reverse_audit["review_slide_mapping_failure_count"] == 0
    registry = SchemaRegistry(ROOT / "thesis-deck-system/schemas", include_cp5hi=True)
    registry.validate("planner-application-acceptance", acceptance)
    registry.validate("composition-review-selections", json.loads(outputs["selections"].read_text(encoding="utf-8")))


def test_generated_physical_planner_artifacts_close_all_registered_contracts(tmp_path: Path):
    """The persisted review artifacts, not just in-memory plans, are schema-closed."""
    import json

    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    registry = SchemaRegistry(
        ROOT / "thesis-deck-system/schemas",
        schema_names=(
            "body-composition-recipe-registry",
            "physical-composition-plans",
            "presentation-review-overlay",
            "planner-physical-reverse-audit",
            "incremental-planner-physical-application-audit",
        ),
    )
    for artifact_name, schema_name in (
        ("recipe_registry", "body-composition-recipe-registry"),
        ("physical_plans", "physical-composition-plans"),
        ("review_overlays", "presentation-review-overlay"),
        ("reverse_audit", "planner-physical-reverse-audit"),
        ("incremental", "incremental-planner-physical-application-audit"),
    ):
        registry.validate(schema_name, json.loads(outputs[artifact_name].read_text(encoding="utf-8")))


def test_physical_realization_qa_is_derived_from_generated_artifacts(tmp_path: Path):
    """Physical planner QA must retain execution facts instead of self-certifying PASS."""
    import json

    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    qa = json.loads(outputs["physical_qa"].read_text(encoding="utf-8"))

    assert qa["aggregate_status"] == "pass"
    assert qa["recipe_coverage"]["recipe_count"] == 10
    assert qa["recipe_coverage"]["physical_family_coverage"] == "8/10"
    assert qa["recipe_coverage"]["recipe_registry_family_coverage"] == "10/10"
    assert qa["recipe_coverage"]["review_deck_materialized_family_coverage"] == "10/10"
    assert qa["reverse_physical_audit"]["missing_required_region_count"] == 0
    assert qa["review_overlay"]["stale_review_applied_count"] == 0
    assert qa["incremental_physical_application"]["historical_reused"] == 20
    assert qa["body_reference"]["shell_override_count"] == 0


def test_writer_returns_only_schema_valid_physical_planner_artifacts(tmp_path: Path):
    from thesis_deck_system.presentation_planner_application import (
        validate_planner_physical_realization_artifacts,
        write_planner_application_artifacts,
    )

    outputs = write_planner_application_artifacts(ROOT, tmp_path)

    assert validate_planner_physical_realization_artifacts(ROOT, outputs) == 0


def test_physical_realization_candidate_state_hashes_the_exact_execution_surface():
    from thesis_deck_system.presentation_planner_application import physical_realization_candidate_state

    state = physical_realization_candidate_state(ROOT)

    assert state["component_count"] == len(state["component_hashes"])
    assert state["component_count"] >= 12
    assert "thesis-deck-system/schemas/planner-visual-calibration-qa.schema.json" in state["component_hashes"]
    assert "thesis-deck-system/artifacts/phase3/planner-visual-calibration-qa.json" in state["component_hashes"]
    assert len(state["candidate_state_sha256"]) == 64
    assert all(len(value) == 64 for value in state["component_hashes"].values())


def test_candidate_component_hash_normalizes_text_line_endings_but_not_binary_content(tmp_path: Path):
    from thesis_deck_system.presentation_planner_application import _candidate_component_hash

    lf = tmp_path / "fixture.py"
    crlf = tmp_path / "fixture-copy.py"
    binary = tmp_path / "fixture.pptx"
    lf.write_bytes(b"x = 1\ny = 2\n")
    crlf.write_bytes(b"x = 1\r\ny = 2\r\n")
    binary.write_bytes(b"x\r\ny\r\n")

    assert _candidate_component_hash(lf) == _candidate_component_hash(crlf)
    assert _candidate_component_hash(binary) != _candidate_component_hash(lf)


def test_writer_persists_the_candidate_state_used_for_physical_review_evidence(tmp_path: Path):
    import json

    from thesis_deck_system.presentation_planner_application import write_planner_application_artifacts

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    state = json.loads(outputs["candidate_state"].read_text(encoding="utf-8"))

    assert state["candidate_id"] == "PPA-PHYSICAL-CANDIDATE-001"
    assert state["component_count"] == len(state["component_hashes"])


def test_persisted_candidate_state_is_not_self_referential(tmp_path: Path):
    import json

    from thesis_deck_system.presentation_planner_application import (
        physical_realization_candidate_state,
        write_planner_application_artifacts,
    )

    outputs = write_planner_application_artifacts(ROOT, tmp_path)
    persisted = json.loads(outputs["candidate_state"].read_text(encoding="utf-8"))

    assert persisted["candidate_state_sha256"] == physical_realization_candidate_state(ROOT)["candidate_state_sha256"]
    assert "thesis-deck-system/artifacts/phase3/planner-physical-realization-candidate-state.json" not in persisted["component_hashes"]
