"""Sanitized professor-shell authority and native-template reverse audit.

This module deliberately consumes only committed sanitized resolver artifacts.
It never opens a professor template or resolves a private source.
"""

from __future__ import annotations

from hashlib import sha256
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


class ProfessorShellError(ValueError):
    """Raised when a shell profile or its physical projection is not truthful."""


_INCHES_PER_EMU = 914400
_SHELL_ARTIFACT = Path("thesis-deck-system/artifacts/phase3/professor-template-resolved.json")
_ALLOWED_EVIDENCE = frozenset({
    "measured_sanitized", "repository_observed", "synthetic_system_owned",
    "insufficient_evidence", "not_applicable",
})


def _canonical(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _hash(value: Any) -> str:
    return sha256(_canonical(value).encode("utf-8")).hexdigest()


def _as_inches(geometry: dict[str, float], canvas: dict[str, float]) -> dict[str, float]:
    return {
        "left": round(geometry["x"] * canvas["width_inches"], 6),
        "top": round(geometry["y"] * canvas["height_inches"], 6),
        "width": round(geometry["w"] * canvas["width_inches"], 6),
        "height": round(geometry["h"] * canvas["height_inches"], 6),
    }


def _region(
    geometry: dict[str, float] | None,
    canvas: dict[str, float],
    *,
    evidence_level: str,
    fidelity_status: str,
    source_ids: list[str],
) -> dict[str, Any]:
    if evidence_level not in _ALLOWED_EVIDENCE:
        raise ProfessorShellError("uncontrolled shell evidence level")
    return {
        "geometry_inches": _as_inches(geometry, canvas) if geometry is not None else None,
        "evidence_level": evidence_level,
        "fidelity_status": fidelity_status,
        "source_evidence_ids": source_ids,
    }


def build_professor_shell_profile(root: Path) -> dict[str, Any]:
    """Build physical-shell authority from the committed CP3 resolver output."""
    root = Path(root).resolve()
    source = json.loads((root / _SHELL_ARTIFACT).read_text(encoding="utf-8"))
    tokens = {item["token_family"]: item for item in source["shell_tokens"]}
    canvas_token = tokens.get("canvas", {}).get("value", {})
    if canvas_token.get("kind") != "canvas":
        raise ProfessorShellError("sanitized resolver has no canvas authority")
    canvas = {
        "width_inches": float(canvas_token["width"]),
        "height_inches": float(canvas_token["height"]),
        "aspect_ratio": "16:9",
        "evidence_level": "measured_sanitized",
        "fidelity_status": "partial",
        "source_evidence_ids": list(tokens["canvas"].get("supporting_ids", [])),
    }

    def measured_region(token_name: str) -> dict[str, Any]:
        token = tokens[token_name]
        value = token["value"]
        if value.get("kind") != "geometry":
            return _region(None, canvas, evidence_level="insufficient_evidence", fidelity_status="insufficient_evidence", source_ids=[])
        return _region(
            value["geometry"], canvas,
            evidence_level="measured_sanitized",
            fidelity_status="partial",
            source_ids=list(token.get("supporting_ids", [])),
        )

    # CP3 intentionally found no defensible professor-derived safe bounds.  This
    # system-owned region is an explicit assembly fallback, never a measurement.
    # This is deliberately an assembly-owned fallback because CP3 found no
    # defensible professor-derived body safe area.  It starts below the
    # measured title region, so a physical planner cannot silently overlap
    # title and body while calling the resulting bounds professor evidence.
    body_fallback = {"left": 0.7, "top": 2.0, "width": 11.85, "height": 4.3}
    body_region = {
        "geometry_inches": body_fallback,
        "evidence_level": "synthetic_system_owned",
        "fidelity_status": "insufficient_evidence",
        "source_evidence_ids": [],
    }
    topology = source.get("content_master_layout_topology", [])
    profile = {
        "schema_version": "1.0.0",
        "shell_profile_id": "PSP-001",
        "version": "1.0.0",
        "canvas": canvas,
        "master_identity": {
            "master_ids": sorted({item["target_id"] for item in topology}),
            "evidence_level": "measured_sanitized",
            "fidelity_status": "insufficient_evidence",
        },
        "layout_identities": {
            "layout_ids": sorted(item["source_id"] for item in topology),
            "topology": topology,
            "evidence_level": "measured_sanitized",
            "fidelity_status": "insufficient_evidence",
        },
        "title_safe_region": measured_region("content_title"),
        "body_content_safe_region": body_region,
        "footer_region": measured_region("footer"),
        "header_region": _region(None, canvas, evidence_level="insufficient_evidence", fidelity_status="insufficient_evidence", source_ids=[]),
        "page_number_region": measured_region("page_number"),
        "recurring_marker_regions": [],
        "background_treatment": {"role": "background", "evidence_level": "repository_observed", "fidelity_status": "partial"},
        "shell_color_roles": ["background", "foreground", "muted", "border"],
        "shell_font_role_bindings": {"status": "insufficient_evidence", "roles": []},
        "content_exclusion_regions": [],
        "source_evidence_ids": sorted({source["profile_id"], *[item["source_id"] for item in topology], *[item["target_id"] for item in topology]}),
        "measurement_evidence_level": "partial",
        "fidelity_status": "insufficient_evidence",
    }
    profile["shell_profile_sha256"] = _hash({key: value for key, value in profile.items() if key != "shell_profile_sha256"})
    return profile


def apply_professor_shell_profile(presentation: Presentation, profile: dict[str, Any]) -> None:
    """Apply only physically realizable, profile-declared shell features."""
    if profile.get("shell_profile_id") != "PSP-001":
        raise ProfessorShellError("unknown professor shell profile")
    canvas = profile["canvas"]
    presentation.slide_width = int(round(canvas["width_inches"] * _INCHES_PER_EMU))
    presentation.slide_height = int(round(canvas["height_inches"] * _INCHES_PER_EMU))
    title = profile["title_safe_region"]["geometry_inches"]
    body = profile["body_content_safe_region"]["geometry_inches"]
    if title is None or body is None:
        raise ProfessorShellError("physical template requires title and body regions")
    for layout in presentation.slide_layouts:
        for placeholder in layout.placeholders:
            kind = str(placeholder.placeholder_format.type).lower()
            if "title" in kind:
                placeholder.left = int(round(title["left"] * _INCHES_PER_EMU))
                placeholder.top = int(round(title["top"] * _INCHES_PER_EMU))
                placeholder.width = int(round(title["width"] * _INCHES_PER_EMU))
                placeholder.height = int(round(title["height"] * _INCHES_PER_EMU))
            elif "body" in kind or "object" in kind:
                placeholder.left = int(round(body["left"] * _INCHES_PER_EMU))
                placeholder.top = int(round(body["top"] * _INCHES_PER_EMU))
                placeholder.width = int(round(body["width"] * _INCHES_PER_EMU))
                placeholder.height = int(round(body["height"] * _INCHES_PER_EMU))


def _matches(actual: dict[str, float], expected: dict[str, float], tolerance: float = 0.003) -> bool:
    return all(abs(actual[name] - expected[name]) <= tolerance for name in expected)


def audit_professor_shell_template(path: Path, profile: dict[str, Any]) -> dict[str, Any]:
    """Reverse-audit only physical claims the profile says are reproducible."""
    presentation = Presentation(path)
    canvas = profile["canvas"]
    canvas_mismatch = int(
        presentation.slide_width != int(round(canvas["width_inches"] * _INCHES_PER_EMU))
        or presentation.slide_height != int(round(canvas["height_inches"] * _INCHES_PER_EMU))
    )
    expected_title = profile["title_safe_region"]["geometry_inches"]
    expected_body = profile["body_content_safe_region"]["geometry_inches"]
    title_matches = body_matches = 0
    title_seen = body_seen = 0
    for layout in presentation.slide_layouts:
        for placeholder in layout.placeholders:
            kind = str(placeholder.placeholder_format.type).lower()
            geometry = {
                "left": placeholder.left / _INCHES_PER_EMU,
                "top": placeholder.top / _INCHES_PER_EMU,
                "width": placeholder.width / _INCHES_PER_EMU,
                "height": placeholder.height / _INCHES_PER_EMU,
            }
            if "title" in kind:
                title_seen += 1
                title_matches += int(_matches(geometry, expected_title))
            elif "body" in kind or "object" in kind:
                body_seen += 1
                body_matches += int(_matches(geometry, expected_body))
    mismatches = canvas_mismatch + int(title_seen == 0 or title_seen != title_matches) + int(body_seen == 0 or body_seen != body_matches)
    # Master and layout identity reconstruction is explicitly not claimed: CP3
    # measured the source topology, while this is a fresh sanitized template.
    unsupported_claimed = 0
    return {
        "audit_id": "PSP-TEMPLATE-AUDIT-001",
        "shell_profile_id": profile["shell_profile_id"],
        "slide_size": {"width_emu": presentation.slide_width, "height_emu": presentation.slide_height},
        "master_count": len(presentation.slide_masters),
        "layout_count": len(presentation.slide_layouts),
        "title_placeholder_count": title_seen,
        "body_placeholder_count": body_seen,
        "canvas_mismatch_count": canvas_mismatch,
        "title_safe_region_mismatch_count": int(title_seen == 0 or title_seen != title_matches),
        "body_safe_region_mismatch_count": int(body_seen == 0 or body_seen != body_matches),
        "shell_profile_to_pptx_mismatch_count": mismatches,
        "unsupported_claimed_shell_feature_count": unsupported_claimed,
        "aggregate_status": "pass" if mismatches == 0 and unsupported_claimed == 0 else "fail",
    }


def write_professor_shell_artifacts(root: Path, destination: Path | None = None) -> dict[str, Path]:
    """Persist the profile, a fresh shell-only template, and its reverse audit."""
    from .template import create_sanitized_native_template

    root = Path(root).resolve()
    destination = Path(destination or root / "thesis-deck-system" / "artifacts" / "phase3")
    destination.mkdir(parents=True, exist_ok=True)
    profile = build_professor_shell_profile(root)
    profile_path = destination / "professor-shell-profile.json"
    template_path = destination / "professor-shell-template.pptx"
    audit_path = destination / "professor-shell-template-audit.json"
    create_sanitized_native_template(template_path, shell_profile=profile)
    audit = audit_professor_shell_template(template_path, profile)
    if audit["aggregate_status"] != "pass":
        raise ProfessorShellError("fresh professor shell template did not satisfy its profile")
    profile_path.write_text(json.dumps(profile, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    audit_path.write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return {"profile": profile_path, "template": template_path, "audit": audit_path}
