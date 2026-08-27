"""Synthetic native PowerPoint template creation and OpenXML profiling."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches


NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def create_synthetic_template(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Synthetic Native Template"
    slide.placeholders[1].text = "Redistributable test fixture — not laboratory data"
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Representative content layout"
    slide.placeholders[1].text = "Native title/content placeholders"
    prs.save(path)
    return path


def profile_template(path: Path, output_path: Path) -> dict:
    prs = Presentation(path)
    masters = [
        {
            "master_index": index,
            "master_path": master.part.partname.lstrip("/"),
            "name": f"master-{index}",
            "relationship_ids": sorted(master.part.rels),
        }
        for index, master in enumerate(prs.slide_masters)
    ]
    layouts = []
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for placeholder in layout.placeholders:
            placeholders.append(
                {
                    "type": str(placeholder.placeholder_format.type).split(" (")[0].lower(),
                    "idx": placeholder.placeholder_format.idx,
                    "geometry": {
                        "left": placeholder.left,
                        "top": placeholder.top,
                        "width": placeholder.width,
                        "height": placeholder.height,
                    },
                }
            )
        layouts.append(
            {
                "layout_index": index,
                "layout_path": layout.part.partname.lstrip("/"),
                "master_path": layout.slide_master.part.partname.lstrip("/"),
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        theme_fonts = []
        theme_colors = []
        theme_names = sorted(n for n in names if n.startswith("ppt/theme/") and n.endswith(".xml"))
        if theme_names:
            root = ET.fromstring(archive.read(theme_names[0]))
            theme_fonts = [node.attrib.get("typeface") for node in root.findall(".//a:latin", NS) if node.attrib.get("typeface")]
            theme_colors = [node.tag.rsplit("}", 1)[-1] for node in root.findall(".//a:clrScheme/*", NS)]
    content_layout = layouts[1]
    role = {
        "layout_index": content_layout["layout_index"],
        "layout_path": content_layout["layout_path"],
        "master_path": content_layout["master_path"],
        "required_placeholders": ["title", "body"],
    }
    profile = {"schema_version": "1.0.0", "profile_id": "TP-SYNTH-001", "version": "1.0.0", "source_path": path.as_posix(), "source_sha256": hashlib.sha256(path.read_bytes()).hexdigest(), "slide_size": {"width_emu": prs.slide_width, "height_emu": prs.slide_height, "aspect_ratio": "16:9"}, "masters": masters, "layouts": layouts, "theme": {"major_fonts": theme_fonts[:1], "minor_fonts": theme_fonts[1:2], "colors": theme_colors}, "semantic_roles": {"photo_observation": dict(role), "hero_plot_discussion": dict(role)}, "created_at": "2026-08-27T00:00:00Z"}
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(profile, indent=2), encoding="utf-8")
    return profile
