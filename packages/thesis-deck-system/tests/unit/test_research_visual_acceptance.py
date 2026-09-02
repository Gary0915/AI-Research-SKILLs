"""Real-research visual-acceptance fixtures are source-closed and non-inventive."""

from pathlib import Path

from copy import deepcopy


ROOT = Path(__file__).resolve().parents[4]


def test_real_research_fixture_pack_is_source_closed_chinese_primary_and_non_inventive():
    from thesis_deck_system.research_visual_acceptance import build_real_research_fixture_pack

    pack = build_real_research_fixture_pack(ROOT)

    assert pack["fixture_pack_id"] == "RRVFP-001"
    assert len(pack["fixtures"]) >= 14
    assert {item["fixture_id"] for item in pack["fixtures"]} >= {
        "R01", "R02", "R03", "R04", "R05", "R06", "R07", "R08", "R09", "R10", "R11", "R12", "R13", "R14",
    }
    assert all(item["canonical_source_refs"] for item in pack["fixtures"])
    assert all(len(item["dependency_hash"]) == 64 for item in pack["fixtures"])
    assert all(item["traditional_chinese_primary"] is True for item in pack["fixtures"])
    assert pack["invented_scientific_claim_count"] == 0
    assert pack["invented_measured_value_count"] == 0


def test_real_result_fixture_marks_layout_only_plot_as_synthetic_non_evidence():
    from thesis_deck_system.research_visual_acceptance import build_real_research_fixture_pack

    result = next(item for item in build_real_research_fixture_pack(ROOT)["fixtures"] if item["fixture_id"] == "R11")

    assert result["scientific_evidence_status"] == "synthetic_non_evidence"
    assert "不代表實驗結果" in result["visible_text"]
    assert result["source_only_role_indicators"] == ["visual_layout_fixture"]


def test_fixture_writer_persists_schema_valid_source_closed_artifact(tmp_path: Path):
    import json

    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.research_visual_acceptance import write_real_research_fixture_pack

    path = write_real_research_fixture_pack(ROOT, tmp_path)
    payload = json.loads(path.read_text(encoding="utf-8"))
    SchemaRegistry(ROOT / "thesis-deck-system/schemas", schema_names=("real-research-visual-fixture-pack",)).validate(
        "real-research-visual-fixture-pack", payload
    )
    assert payload["aggregate_status"] == "source_closed_review_fixture_pack"


def test_visual_acceptance_profile_and_manifest_preserve_pending_human_choices():
    from thesis_deck_system.research_visual_acceptance import (
        build_professor_visual_review_manifest,
        build_research_presentation_visual_acceptance_profile,
    )

    profile = build_research_presentation_visual_acceptance_profile(ROOT)
    manifest = build_professor_visual_review_manifest(ROOT)

    assert profile["human_visual_acceptance"] == "not_reviewed"
    assert profile["traditional_chinese_primary_language"] == "pass"
    assert profile["main_content_minimum_font_pt"] == 16
    assert profile["rule_counts"]["source_observed"] > 0
    assert profile["rule_counts"]["system_calibrated"] > 0
    assert len(manifest["cases"]) == 14
    assert sum(len(case["candidates"]) for case in manifest["cases"]) == 21
    assert all(case["human_selection"] is None and case["human_status"] == "pending" for case in manifest["cases"])


def test_profile_and_manifest_writers_emit_closed_contracts(tmp_path: Path):
    import json

    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.research_visual_acceptance import write_visual_acceptance_review_artifacts

    outputs = write_visual_acceptance_review_artifacts(ROOT, tmp_path)
    registry = SchemaRegistry(ROOT / "thesis-deck-system/schemas", schema_names=(
        "research-presentation-visual-acceptance-profile", "professor-visual-review-manifest",
    ))
    registry.validate("research-presentation-visual-acceptance-profile", json.loads(outputs["profile"].read_text(encoding="utf-8")))
    registry.validate("professor-visual-review-manifest", json.loads(outputs["manifest"].read_text(encoding="utf-8")))


def test_real_research_review_application_materializes_all_fixtures_with_meaningful_candidates():
    from thesis_deck_system.presentation_planner_application import build_physical_composition_plans
    from thesis_deck_system.professor_shell import build_professor_shell_profile
    from thesis_deck_system.research_visual_acceptance import build_real_research_review_application

    application = build_real_research_review_application(ROOT)

    assert application["logical_fixture_count"] == 14
    assert len(application["cases"]) == 14
    assert application["real_candidate_slide_count"] >= 20
    assert application["multi_candidate_fixture_count"] >= 6
    assert application["fake_candidate_variant_count"] == 0
    assert all(case["body_source_fit_status"] == "pass" for case in application["cases"])
    assert all(candidate["content_items"] for case in application["cases"] for candidate in case["candidates"])
    assert any("不代表實驗結果" in item["visible_text"] for case in application["cases"] for candidate in case["candidates"] for item in candidate["content_items"])

    physical = build_physical_composition_plans(application, shell_profile=build_professor_shell_profile(ROOT))
    assert len(physical) == application["real_candidate_slide_count"]
    assert all(plan["required_role_coverage_status"] == "pass" for plan in physical)
    assert any("接觸壓力" in region.get("visible_text", "") for plan in physical for region in plan["physical_regions"])


def test_real_research_review_writer_creates_real_first_pptx_and_pending_human_mapping(tmp_path: Path):
    import json

    from pptx import Presentation
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.research_visual_acceptance import write_real_research_visual_review_artifacts

    outputs = write_real_research_visual_review_artifacts(ROOT, tmp_path)
    presentation = Presentation(outputs["review_pptx"])
    manifest = json.loads(outputs["review_manifest"].read_text(encoding="utf-8"))
    texts = [shape.text for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame]

    assert outputs["review_pptx"].is_file()
    assert len(presentation.slides) == 31  # 21 real candidates + ten golden appendices.
    assert any("研究缺口" in text for text in texts)
    assert any("不代表實驗結果" in text for text in texts)
    assert all(candidate["pptx_slide_index"] >= 1 for case in manifest["cases"] for candidate in case["candidates"])
    assert len(manifest["cases"]) == 14
    assert sum(len(case["candidates"]) for case in manifest["cases"]) == 21
    assert all(case["human_selection"] is None and case["human_status"] == "pending" for case in manifest["cases"])
    assert all(shape.text_frame.paragraphs[0].runs[0].font.size.pt == 30 for slide in list(presentation.slides)[:21] for shape in slide.shapes if shape.name.startswith("tds-title:"))
    SchemaRegistry(ROOT / "thesis-deck-system/schemas", schema_names=("professor-visual-review-manifest", "physical-composition-plans")).validate(
        "professor-visual-review-manifest", manifest
    )
    assert outputs["visual_qa"].is_file()
    assert outputs["render_discovery"].is_file()


def test_real_research_visual_qa_is_derived_from_materialized_pptx(tmp_path: Path):
    from thesis_deck_system.research_visual_acceptance import (
        build_real_research_visual_qa,
        write_real_research_visual_review_artifacts,
    )

    outputs = write_real_research_visual_review_artifacts(ROOT, tmp_path)
    qa = build_real_research_visual_qa(outputs["review_pptx"], outputs["application"])

    assert qa["aggregate_status"] == "pass"
    assert qa["traditional_chinese_primary_language"] == "pass"
    assert qa["main_content_below_16pt_count"] == 0
    assert qa["title_typography_violation_count"] == 0
    assert qa["hard_overlap_violation_count"] == 0
    assert qa["dashboard_style_violation_count"] == 0
    assert qa["fixed_four_box_footer_count"] == 0
    assert qa["shell_override_count"] == 0
    assert qa["scientific_truth_override_count"] == 0


def test_real_research_review_outputs_have_closed_machine_contracts(tmp_path: Path):
    import json

    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.research_visual_acceptance import write_real_research_visual_review_artifacts

    outputs = write_real_research_visual_review_artifacts(ROOT, tmp_path)
    registry = SchemaRegistry(ROOT / "thesis-deck-system/schemas", schema_names=(
        "real-research-visual-review-application", "real-research-visual-qa", "render-capability-discovery",
    ))
    for key, schema_name in (("application", "real-research-visual-review-application"), ("visual_qa", "real-research-visual-qa"), ("render_discovery", "render-capability-discovery")):
        registry.validate(schema_name, json.loads(outputs[key].read_text(encoding="utf-8")))


def test_candidate_bound_visual_acceptance_schemas_are_recursively_closed_and_fail_closed():
    """The source-ref and planner-score contracts admit no unknown members."""
    import json

    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.phase3_final_visual_composition import build_candidate_schema_closure_inventory

    registry = SchemaRegistry(ROOT / "thesis-deck-system/schemas", schema_names=(
        "real-research-visual-review-application",
        "composition-selection-audit",
        "real-research-visual-acceptance-schema-closure-inventory",
    ))
    application = json.loads((ROOT / "thesis-deck-system/artifacts/phase3/real-research-visual-review-application.json").read_text(encoding="utf-8"))
    planner_audit = json.loads((ROOT / "thesis-deck-system/artifacts/phase3/composition-selection-audit.json").read_text(encoding="utf-8"))

    # All fourteen real source-reference variants remain valid.
    assert registry.errors("real-research-visual-review-application", application) == []
    assert registry.errors("composition-selection-audit", planner_audit) == []

    unknown_source_ref = deepcopy(application)
    unknown_source_ref["cases"][0]["fixture"]["canonical_source_refs"]["unexpected_source_ref"] = "X001"
    assert registry.errors("real-research-visual-review-application", unknown_source_ref)

    unknown_nested_score_member = deepcopy(planner_audit)
    unknown_nested_score_member["selections"][0]["candidate_component_scores"][0]["score"]["body_recurrence_evidence"]["unexpected_member"] = "forbidden"
    assert registry.errors("composition-selection-audit", unknown_nested_score_member)

    inventory = build_candidate_schema_closure_inventory(ROOT)
    assert inventory["checked_schema_count"] >= 26
    assert inventory["open_node_count"] == 0
    registry.validate("real-research-visual-acceptance-schema-closure-inventory", inventory)


def test_render_capability_discovery_is_honest_and_non_rendering():
    from thesis_deck_system.research_visual_acceptance import discover_review_render_capability

    discovery = discover_review_render_capability()

    assert discovery["render_attempt_count"] == 0
    assert discovery["renderer_install_attempt_count"] == 0
    assert discovery["candidate_preview_status"] in {"renderer_available_not_run", "blocked_environment"}
