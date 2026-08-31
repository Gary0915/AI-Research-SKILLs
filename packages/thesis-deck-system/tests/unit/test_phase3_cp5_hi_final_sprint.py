"""CP5-H/I final-sprint contracts, starting with H0 preflight."""
from __future__ import annotations

from pathlib import Path


ROOT = Path(__file__).resolve().parents[4]


def test_h0_backend_uniqueness_and_environment_probe_are_execution_owned(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import (
        ScientificSvgNativeCompiler,
        audit_single_pptx_backend,
        probe_native_environment,
    )

    audit = audit_single_pptx_backend(ROOT)
    environment = probe_native_environment()

    assert audit["status"] == "pass"
    assert audit["public_pptx_backends"] == ["PythonPptxAssembler"]
    assert audit["compiler_pptx_writer_methods"] == []
    assert audit["template_scientific_slide_assembly_bypasses"] == []
    assert audit["bypass_count"] == 0
    assert environment["python_pptx_version"]
    assert environment["native_powerpoint_status"] in {"available", "blocked_environment"}
    assert environment["host_pptx_renderer_status"] in {"available", "blocked_environment"}
    assert not hasattr(ScientificSvgNativeCompiler, "save_pptx")
    assert not hasattr(ScientificSvgNativeCompiler, "export_pptx")


def test_h0_artifacts_are_schema_closed_and_bind_real_audit_facts(tmp_path: Path):
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_h0_artifacts

    artifacts = build_h0_artifacts(ROOT, tmp_path)
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)

    assert artifacts["backend_uniqueness"]["status"] == "pass"
    assert artifacts["execution_evidence"]["private_alias_resolution_attempts"] == 0
    assert registry.errors("cp5-hi-backend-uniqueness-audit", artifacts["backend_uniqueness"]) == []
    assert registry.errors("cp5-hi-execution-evidence", artifacts["execution_evidence"]) == []


def test_h1_compiler_requires_reverified_handle_and_returns_deterministic_native_plan():
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.phase3_cp5_hi_final_sprint import NativeCompilationError, ScientificSvgNativeCompiler
    from thesis_deck_system.phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure

    produced = build_representative_director_output(ROOT, "mechanism")
    handle = reverify_approved_figure(produced["manifest"], produced["critic"]["report"], produced["critic"]["approval"], ROOT)
    compiler = ScientificSvgNativeCompiler()

    plan = compiler.compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.2, "width": 8.0, "height": 4.5})
    repeat = compiler.compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.2, "width": 8.0, "height": 4.5})

    assert plan["schema_version"] == "1.0.0"
    assert plan["figure_id"] == handle.figure_id
    assert plan["plan_sha256"] == repeat["plan_sha256"]
    assert plan["objects"]
    assert all(item["outcome"] in {"DRAWINGML_EMITTED", "SVG_VECTOR_FALLBACK", "BLOCKED_UNSUPPORTED", "BLOCKED_UNKNOWN_MAPPING"} for item in plan["objects"])
    assert any(item["text"] for item in plan["objects"] if item["shape_kind"] == "text")
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)
    assert registry.errors("native-figure-compilation-plan", plan) == []
    with __import__("pytest").raises(NativeCompilationError):
        compiler.compile({}, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.2, "width": 8.0, "height": 4.5})


def test_h1_artifacts_cover_all_registered_features_and_approved_figure_representatives(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_h1_artifacts

    result = build_h1_artifacts(ROOT, tmp_path)

    assert result["mapping_manifest"]["unmapped_feature_count"] == 0
    assert result["mapping_manifest"]["feature_count"] >= 30
    assert len(result["plans"]) == 8
    assert all(plan["approved_figure"]["manifest_id"] for plan in result["plans"])
    assert (tmp_path / "native-figure-compilation-plans.json").exists()


def test_h2_existing_assembler_consumes_native_plan_as_named_editable_shapes(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.phase3_cp5_hi_final_sprint import ScientificSvgNativeCompiler
    from thesis_deck_system.phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure
    from thesis_deck_system.pptx import PythonPptxAssembler
    from thesis_deck_system.template import create_synthetic_template

    produced = build_representative_director_output(ROOT, "mechanism")
    handle = reverify_approved_figure(produced["manifest"], produced["critic"]["report"], produced["critic"]["approval"], ROOT)
    plan = ScientificSvgNativeCompiler().compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.3, "width": 8.0, "height": 4.5})
    template = create_synthetic_template(tmp_path / "template.pptx")
    presentation = Presentation(template)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    facts = PythonPptxAssembler().add_compiled_figure(slide, handle, plan)

    assert facts["native_object_count"] > 0
    assert facts["fallback_object_count"] >= 0
    assert all(shape.name.startswith(f"tds-fig:{handle.figure_id}/") for shape in slide.shapes if shape.name.startswith("tds-fig:"))
    assert any(shape.has_text_frame and shape.text for shape in slide.shapes)


def test_h2_benchmark_deck_uses_sole_assembler_and_audits_native_identity(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_h2_native_vector_benchmark

    result = build_h2_native_vector_benchmark(ROOT, tmp_path)

    assert result["benchmark"]["backend"] == "PythonPptxAssembler"
    assert result["benchmark"]["figure_count"] == 8
    assert result["audit"]["orphan_parts"] == []
    assert result["audit"]["has_editable_text"] is True
    assert Path(result["benchmark"]["pptx_path"]).exists()


def test_i0_reconstructs_a_fresh_sanitized_template_with_closed_part_manifest(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_i0_sanitized_native_template

    result = build_i0_sanitized_native_template(ROOT, tmp_path)

    assert Path(result["template_path"]).exists()
    presentation = Presentation(result["template_path"])
    assert presentation.slide_width > 0
    assert len(presentation.slides) == 0
    assert result["template_profile"]["fresh_lineage_status"] == "pass"
    assert result["template_profile"]["safe_content_bounds"]["status"] == "insufficient_evidence"
    assert {"formal_cover", "content_academic", "fishbone", "comparison_result", "summary_decision"} <= set(result["template_profile"]["semantic_roles"])
    assert result["reconstruction_manifest"]["unclassified_part_count"] == 0
    assert result["reconstruction_manifest"]["forbidden_part_count"] == 0


def test_i1_builds_fresh_twenty_slide_h001_h002_acceptance_deck(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_i0_sanitized_native_template, build_i1_acceptance_deck

    build_i0_sanitized_native_template(ROOT, tmp_path)
    result = build_i1_acceptance_deck(ROOT, tmp_path)

    assert Path(result["acceptance_deck_path"]).exists()
    assert len(Presentation(result["acceptance_deck_path"]).slides) == 20
    assert result["deck_manifest"]["source_slide_count"] == 19
    assert result["deck_manifest"]["source_slide_mapping_count"] == 19
    assert result["deck_manifest"]["hypothesis_layer_order"] == ["H001", "H002"]
    assert result["deck_manifest"]["h003_slide_count"] == 0
    assert result["deck_manifest"]["governed_figure_bypass_count"] == 0


def test_i2_computes_independent_release_gates_without_promoting_blocked_evidence(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import (
        build_i0_sanitized_native_template,
        build_i1_acceptance_deck,
        build_i2_release_qa,
    )

    build_i0_sanitized_native_template(ROOT, tmp_path)
    build_i1_acceptance_deck(ROOT, tmp_path)
    result = build_i2_release_qa(ROOT, tmp_path)

    by_id = {item["gate_id"]: item for item in result["release_gates"]["gates"]}
    assert len(by_id) == 16
    assert by_id["RG-07"]["status"] == "pass"
    assert by_id["RG-10"]["status"] in {"blocked_environment", "not_run"}
    assert by_id["RG-12"]["status"] == "blocked_visual_review"
    assert by_id["RG-13"]["status"] == "blocked_environment"
    assert result["release_gates"]["production_release_status"] != "pass"
    assert result["release_gates"]["production_group_meeting_ready"] is False


def test_i2_accepts_a_relative_repository_root_for_artifact_rebuild(tmp_path: Path, monkeypatch):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import (
        build_i0_sanitized_native_template,
        build_i1_acceptance_deck,
        build_i2_release_qa,
    )

    build_i0_sanitized_native_template(ROOT, tmp_path)
    build_i1_acceptance_deck(ROOT, tmp_path)
    monkeypatch.chdir(ROOT)
    assert build_i2_release_qa(Path("."), tmp_path)["release_gates"]["release_id"] == "CP5-I2-RELEASE-FACTS-001"


def test_i2_release_and_package_artifacts_validate_against_closed_cp5hi_schemas(tmp_path: Path):
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.phase3_cp5_hi_final_sprint import build_i0_sanitized_native_template, build_i1_acceptance_deck, build_i2_release_qa

    build_i0_sanitized_native_template(ROOT, tmp_path)
    build_i1_acceptance_deck(ROOT, tmp_path)
    result = build_i2_release_qa(ROOT, tmp_path)
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)

    assert registry.errors("cp5-hi-release-gates", result["release_gates"]) == []
    assert registry.errors("cp5-hi-package-manifest", result["package_manifest"]) == []


def test_hi_candidate_state_hash_binds_h_i_sources_tests_schemas_and_artifacts():
    from thesis_deck_system.phase3_cp5_hi_final_sprint import compute_hi_candidate_state

    state = compute_hi_candidate_state(ROOT)

    assert state["component_count"] >= 12
    assert "packages/thesis-deck-system/src/thesis_deck_system/phase3_cp5_hi_final_sprint.py" in state["component_hashes"]
    assert "packages/thesis-deck-system/tests/unit/test_phase3_cp5_hi_final_sprint.py" in state["component_hashes"]
    assert "thesis-deck-system/schemas/cp5-hi-release-gates.schema.json" in state["component_hashes"]
    assert len(state["candidate_state_sha256"]) == 64


def test_hi_cross_gate_acceptance_derives_all_required_h_i_facts(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import (
        build_i0_sanitized_native_template,
        build_i1_acceptance_deck,
        build_i2_release_qa,
        build_hi_cross_gate_acceptance,
    )

    build_i0_sanitized_native_template(ROOT, tmp_path)
    build_i1_acceptance_deck(ROOT, tmp_path)
    build_i2_release_qa(ROOT, tmp_path)
    evidence = build_hi_cross_gate_acceptance(ROOT, tmp_path)

    assert evidence["status"] == "pass"
    assert evidence["check_count"] >= 17
    assert evidence["private_access_counters"] == {"private_alias_resolution_attempts": 0, "private_source_open_attempts": 0, "private_render_attempts": 0}


def test_hi_candidate_hash_normalizes_checkout_line_endings_without_normalizing_binary_inputs():
    from thesis_deck_system.phase3_cp5_hi_final_sprint import _candidate_component_digest

    assert _candidate_component_digest("example.py", b"a\r\nb\r\n") == _candidate_component_digest("example.py", b"a\nb\n")
    assert _candidate_component_digest("example.pptx", b"a\r\nb\r\n") != _candidate_component_digest("example.pptx", b"a\nb\n")
