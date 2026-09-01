"""Final visual-composition closure acceptance tests."""

from pathlib import Path

import pytest


ROOT = Path(__file__).resolve().parents[4]


@pytest.fixture(scope="module")
def final_composition_build(tmp_path_factory: pytest.TempPathFactory):
    """One immutable expensive build shared by read-only composition audits."""
    from thesis_deck_system.phase3_final_visual_composition import build_final_visual_composition

    return build_final_visual_composition(ROOT, tmp_path_factory.mktemp("final-composition"))


def test_final_projection_corrects_results_deduplicates_and_hides_backend_fields():
    from thesis_deck_system.phase3_final_visual_composition import build_final_projection

    projection = build_final_projection(ROOT)
    by_slide = {item["slide_id"]: item for item in projection["slides"]}

    assert len(projection["slides"]) == 19
    assert projection["h003_slide_count"] == 0
    assert by_slide["S-H001-RESULT-SINGLE-08"]["result"]["metric"]["name"] == "mean conductivity increase"
    assert by_slide["S-H001-RESULT-SINGLE-08"]["result"]["metric"]["uncertainty"] == 5
    assert by_slide["S-H001-RESULT-SINGLE-09"]["result"]["metric"]["uncertainty"] == 6
    assert by_slide["S-H002-EXPERIMENT-DESIGN-06"]["result"]["metric"]["value"] == 38
    assert by_slide["S-H002-EXPERIMENT-DESIGN-06"]["result"]["secondary_metric"]["status"] == "qualitative_supported"
    summary = by_slide["S-H002-LAYER-SUMMARY-DECISION-09"]
    assert summary["deduplicated_field_count"] >= 3
    assert all("{'" not in value for value in summary["visible_text"])
    assert all("Metric｜" not in value for item in projection["slides"] for value in item["visible_text"])


def test_final_composition_plan_has_diverse_layouts_and_primary_governed_figures():
    from thesis_deck_system.phase3_final_visual_composition import build_final_projection, build_final_composition_plan

    plan = build_final_composition_plan(ROOT, build_final_projection(ROOT))

    assert len(plan["slides"]) == 20
    assert plan["h003_slide_count"] == 0
    assert len(plan["layout_role_distribution"]) >= 5
    figure_slides = [item for item in plan["slides"] if item["governed_figure"]]
    assert {item["semantic_stage"] for item in figure_slides} >= {"fishbone_locator", "experiment_design", "result_single", "result_comparison"}
    assert all(item["primary_visual_region"]["width"] > item["secondary_text_region"]["width"] for item in figure_slides)


def test_final_composed_deck_is_fresh_and_has_semantic_layout_and_figure_audits(final_composition_build):
    from pptx import Presentation
    result = final_composition_build
    deck = Presentation(result["deck_path"])

    assert len(deck.slides) == 20
    assert result["semantic_audit"]["aggregate_status"] == "pass"
    assert result["layout_audit"]["aggregate_status"] == "pass"
    assert result["figure_audit"]["aggregate_status"] == "pass"
    assert result["semantic_audit"]["visible_raw_backend_field_count"] == 0
    assert result["figure_audit"]["governed_figure_placement_count"] >= 7
    assert any(shape.name.startswith("tds-fig:") for slide in deck.slides for shape in slide.shapes)
    assert (result["deck_path"].parent / "final-visual-composition-candidate-state.json").exists()
    assert result["template_lineage"]["template_unchanged"] is True


def test_final_audits_preserve_per_slide_provenance_and_real_fishbone_revision_bindings(final_composition_build):
    result = final_composition_build
    audit_slides = result["semantic_audit"]["slides"]
    required = {
        "source_slide_spec_id", "source_cursor", "hypothesis_layer", "semantic_stage",
        "canonical_archetype_id", "selected_pptx_layout_id", "title_region",
        "primary_visual_region", "secondary_text_region", "visible_presentation_fields",
        "notes_only_fields", "deduplicated_or_suppressed_fields", "safe_bounds_status",
        "text_occupancy", "figure_occupancy",
    }
    assert len(audit_slides) == 20
    assert all(required <= set(slide) for slide in audit_slides)
    result_trace = next(item for item in audit_slides if item["slide_id"] == "S-H001-RESULT-SINGLE-08")
    assert result_trace["result_trace"]["materialized_result_artifact"].endswith("materialized-h02.json")
    assert result_trace["result_trace"]["presentation_fields"] == result_trace["visible_presentation_fields"]
    fishbone = [item for item in result["figure_audit"]["placements"] if item["route"] == "fishbone"]
    assert {item["fishbone_revision_ref"] for item in fishbone} == {"FB001 rev1", "FB001 rev2"}
    assert all(item["binding_kind"] == "explicit_svg_fallback" for item in fishbone)


def test_final_release_qa_is_execution_derived_and_preserves_blocked_truth(tmp_path: Path):
    from thesis_deck_system.phase3_final_visual_composition import (
        build_final_visual_composition,
        finalize_final_visual_composition_release,
    )

    build_final_visual_composition(ROOT, tmp_path)
    release = finalize_final_visual_composition_release(
        ROOT,
        tmp_path,
        candidate_state={"candidate_state_sha256": "a" * 64, "component_count": 1},
        privacy_evidence={"aggregate_status": "pass", "repository_findings": 0, "staged_findings": 0,
                          "approved_legacy_exceptions": 1,
                          "private_alias_resolution_attempts": 0, "private_source_open_attempts": 0,
                          "private_render_attempts": 0},
        render_evidence={"status": "blocked_environment", "rendered_slide_count": 0},
    )

    assert len(release["gates"]) == 16
    assert release["acceptance_deck_build_status"] == "pass"
    assert release["production_group_meeting_ready"] is False
    assert release["private_alias_resolution_attempts"] == 0
    assert (tmp_path / "final-visual-composition-release-qa.json").exists()


def test_final_candidate_state_binds_composition_source_fixture_and_approved_figure_dependencies():
    from thesis_deck_system.phase3_final_visual_composition import compute_final_visual_composition_candidate_state

    candidate = compute_final_visual_composition_candidate_state(ROOT)

    assert candidate["component_count"] >= 12
    assert len(candidate["candidate_state_sha256"]) == 64
    assert "packages/thesis-deck-system/src/thesis_deck_system/phase3_final_visual_composition.py" in candidate["component_hashes"]
    assert "thesis-deck-system/artifacts/phase2/materialized-h02.json" in candidate["component_hashes"]
    assert "thesis-deck-system/artifacts/phase2/fishbone/FB001-rev1.svg" in candidate["component_hashes"]
