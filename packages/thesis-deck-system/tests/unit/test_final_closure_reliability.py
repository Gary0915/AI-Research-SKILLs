"""Bounded final-closure reliability contracts.

These tests deliberately exercise only synthetic packages and values.  They
prove that a raw PPTX candidate remains visible to the scanner until a
candidate-bound generated-artifact attestation adjudicates it.
"""

from __future__ import annotations

import copy
import json
from pathlib import Path
import subprocess
import zipfile

import pytest


ROOT = Path(__file__).resolve().parents[4]


def _synthetic_pptx(tmp_path: Path) -> Path:
    from thesis_deck_system.template import create_synthetic_template

    return create_synthetic_template(tmp_path / "generated.pptx")


def test_generated_pptx_attestation_is_bound_to_exact_staged_bytes_not_working_tree(tmp_path: Path):
    """FEC-01: a staged package cannot be certified after working-tree drift."""
    from thesis_deck_system.final_closure_reliability import (
        GeneratedArtifactAdjudicationError,
        GeneratedArtifactAdjudicator,
    )
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    producer = tmp_path / "producer.py"
    producer.write_text("# synthetic producer\n", encoding="utf-8")
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    subprocess.run(["git", "add", "generated.pptx", "producer.py"], cwd=tmp_path, check=True)
    adjudicator = GeneratedArtifactAdjudicator(
        root=tmp_path,
        candidate_state_hash="a" * 64,
        approved_producers={"synthetic-template-builder": producer},
        generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")},
        privacy_scanner=RepositoryPrivacyScanner(),
    )

    record = adjudicator.attest_staged_generated_pptx(
        package,
        artifact_class="sanitized_native_template",
        producer_id="synthetic-template-builder",
        declared_input_paths=[],
        execution_id="FC-FEC-01-TEST",
    )
    assert record["working_tree_matches_staged"] is True
    assert record["staged_bytes_sha256"] == record["artifact_sha256"]

    package.write_bytes(package.read_bytes() + b"working-tree-drift")
    with pytest.raises(GeneratedArtifactAdjudicationError):
        adjudicator.adjudicate(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder")


def test_generated_source_closure_and_media_lineage_are_closed_and_hash_bound(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import (
        GeneratedArtifactAdjudicationError,
        build_generated_source_closure,
        build_package_media_lineage,
    )

    source = tmp_path / "source.png"
    source.write_bytes(b"synthetic-media")
    package = tmp_path / "generated.pptx"
    with zipfile.ZipFile(package, "w") as archive:
        archive.writestr("ppt/media/image1.png", source.read_bytes())

    closure = build_generated_source_closure(
        tmp_path,
        [
            {
                "input_id": "SRC-001",
                "repository_relative_path": "source.png",
                "input_class": "canonical_asset",
                "producer_id": "synthetic-producer",
                "input_role": "media_source",
                "source_kind": "repository_canonical",
                "privacy_status": "sanitized",
            }
        ],
    )
    assert closure["unresolved_input_count"] == 0
    media = build_package_media_lineage(package, closure["input_records"], {"ppt/media/image1.png": "SRC-001"})
    assert media["unresolved_media_part_count"] == 0
    assert media["media_lineage_records"][0]["media_sha256"] == media["media_lineage_records"][0]["source_sha256"]
    with pytest.raises(GeneratedArtifactAdjudicationError):
        build_package_media_lineage(package, closure["input_records"], {})


def test_repository_declares_lf_checkout_for_exact_svg_media_lineage():
    """FEC-02 must be invariant to a developer's global autocrlf setting."""
    policy = (ROOT / ".gitattributes").read_text(encoding="utf-8")
    assert "thesis-deck-system/artifacts/phase2/observation/observation_visual.svg text eol=lf" in policy


def test_generated_source_and_media_lineage_evidence_are_schema_closed(tmp_path: Path):
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.final_closure_reliability import build_generated_source_closure, build_package_media_lineage

    source = tmp_path / "source.png"
    source.write_bytes(b"synthetic-media")
    package = tmp_path / "generated.pptx"
    with zipfile.ZipFile(package, "w") as archive:
        archive.writestr("ppt/media/image1.png", source.read_bytes())
    closure = build_generated_source_closure(tmp_path, [{
        "input_id": "SRC-001", "repository_relative_path": "source.png", "input_class": "canonical_asset",
        "producer_id": "synthetic-producer", "input_role": "media_source", "source_kind": "repository_canonical", "privacy_status": "sanitized",
    }])
    media = build_package_media_lineage(package, closure["input_records"], {"ppt/media/image1.png": "SRC-001"})
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)
    assert registry.errors("generated-pptx-source-closure", closure) == []
    assert registry.errors("generated-pptx-media-lineage", media) == []
    assert registry.errors("generated-pptx-source-closure", closure | {"unexpected": True})
    assert registry.errors("generated-pptx-media-lineage", media | {"unexpected": True})


def test_staged_attestation_seals_source_and_media_closure_identity(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import (
        GeneratedArtifactAdjudicationError,
        GeneratedArtifactAdjudicator,
        build_generated_source_closure,
        build_package_media_lineage,
    )
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    source = tmp_path / "source.png"
    source.write_bytes(b"synthetic-media")
    package = tmp_path / "generated.pptx"
    with zipfile.ZipFile(package, "w") as archive:
        archive.writestr("ppt/media/image1.png", source.read_bytes())
    producer = tmp_path / "producer.py"
    producer.write_text("# synthetic producer\n", encoding="utf-8")
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    subprocess.run(["git", "add", "generated.pptx", "source.png", "producer.py"], cwd=tmp_path, check=True)
    closure = build_generated_source_closure(tmp_path, [{
        "input_id": "SRC-001", "repository_relative_path": "source.png", "input_class": "canonical_asset",
        "producer_id": "synthetic-producer", "input_role": "media_source", "source_kind": "repository_canonical", "privacy_status": "sanitized",
    }])
    media = build_package_media_lineage(package, closure["input_records"], {"ppt/media/image1.png": "SRC-001"})
    adjudicator = GeneratedArtifactAdjudicator(
        root=tmp_path, candidate_state_hash="a" * 64,
        approved_producers={"synthetic-template-builder": producer},
        generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")},
        privacy_scanner=RepositoryPrivacyScanner(),
    )
    record = adjudicator.attest_staged_generated_pptx(
        package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder",
        declared_input_paths=[source], execution_id="FC-FEC-02-TEST", source_closure=closure, media_lineage=media,
    )
    assert record["source_closure_sha256"] == closure["source_closure_sha256"]
    assert record["media_lineage_sha256"] == media["media_lineage_sha256"]
    record["media_lineage_sha256"] = "b" * 64
    with pytest.raises(GeneratedArtifactAdjudicationError):
        adjudicator.adjudicate_record(record)


def test_final_figure_slots_bind_slide_specific_canonical_scientific_inputs():
    from thesis_deck_system.phase3_final_visual_composition import (
        build_final_projection,
        build_slide_scientific_figure_bindings,
    )

    bindings = build_slide_scientific_figure_bindings(ROOT, build_final_projection(ROOT))
    by_slide = {item["slide_id"]: item for item in bindings}
    result_bindings = [by_slide[key] for key in ("S-H001-RESULT-SINGLE-08", "S-H001-RESULT-SINGLE-09", "S-H002-EXPERIMENT-DESIGN-06")]
    assert [item["result_id"] for item in result_bindings] == ["RES101", "RES102", "RES201"]
    assert len({item["scientific_input_sha256"] for item in result_bindings}) == 3
    assert all(item["binding_method"] == "canonical_slide_scientific_input" for item in bindings)
    assert not any(item.get("representative_fixture") for item in bindings)


def test_final_native_figure_build_uses_bound_scientific_input_not_representative_route(monkeypatch):
    """A final acceptance figure must not fall back to a route fixture."""
    import thesis_deck_system.phase3_cp5bcd_integrated as structured
    import thesis_deck_system.phase3_cp5efg_integrated as evidence_bound
    from thesis_deck_system.phase3_final_visual_composition import (
        build_bound_native_bundle,
        build_final_composition_plan,
        build_final_projection,
    )

    def forbidden(*_args, **_kwargs):
        raise AssertionError("representative route fixture was invoked")

    monkeypatch.setattr(structured, "build_representative_director_output", forbidden)
    monkeypatch.setattr(evidence_bound, "build_evidence_bound_outputs", forbidden)
    plan = build_final_composition_plan(ROOT, build_final_projection(ROOT))
    slide = next(
        item
        for item in plan["slides"]
        if isinstance(item.get("governed_figure"), dict)
        and item["governed_figure"].get("route") == "scientific_plot"
    )
    governed = slide["governed_figure"]
    binding = next(item for item in plan["slide_scientific_figure_bindings"] if item["binding_id"] == governed["slide_scientific_figure_binding_id"])
    handle, native_plan = build_bound_native_bundle(ROOT, binding, slide["primary_visual_region"])
    assert handle.figure_id == "FIG001"
    assert native_plan["plan_sha256"]


def test_final_composition_build_never_calls_representative_route_fixtures(monkeypatch, tmp_path: Path):
    """The complete final-composition path must preserve the binding boundary."""
    import thesis_deck_system.phase3_cp5bcd_integrated as structured
    import thesis_deck_system.phase3_cp5efg_integrated as evidence_bound
    from thesis_deck_system.phase3_final_visual_composition import _build_final_visual_composition_in_directory

    def forbidden(*_args, **_kwargs):
        raise AssertionError("representative route fixture was invoked")

    monkeypatch.setattr(structured, "build_representative_director_output", forbidden)
    monkeypatch.setattr(evidence_bound, "build_evidence_bound_outputs", forbidden)
    built = _build_final_visual_composition_in_directory(ROOT, tmp_path)
    assert built["figure_audit"]["unapproved_figure_bypass_count"] == 0
    assert built["figure_audit"]["route_only_representative_final_figure_count"] == 0
    assert built["figure_audit"]["unbound_scientific_figure_count"] == 0
    assert built["figure_audit"]["scientific_input_mismatch_count"] == 0
    assert built["figure_audit"]["untruthful_vector_fallback_count"] == 0


def test_figure_audit_counts_materialized_figures_without_governed_bindings_as_bypasses(tmp_path: Path):
    from thesis_deck_system.phase3_final_visual_composition import _audit_figure_composition, _build_final_visual_composition_in_directory

    built = _build_final_visual_composition_in_directory(ROOT, tmp_path)
    ungoverned_plan = copy.deepcopy(built["plan"])
    for slide in ungoverned_plan["slides"]:
        slide["governed_figure"] = None
    audit = _audit_figure_composition(built["deck_path"], ungoverned_plan, {}, [])

    assert audit["unapproved_figure_bypass_count"] > 0


def test_fishbone_fallback_truth_matches_png_physical_representation():
    from thesis_deck_system.phase3_final_visual_composition import build_final_composition_plan, build_final_projection

    plan = build_final_composition_plan(ROOT, build_final_projection(ROOT))
    fishbone = next(item["governed_figure"] for item in plan["slides"] if item["governed_figure"] and item["governed_figure"]["route"] == "fishbone")
    assert fishbone["binding_kind"] == "raster_fallback_explicit"
    assert fishbone["source_authority"] == "fishbone_svg"
    assert fishbone["physical_pptx_representation"] == "png_raster"
    assert fishbone["vector_fallback"] is False
    assert fishbone["raster_fallback"] is True


def test_generated_pptx_adjudication_keeps_raw_scanner_finding_and_requires_execution_attestation(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import (
        GeneratedArtifactAdjudicator,
        GeneratedArtifactAdjudicationError,
    )
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    scanner = RepositoryPrivacyScanner()
    raw = scanner.scan_paths([package], location_root=tmp_path)
    assert [item.classification for item in raw] == ["private_pptx_candidate"]

    adjudicator = GeneratedArtifactAdjudicator(
        root=tmp_path,
        candidate_state_hash="a" * 64,
        approved_producers={"synthetic-template-builder": __file__},
        generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")},
        privacy_scanner=scanner,
    )
    with pytest.raises(GeneratedArtifactAdjudicationError):
        adjudicator.adjudicate(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder")

    record = adjudicator.attest_generated_pptx(
        package,
        artifact_class="sanitized_native_template",
        producer_id="synthetic-template-builder",
        declared_input_paths=[],
        execution_id="FC-R1-TEST-001",
    )
    assert record["status"] == "attested_generated_artifact"
    assert record["raw_scanner_finding"] == "private_pptx_candidate"
    assert record["private_input_count"] == 0
    assert adjudicator.adjudicate(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder")["status"] == "adjudicated_safe_generated_artifact"


@pytest.mark.parametrize("mutation", ["candidate", "sha", "producer", "private_input", "external_relationship", "unknown_part"])
def test_generated_pptx_adjudication_fails_closed_for_mutated_or_unowned_evidence(tmp_path: Path, mutation: str):
    from thesis_deck_system.final_closure_reliability import GeneratedArtifactAdjudicator, GeneratedArtifactAdjudicationError
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    adjudicator = GeneratedArtifactAdjudicator(
        root=tmp_path, candidate_state_hash="b" * 64,
        approved_producers={"synthetic-template-builder": __file__},
        generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")}, privacy_scanner=RepositoryPrivacyScanner(),
    )
    record = adjudicator.attest_generated_pptx(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder", declared_input_paths=[], execution_id="FC-R1-TEST-002")
    mutated = copy.deepcopy(record)
    if mutation == "candidate":
        mutated["candidate_state_hash"] = "c" * 64
    elif mutation == "sha":
        mutated["artifact_sha256"] = "d" * 64
    elif mutation == "producer":
        mutated["producer_id"] = "unknown-producer"
    elif mutation == "private_input":
        mutated["private_input_count"] = 1
    elif mutation == "external_relationship":
        mutated["external_relationship_count"] = 1
    else:
        mutated["unknown_package_part_count"] = 1
    with pytest.raises(GeneratedArtifactAdjudicationError):
        adjudicator.adjudicate_record(mutated)


def test_authoritative_privacy_adjudication_clears_only_matching_attested_generated_pptx(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import GeneratedArtifactAdjudicator, authoritative_privacy_adjudication
    from thesis_deck_system.phase3_privacy import PrivacyFinding, RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    adjudicator = GeneratedArtifactAdjudicator(root=tmp_path, candidate_state_hash="a" * 64, approved_producers={"synthetic-template-builder": __file__}, generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")}, privacy_scanner=RepositoryPrivacyScanner())
    record = adjudicator.attest_generated_pptx(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder", declared_input_paths=[], execution_id="FC-R1-TEST-003")
    accepted = authoritative_privacy_adjudication(
        raw_findings=[PrivacyFinding("private_pptx_candidate", "generated.pptx")],
        adjudicator=adjudicator,
        attestations=[record],
    )
    assert accepted["raw_pptx_candidate_count"] == 1
    assert accepted["attested_generated_pptx_count"] == 1
    assert accepted["unexcepted_final_finding_count"] == 0
    rejected = authoritative_privacy_adjudication(
        raw_findings=[PrivacyFinding("private_pptx_candidate", "unowned.pptx")], adjudicator=adjudicator, attestations=[record],
    )
    assert rejected["unexcepted_final_finding_count"] == 1


def test_final_generated_pptx_contract_set_has_execution_owned_records_for_all_five_current_outputs():
    from thesis_deck_system.final_closure_reliability import attest_final_generated_pptx_set
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    records = attest_final_generated_pptx_set(ROOT, candidate_state_hash="1" * 64, privacy_scanner=RepositoryPrivacyScanner(), execution_id="FC-R1-TEST-005")
    assert len(records) == 5
    assert all(record["status"] == "attested_generated_artifact" for record in records)
    assert all(record["private_input_count"] == 0 for record in records)
    assert all(record["source_closure_input_record_count"] >= 1 for record in records)
    assert all(len(record["source_closure_sha256"]) == 64 for record in records)
    assert all(len(record["media_lineage_sha256"]) == 64 for record in records)
    assert all(record["media_lineage_record_count"] == record["media_part_count"] for record in records)


def test_final_generated_pptx_evidence_bundle_preserves_per_artifact_closure_records():
    from thesis_deck_system.final_closure_reliability import build_final_generated_pptx_evidence_bundle
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    bundle = build_final_generated_pptx_evidence_bundle(
        ROOT, candidate_state_hash="1" * 64, privacy_scanner=RepositoryPrivacyScanner(), execution_id="FC-R1-TEST-006"
    )
    assert len(bundle["attestations"]) == len(bundle["source_closures"]) == len(bundle["media_lineages"]) == 5
    by_path = {item["repository_relative_path"]: item for item in bundle["attestations"]}
    assert all(item["source_closure_id"] == bundle["source_closures"][path]["source_closure_id"] for path, item in by_path.items())
    assert all(item["media_lineage_id"] == bundle["media_lineages"][path]["media_lineage_id"] for path, item in by_path.items())


def test_generated_pptx_path_contract_cannot_be_used_as_a_path_only_allowlist(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import GeneratedArtifactAdjudicator, GeneratedArtifactAdjudicationError
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    adjudicator = GeneratedArtifactAdjudicator(root=tmp_path, candidate_state_hash="a" * 64, approved_producers={"synthetic-template-builder": __file__}, generated_contracts={}, privacy_scanner=RepositoryPrivacyScanner())
    with pytest.raises(GeneratedArtifactAdjudicationError):
        adjudicator.attest_generated_pptx(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder", declared_input_paths=[], execution_id="FC-R1-TEST-004")


def test_native_compiler_downgrades_nonmaterializable_objects_before_assembler(tmp_path: Path):
    from thesis_deck_system.phase3_cp5_hi_final_sprint import ScientificSvgNativeCompiler
    from thesis_deck_system.phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure

    produced = build_representative_director_output(ROOT, "mechanism")
    handle = reverify_approved_figure(produced["manifest"], produced["critic"]["report"], produced["critic"]["approval"], ROOT)
    plan = ScientificSvgNativeCompiler().compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.0, "width": 8.0, "height": 4.0})

    unsupported = {"group", "polyline", "polygon", "marker", "svg_vector"}
    assert all(item["outcome"] != "DRAWINGML_EMITTED" for item in plan["objects"] if item["shape_kind"] in unsupported)


def test_assembler_fails_closed_if_a_native_plan_object_is_not_materialized(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.phase3_cp5_hi_final_sprint import NativeCompilationError, ScientificSvgNativeCompiler, _plan_hash
    from thesis_deck_system.phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure
    from thesis_deck_system.pptx import PythonPptxAssembler

    produced = build_representative_director_output(ROOT, "mechanism")
    handle = reverify_approved_figure(produced["manifest"], produced["critic"]["report"], produced["critic"]["approval"], ROOT)
    plan = ScientificSvgNativeCompiler().compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.0, "width": 8.0, "height": 4.0})
    plan = copy.deepcopy(plan)
    plan["objects"][0]["shape_kind"] = "polyline"
    plan["objects"][0]["outcome"] = "DRAWINGML_EMITTED"
    plan["objects"][0]["fallback_reason"] = None
    plan["fallback_records"] = [item for item in plan["fallback_records"] if item["object_id"] != plan["objects"][0]["svg_object_id"]]
    unhashed = dict(plan); unhashed.pop("plan_sha256")
    plan["plan_sha256"] = _plan_hash(unhashed)
    presentation = Presentation(_synthetic_pptx(tmp_path))
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    with pytest.raises(NativeCompilationError):
        PythonPptxAssembler().add_compiled_figure(slide, handle, plan)


def test_native_style_parity_materializes_supported_fill_stroke_and_text_fields(tmp_path: Path):
    from pptx import Presentation
    from thesis_deck_system.phase3_cp5_hi_final_sprint import ScientificSvgNativeCompiler
    from thesis_deck_system.phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure
    from thesis_deck_system.pptx import PythonPptxAssembler

    produced = build_representative_director_output(ROOT, "mechanism")
    handle = reverify_approved_figure(produced["manifest"], produced["critic"]["report"], produced["critic"]["approval"], ROOT)
    plan = ScientificSvgNativeCompiler().compile(handle, produced["manifest"], produced["svg"], target_box={"left": 1.0, "top": 1.0, "width": 8.0, "height": 4.0})
    presentation = Presentation(_synthetic_pptx(tmp_path))
    facts = PythonPptxAssembler().add_compiled_figure(presentation.slides.add_slide(presentation.slide_layouts[1]), handle, plan)
    coverage = [record["style_coverage"] for record in facts["materialization_records"]]
    assert any(item["fill"] == "supported" and item["stroke"] == "supported" and item["stroke_width"] == "supported" for item in coverage)
    assert any(item["font_family"] == "supported" and item["font_size"] == "supported" for item in coverage)
    assert facts["native_mismatch_count"] == 0


def test_schema_registry_cache_invalidates_by_content_hash_and_preserves_format_checker(tmp_path: Path):
    from thesis_deck_system.contracts import SchemaRegistry

    schema = tmp_path / "minimal.schema.json"
    schema.write_text('{"$schema":"https://json-schema.org/draft/2020-12/schema","type":"object","properties":{"value":{"type":"string","format":"uuid"}},"required":["value"],"additionalProperties":false}', encoding="utf-8")
    first = SchemaRegistry(tmp_path, schema_names=("minimal",))
    assert first.errors("minimal", {"value": "not-a-uuid"})
    schema.write_text('{"$schema":"https://json-schema.org/draft/2020-12/schema","type":"object","properties":{"value":{"type":"integer"}},"required":["value"],"additionalProperties":false}', encoding="utf-8")
    second = SchemaRegistry(tmp_path, schema_names=("minimal",))
    assert second.errors("minimal", {"value": "not-a-uuid"})
    assert second.errors("minimal", {"value": 1}) == []


def test_durable_validation_runner_persists_pre_and_post_candidate_hashes(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import DurableValidationRunner

    runner = DurableValidationRunner(root=ROOT, evidence_root=tmp_path, candidate_hash=lambda: "e" * 64)
    record = runner.run("T1-targeted", ["python", "-c", "print('1 passed')"])
    assert record["completion_status"] == "completed"
    assert record["exit_code"] == 0
    assert record["candidate_hash_pre"] == record["candidate_hash_post"] == "e" * 64
    assert Path(record["stdout_path"]).is_file()
    assert record["passed"] == 1
    assert record["failed"] == 0
    assert Path(record["exit_status_path"]).read_text(encoding="utf-8").strip() == "0"


def test_validation_manifest_partitions_a_complete_collection_deterministically(tmp_path: Path):
    from thesis_deck_system.final_closure_reliability import DurableValidationRunner

    runner = DurableValidationRunner(root=ROOT, evidence_root=tmp_path, candidate_hash=lambda: "e" * 64)
    manifest = runner.build_manifest(["test_a", "test_b", "test_c", "test_d", "test_e"], shard_count=3)
    assert manifest["full_collection_count"] == 5
    assert manifest["executed_unique_count"] == 5
    assert manifest["missing_count"] == 0
    assert manifest["duplicate_count"] == 0
    assert [node for shard in manifest["shards"] for node in shard["node_ids"]] == ["test_a", "test_d", "test_b", "test_e", "test_c"]


def test_final_closure_machine_evidence_is_schema_closed(tmp_path: Path):
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.final_closure_reliability import DurableValidationRunner, GeneratedArtifactAdjudicator, native_materialization_parity
    from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner

    package = _synthetic_pptx(tmp_path)
    adjudicator = GeneratedArtifactAdjudicator(root=tmp_path, candidate_state_hash="f" * 64, approved_producers={"synthetic-template-builder": __file__}, generated_contracts={"generated.pptx": ("sanitized_native_template", "synthetic-template-builder")}, privacy_scanner=RepositoryPrivacyScanner())
    attestation = adjudicator.attest_generated_pptx(package, artifact_class="sanitized_native_template", producer_id="synthetic-template-builder", declared_input_paths=[], execution_id="FC-SCHEMA-001")
    runner = DurableValidationRunner(root=ROOT, evidence_root=tmp_path / "runs", candidate_hash=lambda: "f" * 64)
    run = runner.run("T0-preflight", ["python", "-c", "print('ok')"])
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)
    assert registry.errors("generated-pptx-attestation", attestation) == []
    assert registry.errors("final-closure-validation-run", run) == []
    assert registry.errors("native-materialization-parity", native_materialization_parity([{"planned_native_object_count": 1, "native_object_count": 1, "fallback_object_count": 0, "native_mismatch_count": 0, "materialization_records": [{"style_coverage": {"geometry": "supported"}}]}])) == []


def test_final_closure_qa_is_projected_from_real_parity_privacy_and_validation_facts():
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.final_closure_reliability import build_final_closure_qa

    qa = build_final_closure_qa(
        candidate_state_hash="a" * 64,
        parity={"aggregate_status": "pass", "native_mismatch_count": 0},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
        validation={"exit_code": 0, "candidate_hash_pre": "a" * 64, "candidate_hash_post": "a" * 64, "completion_status": "completed"},
    )
    assert qa["aggregate_status"] == "pass"
    assert all(item["status"] == "pass" for item in qa["owning_checks"])
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)
    assert registry.errors("final-closure-reliability-qa", qa) == []


def test_final_evidence_facts_reject_stale_candidate_or_figure_counts():
    from thesis_deck_system.final_closure_reliability import build_final_closure_qa, build_final_evidence_facts

    facts = build_final_evidence_facts(
        candidate_state_hash="a" * 64,
        focused={"candidate_hash_pre": "a" * 64, "candidate_hash_post": "a" * 64, "exit_code": 0, "passed": 40, "failed": 0},
        figure_audit={"route_only_representative_final_figure_count": 0, "unbound_scientific_figure_count": 0, "scientific_input_mismatch_count": 0, "unapproved_figure_bypass_count": 0, "native_mismatch_count": 0, "untruthful_vector_fallback_count": 0},
        incremental_audit={"stale_mixed_generation_slide_count": 0, "shell_override_by_body_reference_count": 0},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
    )
    assert facts["aggregate_status"] == "pass"
    stale = build_final_evidence_facts(
        candidate_state_hash="a" * 64,
        focused={"candidate_hash_pre": "b" * 64, "candidate_hash_post": "a" * 64, "exit_code": 0, "passed": 40, "failed": 0},
        figure_audit={"route_only_representative_final_figure_count": 0, "unbound_scientific_figure_count": 0, "scientific_input_mismatch_count": 0, "unapproved_figure_bypass_count": 0, "native_mismatch_count": 0, "untruthful_vector_fallback_count": 0},
        incremental_audit={"stale_mixed_generation_slide_count": 0, "shell_override_by_body_reference_count": 0},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
    )
    assert stale["aggregate_status"] == "fail"
    vector_fallback = build_final_evidence_facts(
        candidate_state_hash="a" * 64,
        focused={"candidate_hash_pre": "a" * 64, "candidate_hash_post": "a" * 64, "exit_code": 0, "passed": 40, "failed": 0},
        figure_audit={"route_only_representative_final_figure_count": 0, "unbound_scientific_figure_count": 0, "scientific_input_mismatch_count": 0, "unapproved_figure_bypass_count": 0, "native_mismatch_count": 0, "untruthful_vector_fallback_count": 1},
        incremental_audit={"stale_mixed_generation_slide_count": 0, "shell_override_by_body_reference_count": 0},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
    )
    assert vector_fallback["aggregate_status"] == "fail"
    qa = build_final_closure_qa(
        candidate_state_hash="a" * 64,
        parity={"aggregate_status": "pass", "native_mismatch_count": 1},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
        validation={"exit_code": 0, "candidate_hash_pre": "a" * 64, "candidate_hash_post": "a" * 64, "completion_status": "completed"},
    )
    assert qa["aggregate_status"] == "fail"


def test_final_evidence_facts_are_registered_and_closed():
    """FEC-05 current facts are a closed, schema-owned final artifact."""
    from thesis_deck_system.contracts import SchemaRegistry
    from thesis_deck_system.final_closure_reliability import build_final_evidence_facts

    facts = build_final_evidence_facts(
        candidate_state_hash="a" * 64,
        focused={"candidate_hash_pre": "a" * 64, "candidate_hash_post": "a" * 64, "exit_code": 0, "passed": 42, "failed": 0},
        figure_audit={
            "route_only_representative_final_figure_count": 0,
            "unbound_scientific_figure_count": 0,
            "scientific_input_mismatch_count": 0,
            "unapproved_figure_bypass_count": 0,
            "native_mismatch_count": 0,
            "governed_figure_placement_count": 5,
        },
        incremental_audit={"stale_mixed_generation_slide_count": 0, "shell_override_by_body_reference_count": 0},
        privacy={"unexcepted_final_finding_count": 0, "attested_generated_pptx_count": 5, "raw_pptx_candidate_count": 5},
    )
    registry = SchemaRegistry(ROOT / "thesis-deck-system" / "schemas", include_phase3=True, include_cp5hi=True)
    assert registry.errors("final-evidence-current-facts", facts) == []
    assert registry.errors("final-evidence-current-facts", facts | {"unexpected": True})


def test_final_composition_preserves_source_tree_and_states_template_truth(tmp_path: Path):
    from thesis_deck_system.phase3_final_visual_composition import build_final_visual_composition

    source_tree = ROOT / "thesis-deck-system" / "artifacts" / "phase2" / "fishbone"
    before = {path.name: path.read_bytes() for path in source_tree.iterdir() if path.is_file()}
    result = build_final_visual_composition(ROOT, tmp_path)
    after = {path.name: path.read_bytes() for path in source_tree.iterdir() if path.is_file()}

    assert before == after
    assert result["template_lineage"]["fresh_sanitized_base_template_status"] == "pass"
    assert result["template_lineage"]["professor_shell_tokens_consumed_status"] == "partial_structural"
    assert result["template_lineage"]["physical_professor_template_reconstruction_status"] == "insufficient_evidence"
    assert result["layout_audit"]["legacy_hardcoded_final_composition_bypass_count"] == 0
    assert result["native_parity"]["aggregate_status"] == "pass"
    assert result["native_parity"]["native_mismatch_count"] == 0
    proof = json.loads((tmp_path / "incremental-lineage-acceptance-proof.json").read_text(encoding="utf-8"))
    assert proof["aggregate_status"] == "pass"


def test_svg_compatibility_preview_never_writes_beside_its_source(tmp_path: Path):
    from thesis_deck_system.pptx import _svg_compatibility_preview

    source = tmp_path / "fixtures" / "source.svg"
    source.parent.mkdir()
    source.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 10 10"><rect width="10" height="10"/></svg>', encoding="utf-8")
    output_root = tmp_path / "build" / "previews"
    preview = _svg_compatibility_preview(source, output_root)
    assert preview.is_file()
    assert preview.parent == output_root
    assert not source.with_suffix(".png").exists()


def test_final_composition_transaction_rolls_back_when_assembly_fails(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    from thesis_deck_system.phase3_final_visual_composition import ResultSemanticError, build_final_visual_composition
    from thesis_deck_system.pptx import PythonPptxAssembler

    prior = tmp_path / "cp5-final-visual-composition-acceptance-deck.pptx"
    prior.write_bytes(b"previous-release")
    monkeypatch.setattr(PythonPptxAssembler, "assemble_final_visual_composition", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("synthetic failure")))
    with pytest.raises(RuntimeError):
        build_final_visual_composition(ROOT, tmp_path)
    assert prior.read_bytes() == b"previous-release"


def test_final_candidate_state_binds_closure_adjudication_parity_and_validation_dependencies():
    from thesis_deck_system.phase3_final_visual_composition import compute_final_visual_composition_candidate_state

    state = compute_final_visual_composition_candidate_state(ROOT)
    required = {
        "packages/thesis-deck-system/src/thesis_deck_system/final_closure_reliability.py",
        "packages/thesis-deck-system/src/thesis_deck_system/contracts.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_privacy.py",
        "packages/thesis-deck-system/tests/unit/test_final_closure_reliability.py",
        ".gitattributes",
        "thesis-deck-system/schemas/generated-pptx-attestation.schema.json",
        "thesis-deck-system/schemas/native-materialization-parity.schema.json",
        "thesis-deck-system/schemas/final-closure-validation-run.schema.json",
        "thesis-deck-system/schemas/final-evidence-current-facts.schema.json",
    }
    assert required <= set(state["component_hashes"])


def test_final_candidate_component_digest_normalizes_text_toml_line_endings_but_not_binary_inputs():
    from thesis_deck_system.phase3_final_visual_composition import _candidate_component_digest

    assert _candidate_component_digest("packages/thesis-deck-system/pyproject.toml", b"[x]\r\ny=1\r\n") == _candidate_component_digest("packages/thesis-deck-system/pyproject.toml", b"[x]\ny=1\n")
    assert _candidate_component_digest(".gitattributes", b"path text eol=lf\r\n") == _candidate_component_digest(".gitattributes", b"path text eol=lf\n")
    assert _candidate_component_digest("fixture.pptx", b"a\r\nb") != _candidate_component_digest("fixture.pptx", b"a\nb")
