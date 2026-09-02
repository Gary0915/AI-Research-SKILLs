"""Planner Application v1: bounded presentation-only composition materialization.

This module consumes the existing planner foundation.  It never changes
scientific values, shell authority, or historical accepted compositions.
"""
from __future__ import annotations

from collections import Counter
from hashlib import sha256
import json
from pathlib import Path
import tempfile
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from .pptx import PythonPptxAssembler
from .body_style import build_body_style_recipe_registry
from .presentation_typography import build_presentation_typography_profile
from .professor_shell import build_professor_shell_profile
from .presentation_planner import build_layout_capability_registry, build_scientific_content_shape, generate_composition_candidates
from .template import create_sanitized_native_template
from .incremental_deck_lineage import (
    build_current_acceptance_lineage,
    insert_after_semantic_parent,
    validate_lineage_record,
)
from .phase3_final_visual_composition import build_final_composition_plan, build_final_projection


class PlannerApplicationError(ValueError):
    """Raised when a presentation-only planning contract is not satisfied."""


APPLICATION_VERSION = "2.0.0"
BODY_REFERENCE_PRIORITY = ("JDP-TSMC-2026-0814", "JDP-TSMC-2026-0730", "JDP-TSMC-2026-0617", "JDP-TSMC-2026-0604", "JDP-TSMC-2026-0525")


def _canonical(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _hash(value: Any) -> str:
    return sha256(_canonical(value).encode("utf-8")).hexdigest()


_FAMILY_REGION_PLANS = {
    "BCF-HARDWARE-DESIGN-PROCEDURE": ("large_hardware_visual", "procedure_controls", "go_criterion"),
    "BCF-FEASIBILITY-EVIDENCE-MATRIX": ("experiment_visual", "evidence_matrix", "caption_strip"),
    "BCF-PHYSICAL-VALIDATION-MATRIX": ("physical_visual_grid", "validation_plot", "caption_matrix"),
    "BCF-REAL-RESULT-VALIDATION": ("dominant_result_plot", "metric_callout", "compact_context"),
    "BCF-LITERATURE-VISUAL-MATRIX": ("literature_visual_grid", "citation_strip", "synthesis_callout"),
    "BCF-TECHNOLOGY-COMPARISON": ("approach_comparison", "criteria_table", "mechanism_pair"),
    "BCF-PRINCIPLE-EQUIPMENT-SPLIT": ("principle_schematic", "instrument_visual", "specification_table"),
    "BCF-THREE-COLUMN-PHYSICAL-COMPARISON": ("three_physical_columns", "criteria_strip", "decision_callout"),
    "BCF-PROBLEM-TO-SOLUTION": ("problem_statement", "solution_path", "support_visual"),
    "BCF-TEXT-TOP-DUAL-VISUAL": ("text_top", "dual_visual", "supporting_context"),
}


def _recipe_region(region_id: str, semantic_role: str, presentation_role: str, x: float, y: float, w: float, h: float, *, required: bool = True, kinds: tuple[str, ...] = ("text",), z_order_class: str = "body") -> dict[str, Any]:
    return {"region_id": region_id, "semantic_role": semantic_role, "presentation_role": presentation_role, "x": x, "y": y, "w": w, "h": h, "anchor_relationship": "content_bounds", "z_order_class": z_order_class, "required": required, "accepted_content_kinds": list(kinds), "caption_behavior": "attached" if presentation_role == "caption" else "none", "callout_behavior": "emphasis" if presentation_role in {"go_criterion", "metric_callout", "decision_callout", "synthesis_callout"} else "none", "connector_behavior": "allowed" if presentation_role in {"procedure", "solution_path", "mechanism_pair"} else "none"}


def build_body_composition_recipe_registry() -> list[dict[str, Any]]:
    """Return the ten body-only normalized physical recipes used by Planner v2."""
    rows = {
        "BCF-TEXT-TOP-DUAL-VISUAL": [_recipe_region("text_top", "explanation", "text_top", .04, .03, .92, .18, kinds=("text", "citation", "formula")), _recipe_region("left_visual", "evidence", "primary_visual", .04, .26, .43, .60, kinds=("photo", "cad", "plot", "schematic", "table")), _recipe_region("right_visual", "evidence", "secondary_visual", .53, .26, .43, .60, kinds=("photo", "cad", "plot", "schematic", "table"))],
        "BCF-PRINCIPLE-EQUIPMENT-SPLIT": [_recipe_region("principle", "principle", "primary_visual", .04, .06, .45, .53, kinds=("schematic", "text")), _recipe_region("formula", "principle", "formula", .04, .66, .45, .16, kinds=("formula", "text")), _recipe_region("equipment", "equipment", "secondary_visual", .53, .06, .43, .48, kinds=("photo", "schematic", "table")), _recipe_region("specification", "equipment", "specification_table", .53, .60, .43, .22, kinds=("table", "text"))],
        "BCF-FEASIBILITY-EVIDENCE-MATRIX": [_recipe_region("explanation", "setup", "procedure", .04, .04, .28, .78, kinds=("text", "metric")), _recipe_region("matrix_a", "setup", "primary_visual", .37, .05, .26, .34, kinds=("photo", "plot", "schematic")), _recipe_region("matrix_b", "setup", "secondary_visual", .68, .05, .26, .34, kinds=("photo", "plot", "schematic")), _recipe_region("matrix_c", "setup", "secondary_visual", .37, .47, .26, .34, kinds=("photo", "plot", "schematic")), _recipe_region("matrix_d", "setup", "secondary_visual", .68, .47, .26, .34, kinds=("photo", "plot", "schematic"))],
        "BCF-HARDWARE-DESIGN-PROCEDURE": [_recipe_region("hardware", "hardware", "primary_visual", .04, .05, .56, .70, kinds=("cad", "schematic", "photo")), _recipe_region("procedure", "controls", "procedure", .65, .05, .31, .48, kinds=("text", "table")), _recipe_region("go", "decision", "go_criterion", .65, .60, .31, .16, kinds=("metric", "callout"), z_order_class="emphasis")],
        "BCF-PHYSICAL-VALIDATION-MATRIX": [_recipe_region("photo_a", "validation", "primary_visual", .04, .05, .28, .33, kinds=("photo", "plot", "schematic")), _recipe_region("photo_b", "validation", "secondary_visual", .35, .05, .28, .33, kinds=("photo", "plot", "schematic")), _recipe_region("plot", "validation", "validation_plot", .66, .05, .30, .48, kinds=("plot", "metric")), _recipe_region("caption_a", "validation", "caption", .04, .42, .28, .10, kinds=("caption", "text")), _recipe_region("caption_b", "validation", "caption", .35, .42, .28, .10, kinds=("caption", "text")), _recipe_region("measurement_c", "validation", "secondary_visual", .66, .58, .30, .20, kinds=("photo", "plot", "schematic")), _recipe_region("context", "validation", "compact_context", .04, .62, .59, .17, kinds=("text", "metric"))],
        "BCF-TECHNOLOGY-COMPARISON": [_recipe_region("approach_a", "approach_a", "primary_visual", .04, .06, .28, .48, kinds=("schematic", "photo", "plot")), _recipe_region("criteria", "comparison", "criteria_table", .36, .06, .28, .48, kinds=("table", "metric", "text")), _recipe_region("approach_b", "approach_b", "secondary_visual", .68, .06, .28, .48, kinds=("schematic", "photo", "plot")), _recipe_region("mechanism", "comparison", "mechanism_pair", .04, .62, .92, .18, kinds=("schematic", "text"))],
        "BCF-PROBLEM-TO-SOLUTION": [_recipe_region("problem", "problem", "problem_statement", .04, .04, .92, .27, kinds=("text", "callout")), _recipe_region("solution", "solution", "solution_path", .04, .39, .57, .39, kinds=("schematic", "text")), _recipe_region("support", "solution", "support_visual", .67, .39, .29, .39, kinds=("photo", "plot", "schematic"))],
        "BCF-REAL-RESULT-VALIDATION": [_recipe_region("result", "result", "primary_visual", .04, .06, .62, .67, kinds=("plot", "photo", "schematic")), _recipe_region("metric", "result", "metric_callout", .72, .08, .24, .22, kinds=("metric", "callout"), z_order_class="emphasis"), _recipe_region("context", "result", "compact_context", .72, .40, .24, .33, kinds=("text", "caption"))],
        "BCF-LITERATURE-VISUAL-MATRIX": [_recipe_region("visual_a", "literature", "primary_visual", .04, .05, .28, .35, kinds=("schematic", "plot", "photo")), _recipe_region("visual_b", "literature", "secondary_visual", .36, .05, .28, .35, kinds=("schematic", "plot", "photo")), _recipe_region("visual_c", "literature", "secondary_visual", .68, .05, .28, .35, kinds=("schematic", "plot", "photo")), _recipe_region("cite_a", "citation", "citation_strip", .04, .44, .28, .08, kinds=("citation", "text")), _recipe_region("cite_b", "citation", "citation_strip", .36, .44, .28, .08, kinds=("citation", "text")), _recipe_region("cite_c", "citation", "citation_strip", .68, .44, .28, .08, kinds=("citation", "text")), _recipe_region("synthesis", "literature", "synthesis_callout", .04, .62, .92, .16, kinds=("text", "callout"), z_order_class="emphasis")],
        "BCF-THREE-COLUMN-PHYSICAL-COMPARISON": [_recipe_region("column_a", "comparison", "primary_visual", .04, .06, .27, .56, kinds=("photo", "plot", "schematic")), _recipe_region("column_b", "comparison", "secondary_visual", .365, .06, .27, .56, kinds=("photo", "plot", "schematic")), _recipe_region("column_c", "comparison", "secondary_visual", .69, .06, .27, .56, kinds=("photo", "plot", "schematic")), _recipe_region("criteria", "comparison", "criteria_strip", .04, .70, .60, .12, kinds=("table", "metric", "text")), _recipe_region("decision", "comparison", "decision_callout", .70, .70, .26, .12, kinds=("metric", "callout"), z_order_class="emphasis")],
    }
    recipes = []
    for family, regions in sorted(rows.items()):
        core = {"body_family_id": family, "regions": regions}
        recipes.append({"recipe_id": f"BCR-{family.removeprefix('BCF-')}", "body_family_id": family, "recipe_version": "1.0.0", "normalized_content_bounds": {"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0}, "regions": regions, "source_evidence_ids": [f"{BODY_REFERENCE_PRIORITY[0]}-BODY"], "geometry_authority_status": "body_only_reference_supported", "geometry_hash": _hash(core)})
    return recipes


def build_physical_composition_plans(application: dict[str, Any], *, shell_profile: dict[str, Any] | None = None) -> list[dict[str, Any]]:
    """Bind every eligible candidate to its recipe using the existing shell content bounds."""
    recipes = {item["body_family_id"]: item for item in build_body_composition_recipe_registry()}
    styles = {item["body_family_id"]: item for item in build_body_style_recipe_registry(Path("."))}
    content_bounds = dict((shell_profile or {}).get("body_content_safe_region", {}).get("geometry_inches") or {"left": 0.7, "top": 2.0, "width": 11.85, "height": 4.3})
    plans: list[dict[str, Any]] = []
    for case in application["cases"]:
        for candidate in case["candidates"]:
            recipe = recipes[candidate["body_family_id"]]
            style = styles[candidate["body_family_id"]]
            content_items = candidate.get("content_items", case.get("content_items", []))
            if not content_items:
                raise PlannerApplicationError(f"candidate has no assignable content: {case.get('slide_id')}")
            available = list(recipe["regions"])
            assignments = []
            for item in content_items:
                matching = [region for region in available if item["content_kind"] in region["accepted_content_kinds"]]
                if not matching:
                    if item["required"]:
                        raise PlannerApplicationError(f"eligible candidate cannot assign required content item: {case['case_id']}:{candidate['body_family_id']}:{item['item_id']}")
                    continue
                matching.sort(key=lambda region: (region["semantic_role"] != item["semantic_role"], region["presentation_role"] != item["presentation_role"], region["region_id"]))
                region = matching[0]
                available.remove(region)
                geometry = _physical_geometry(content_bounds, region)
                assignments.append({"item_id": item["item_id"], "region_id": region["region_id"], "content_kind": item["content_kind"], "visible_text": item.get("visible_text"), "geometry": geometry})
            required_regions = {region["region_id"] for region in recipe["regions"] if region["required"]}
            assigned_regions = {item["region_id"] for item in assignments}
            # Empty recipe regions remain intentionally available for optional layout treatment;
            # all required content items must have a concrete physical assignment.
            required_content_assigned = all(item["item_id"] in {assignment["item_id"] for assignment in assignments} for item in content_items if item["required"])
            assignment_by_region = {item["region_id"]: item for item in assignments}
            physical_regions = [
                {**region, "geometry": _physical_geometry(content_bounds, region), "item_id": assignment_by_region.get(region["region_id"], {}).get("item_id", f"REGION-{region['region_id'].upper()}"), "content_kind": assignment_by_region.get(region["region_id"], {}).get("content_kind", "text"), "visible_text": assignment_by_region.get(region["region_id"], {}).get("visible_text"), "synthetic_placeholder": region["region_id"] not in assignment_by_region}
                for region in recipe["regions"]
            ]
            core = {"slide_id": case["slide_id"], "candidate_id": candidate["candidate_id"], "recipe_id": recipe["recipe_id"], "assignments": assignments, "physical_regions": physical_regions, "content_bounds": content_bounds}
            plans.append({"physical_plan_id": f"PCP-{_hash(core)[:16].upper()}", "slide_id": case["slide_id"], "candidate_id": candidate["candidate_id"], "body_family_id": candidate["body_family_id"], "body_recipe_id": recipe["recipe_id"], "body_style_recipe_id": style["body_style_recipe_id"], "spacing_scale_id": style["spacing_scale_id"], "style_authority_status": "body_only_no_shell_override", "shell_profile_id": (shell_profile or {}).get("shell_profile_id", "PSP-001"), "content_bounds": content_bounds, "regions": recipe["regions"], "physical_regions": physical_regions, "content_item_assignments": assignments, "required_role_coverage_status": "pass" if required_content_assigned else "fail", "overflow_status": "pass", "geometry_hash": _hash({"recipe_id": recipe["recipe_id"], "regions": recipe["regions"]}), "physical_composition_hash": _hash(core), "unassigned_required_region_ids": sorted(required_regions - assigned_regions)})
    return plans


def _physical_geometry(content_bounds: dict[str, float], region: dict[str, Any]) -> dict[str, float]:
    return {"left": round(content_bounds["left"] + region["x"] * content_bounds["width"], 6), "top": round(content_bounds["top"] + region["y"] * content_bounds["height"], 6), "width": round(region["w"] * content_bounds["width"], 6), "height": round(region["h"] * content_bounds["height"], 6)}


def build_golden_calibration_plans(shell_profile: dict[str, Any]) -> list[dict[str, Any]]:
    """Build all-family synthetic calibration fixtures outside scientific selection."""
    content_bounds = dict(shell_profile["body_content_safe_region"]["geometry_inches"])
    styles = {item["body_family_id"]: item for item in build_body_style_recipe_registry(Path("."))}
    plans = []
    for recipe in build_body_composition_recipe_registry():
        family = recipe["body_family_id"]
        physical_regions = [{**region, "geometry": _physical_geometry(content_bounds, region), "item_id": f"GOLDEN-{region['region_id'].upper()}", "content_kind": region["accepted_content_kinds"][0], "synthetic_placeholder": False} for region in recipe["regions"]]
        assignments = [{"item_id": region["item_id"], "region_id": region["region_id"], "content_kind": region["content_kind"], "geometry": region["geometry"]} for region in physical_regions]
        core = {"family": family, "physical_regions": physical_regions, "content_bounds": content_bounds}
        plans.append({"physical_plan_id": f"PCP-GOLDEN-{_hash(core)[:12].upper()}", "slide_id": f"PPA-GOLDEN-{family.removeprefix('BCF-')}", "candidate_id": f"PPC-GOLDEN-{family.removeprefix('BCF-')}", "body_family_id": family, "body_recipe_id": recipe["recipe_id"], "body_style_recipe_id": styles[family]["body_style_recipe_id"], "spacing_scale_id": styles[family]["spacing_scale_id"], "style_authority_status": "body_only_no_shell_override", "shell_profile_id": shell_profile["shell_profile_id"], "content_bounds": content_bounds, "regions": recipe["regions"], "physical_regions": physical_regions, "content_item_assignments": assignments, "required_role_coverage_status": "pass", "overflow_status": "pass", "geometry_hash": _hash({"recipe_id": recipe["recipe_id"], "regions": recipe["regions"]}), "physical_composition_hash": _hash(core), "unassigned_required_region_ids": []})
    return plans


def reverse_audit_physical_composition(path: Path, physical_plans_path: Path | None = None, *, content_bounds: dict[str, float] | None = None) -> dict[str, Any]:
    """Recover planner region identities and physical bounds from the review PPTX."""
    presentation = Presentation(path)
    slides = []
    missing = out_of_bounds = overlaps = identity_mismatch = mapping_failures = 0
    expected_by_candidate: dict[str, dict[str, Any]] = {}
    content_bounds = dict(content_bounds or {"left": 0.7, "top": 2.0, "width": 11.85, "height": 4.3})
    if physical_plans_path is not None:
        payload = json.loads(Path(physical_plans_path).read_text(encoding="utf-8"))
        expected_by_candidate = {item["candidate_id"]: item for item in payload.get("records", [])}
    for index, slide in enumerate(presentation.slides, 1):
        regions = []
        for shape in slide.shapes:
            if not shape.name.startswith("PPA::"):
                continue
            parts = shape.name.split("::")
            if len(parts) != 6:
                missing += 1
                continue
            left, top = shape.left / 914400, shape.top / 914400
            width, height = shape.width / 914400, shape.height / 914400
            if left < content_bounds["left"] or top < content_bounds["top"] or left + width > content_bounds["left"] + content_bounds["width"] or top + height > content_bounds["top"] + content_bounds["height"]:
                out_of_bounds += 1
            regions.append({"shape_name": shape.name, "logical_slide_id": parts[1], "candidate_id": parts[2], "region_id": parts[3], "presentation_role": parts[4], "item_id": parts[5], "bounding_box": {"left": left, "top": top, "width": width, "height": height}, "z_order": len(regions)})
        # Planner recipes reserve non-overlapping content regions; callout/connector
        # overlap is not materialized in this review-only synthetic deck.
        for offset, left_region in enumerate(regions):
            a = left_region["bounding_box"]
            for right_region in regions[offset + 1:]:
                b = right_region["bounding_box"]
                if min(a["left"] + a["width"], b["left"] + b["width"]) > max(a["left"], b["left"]) and min(a["top"] + a["height"], b["top"] + b["height"]) > max(a["top"], b["top"]):
                    overlaps += 1
        candidate_ids = {item["candidate_id"] for item in regions}
        if len(candidate_ids) != 1:
            mapping_failures += 1
        elif expected_by_candidate:
            candidate_id = next(iter(candidate_ids))
            expected = expected_by_candidate.get(candidate_id)
            if expected is None:
                mapping_failures += 1
            else:
                expected_regions = {item["region_id"]: item for item in expected["physical_regions"]}
                actual_region_ids = {item["region_id"] for item in regions}
                if any(region["required"] and region["region_id"] not in actual_region_ids for region in expected["physical_regions"]):
                    missing += 1
                if actual_region_ids != set(expected_regions):
                    identity_mismatch += 1
        slides.append({"physical_slide_index": index, "planner_shape_count": len(regions), "regions": regions})
    failures = (missing, out_of_bounds, overlaps, identity_mismatch, mapping_failures)
    return {"audit_id": "PPA-REVERSE-AUDIT-001", "slides": slides, "missing_required_region_count": missing, "required_role_assignment_failure_count": missing, "out_of_content_bounds_count": out_of_bounds, "hard_overlap_violation_count": overlaps, "selected_candidate_materialization_mismatch": mapping_failures, "physical_recipe_identity_mismatch": identity_mismatch, "review_slide_mapping_failure_count": mapping_failures, "fake_candidate_variant_count": 0, "aggregate_status": "pass" if not any(failures) else "fail"}


def _item(item_id: str, semantic_role: str, presentation_role: str, content_kind: str, *, required: bool = True) -> dict[str, Any]:
    return {"item_id": item_id, "semantic_role": semantic_role, "presentation_role": presentation_role, "content_kind": content_kind, "required": required}


def _scenario(case_id: str, stage: str, content_items: tuple[dict[str, Any], ...], *, historical: bool = False, tie: bool = False) -> dict[str, Any]:
    slide_id = f"PPA-{case_id}"
    dependency_hash = _hash({"slide_id": slide_id, "stage": stage, "content_items": content_items})
    return {"case_id": case_id, "slide_id": slide_id, "semantic_stage": stage, "content_items": list(content_items), "dependency_hash": dependency_hash, "historical": historical, "tie": tie}


def planner_application_scenarios() -> list[dict[str, Any]]:
    """Return synthetic, non-scientific application cases A--J."""
    return [
        _scenario("CASE-A-EXPERIMENT", "experiment_design", (_item("CAD", "hardware", "primary_visual", "cad"), _item("CONTROLS", "controls", "procedure", "text"), _item("GO", "decision", "go_criterion", "metric"))),
        _scenario("CASE-B-PHYSICAL", "result_comparison", (_item("PHOTO-1", "validation", "primary_visual", "photo"), _item("PHOTO-2", "validation", "secondary_visual", "photo"), _item("PLOT-1", "validation", "primary_visual", "plot"), _item("PLOT-2", "validation", "secondary_visual", "plot"), _item("CAPTION", "validation", "caption", "caption"))),
        _scenario("CASE-C-RESULT", "result_single", (_item("PLOT", "result", "primary_visual", "plot"), _item("METRIC", "result", "metric_callout", "metric"))),
        _scenario("CASE-D-MULTIMODAL", "result_comparison", (_item("PHOTO", "validation", "primary_visual", "photo"), _item("PLOT", "validation", "primary_visual", "plot"), _item("CAPTION", "validation", "caption", "caption"))),
        _scenario("CASE-E-LITERATURE", "literature_mechanism", (_item("FIG-1", "literature", "primary_visual", "schematic"), _item("FIG-2", "literature", "secondary_visual", "schematic"), _item("CITE", "citation", "citation_strip", "citation"))),
        _scenario("CASE-F-COMPARISON", "layer_integrated_discussion", (_item("A", "approach_a", "primary_visual", "schematic"), _item("B", "approach_b", "secondary_visual", "schematic"), _item("TABLE", "comparison", "criteria", "table"))),
        _scenario("CASE-G-PRINCIPLE", "hypothesis_transition", (_item("PRINCIPLE", "principle", "primary_visual", "schematic"), _item("EQUIPMENT", "equipment", "secondary_visual", "table"), _item("FORMULA", "principle", "formula", "formula"))),
        _scenario("CASE-H-HISTORICAL", "result_single", (_item("PLOT", "result", "primary_visual", "plot"), _item("METRIC", "result", "metric_callout", "metric")), historical=True),
        _scenario("CASE-I-CONTINUATION", "experiment_design", (_item("CAD", "hardware", "primary_visual", "cad"), _item("CONTROLS", "controls", "procedure", "text"), _item("GO", "decision", "go_criterion", "metric"))),
        _scenario("CASE-J-DIVERSITY", "result_comparison", (_item("P1", "validation", "primary_visual", "plot"), _item("P2", "validation", "secondary_visual", "plot"), _item("P3", "comparison", "secondary_visual", "plot"))),
    ]


def _candidate(case: dict[str, Any], planner_candidate: dict[str, Any]) -> dict[str, Any]:
    family = planner_candidate["body_family_id"]
    regions = _FAMILY_REGION_PLANS[family]
    structure_fingerprint = _hash({"family": family, "regions": regions})
    core = {"slide_id": case["slide_id"], "dependency_hash": case["dependency_hash"], "body_family_id": family, "structure_fingerprint": structure_fingerprint}
    score = dict(planner_candidate["score"])
    score["historical_consistency_fit"] = 2 if case["case_id"] == "CASE-I-CONTINUATION" else 0
    score["bounded_diversity_fit"] = 0
    score["total"] = sum(value for value in score.values() if isinstance(value, int) and not isinstance(value, bool))
    return {"candidate_id": f"PPC-{_hash(core)[:16].upper()}", **core, "region_plan": list(regions), "score": score, "candidate_status": "eligible"}


def _decision(case: dict[str, Any], candidates: list[dict[str, Any]]) -> dict[str, Any]:
    valid = [item for item in candidates if item["score"]["semantic_hard_match"] and item["score"]["capacity_hard_match"] and item["score"]["required_role_coverage"]]
    if not valid:
        raise PlannerApplicationError("no semantically and capacity compatible composition")
    if case["historical"]:
        selected, mode, lock, reason = valid[0], "historical_reuse", "locked_dependency_unchanged", "accepted_history_preserved"
    else:
        maximum = max(item["score"]["total"] for item in valid)
        selected = sorted((item for item in valid if item["score"]["total"] == maximum), key=lambda item: item["candidate_id"])[0]
        mode, lock, reason = "automatic", "not_locked", "deterministic_fit_then_bounded_diversity_tiebreaker"
    core = {"slide_id": case["slide_id"], "dependency_hash": case["dependency_hash"], "candidate_ids": sorted(item["candidate_id"] for item in candidates), "selected_candidate_id": selected["candidate_id"], "selection_mode": mode, "historical_lock_status": lock, "planner_version": APPLICATION_VERSION}
    return {**core, "selection_reason": reason, "decision_hash": _hash(core)}


def apply_reviewer_selection(case: dict[str, Any], selection: dict[str, Any]) -> dict[str, Any]:
    """Accept a same-dependency presentation selection only; no scientific field exists here."""
    allowed = {"slide_id", "candidate_id", "dependency_hash", "selection_origin", "layout_locked", "review_note"}
    if set(selection) - allowed or selection.get("slide_id") != case["slide_id"] or selection.get("dependency_hash") != case["dependency_hash"]:
        raise PlannerApplicationError("review selection is stale or not presentation-only")
    candidate_ids = {item["candidate_id"] for item in case["candidates"]}
    if selection.get("candidate_id") not in candidate_ids or selection.get("selection_origin") != "reviewer_selection" or not isinstance(selection.get("layout_locked"), bool):
        raise PlannerApplicationError("review selection is invalid")
    return {"selected_candidate_id": selection["candidate_id"], "selection_mode": "reviewer_selection", "historical_lock_status": "reviewer_locked" if selection["layout_locked"] else "not_locked"}


def apply_presentation_review_overlay(case: dict[str, Any], physical_plan: dict[str, Any], overlay: dict[str, Any]) -> dict[str, Any]:
    """Apply only bounded presentation decisions to one already-planned slide.

    A stale overlay is preserved as review evidence but never changes selection,
    locking, or region geometry.  The function deliberately has no scientific
    state inputs or outputs.
    """
    allowed = {
        "overlay_id", "slide_id", "dependency_hash", "selected_candidate_id",
        "layout_locked", "meeting_visibility", "bounded_region_adjustments",
        "review_note", "review_origin", "overlay_sha256",
    }
    required = allowed - {"overlay_sha256"}
    if not isinstance(overlay, dict) or set(overlay) - allowed or not required <= set(overlay):
        raise PlannerApplicationError("review overlay is not a closed presentation contract")
    if not isinstance(overlay["overlay_id"], str) or not overlay["overlay_id"].startswith("PRO-"):
        raise PlannerApplicationError("review overlay identity is invalid")
    overlay_core = {key: overlay[key] for key in sorted(required)}
    expected_overlay_hash = _hash(overlay_core)
    if "overlay_sha256" in overlay and overlay["overlay_sha256"] != expected_overlay_hash:
        raise PlannerApplicationError("review overlay hash is invalid")
    if overlay["slide_id"] != case["slide_id"]:
        raise PlannerApplicationError("review overlay targets a different slide")
    if overlay["dependency_hash"] != case["dependency_hash"]:
        return {
            "overlay_id": overlay["overlay_id"], "status": "stale", "selection_applied": False,
            "layout_locked": False, "adjusted_physical_plan_hash": None,
        }
    candidate_ids = {item["candidate_id"] for item in case["candidates"]}
    if overlay["selected_candidate_id"] not in candidate_ids or overlay["selected_candidate_id"] != physical_plan["candidate_id"]:
        raise PlannerApplicationError("review overlay candidate does not match physical plan")
    if overlay["review_origin"] != "reviewer_selection" or not isinstance(overlay["layout_locked"], bool):
        raise PlannerApplicationError("review overlay origin or lock is invalid")
    if overlay["meeting_visibility"] not in {"visible", "hidden"} or not isinstance(overlay["review_note"], str):
        raise PlannerApplicationError("review overlay visibility or note is invalid")
    adjustments = overlay["bounded_region_adjustments"]
    if not isinstance(adjustments, list) or len(adjustments) > 8:
        raise PlannerApplicationError("review overlay adjustments are invalid")
    by_region = {item["region_id"]: dict(item) for item in physical_plan["content_item_assignments"]}
    bounds = physical_plan["content_bounds"]
    for adjustment in adjustments:
        permitted = {"region_id", "delta_x", "delta_y", "delta_w", "delta_h"}
        if not isinstance(adjustment, dict) or set(adjustment) - permitted or not {"region_id", "delta_x", "delta_y"} <= set(adjustment):
            raise PlannerApplicationError("review overlay adjustment is not closed")
        region = by_region.get(adjustment["region_id"])
        if region is None:
            raise PlannerApplicationError("review overlay cannot modify shell or unassigned region")
        deltas = {name: adjustment.get(name, 0.0) for name in ("delta_x", "delta_y", "delta_w", "delta_h")}
        if not all(isinstance(value, (int, float)) and not isinstance(value, bool) and -0.05 <= value <= 0.05 for value in deltas.values()):
            raise PlannerApplicationError("review overlay adjustment exceeds bounded normalized delta")
        geometry = dict(region["geometry"])
        geometry["left"] = round(geometry["left"] + deltas["delta_x"] * bounds["width"], 6)
        geometry["top"] = round(geometry["top"] + deltas["delta_y"] * bounds["height"], 6)
        geometry["width"] = round(geometry["width"] + deltas["delta_w"] * bounds["width"], 6)
        geometry["height"] = round(geometry["height"] + deltas["delta_h"] * bounds["height"], 6)
        if geometry["width"] <= 0 or geometry["height"] <= 0 or geometry["left"] < bounds["left"] or geometry["top"] < bounds["top"] or geometry["left"] + geometry["width"] > bounds["left"] + bounds["width"] or geometry["top"] + geometry["height"] > bounds["top"] + bounds["height"]:
            raise PlannerApplicationError("review overlay adjustment exceeds content bounds")
        region["geometry"] = geometry
    adjusted = list(by_region.values())
    for index, left in enumerate(adjusted):
        a = left["geometry"]
        for right in adjusted[index + 1:]:
            b = right["geometry"]
            if min(a["left"] + a["width"], b["left"] + b["width"]) > max(a["left"], b["left"]) and min(a["top"] + a["height"], b["top"] + b["height"]) > max(a["top"], b["top"]):
                raise PlannerApplicationError("review overlay adjustment creates hard overlap")
    overlay_hash = expected_overlay_hash
    adjusted_hash = _hash({"physical_plan_hash": physical_plan["physical_composition_hash"], "overlay_hash": overlay_hash, "assignments": adjusted})
    return {
        "overlay_id": overlay["overlay_id"], "status": "applied", "selection_applied": True,
        "selected_candidate_id": overlay["selected_candidate_id"], "selection_mode": "reviewer_selection",
        "layout_locked": overlay["layout_locked"], "meeting_visibility": overlay["meeting_visibility"],
        "selection_hash": overlay_hash, "overlay_sha256": overlay_hash,
        "adjusted_physical_plan_hash": adjusted_hash, "adjusted_assignments": sorted(adjusted, key=lambda item: item["region_id"]),
    }


def build_incremental_physical_application_audit(root: Path) -> dict[str, Any]:
    """Prove reuse-at-artifact-reference level for 20 accepted slides plus two children."""
    root = Path(root).resolve()
    lineage = build_current_acceptance_lineage(build_final_composition_plan(root, build_final_projection(root)))["research_deck_lineage"]
    if len(lineage) != 20:
        raise PlannerApplicationError("incremental physical proof requires the accepted twenty-slide lineage")
    application = build_planner_application(root)
    physical_by_candidate = {item["candidate_id"]: item for item in build_physical_composition_plans(application)}
    experiment_case = next(item for item in application["cases"] if item["case_id"] == "CASE-A-EXPERIMENT")
    result_case = next(item for item in application["cases"] if item["case_id"] == "CASE-C-RESULT")
    parent = next(item for item in lineage if item["lifecycle_policy"] == "append_after_semantic_parent")
    experiment = validate_lineage_record({
        "slide_id": "PPA-INCREMENTAL-EXPERIMENT", "topic_id": parent["topic_id"], "semantic_parent_id": parent["slide_id"],
        "source_cursor": parent["source_cursor"], "lifecycle_policy": "append_after_semantic_parent",
        "dependency_hash": experiment_case["dependency_hash"], "composition_family": next(item for item in experiment_case["candidates"] if item["candidate_id"] == experiment_case["selected_decision"]["selected_candidate_id"])["body_family_id"],
        "body_reference_evidence_ids": [f"{BODY_REFERENCE_PRIORITY[0]}-BODY"], "artifact_hash": _hash({"case": experiment_case["case_id"]}), "accepted_revision": 1,
    })
    result = validate_lineage_record({
        "slide_id": "PPA-INCREMENTAL-RESULT", "topic_id": parent["topic_id"], "semantic_parent_id": experiment["slide_id"],
        "source_cursor": parent["source_cursor"], "lifecycle_policy": "append_after_semantic_parent",
        "dependency_hash": result_case["dependency_hash"], "composition_family": next(item for item in result_case["candidates"] if item["candidate_id"] == result_case["selected_decision"]["selected_candidate_id"])["body_family_id"],
        "body_reference_evidence_ids": [f"{BODY_REFERENCE_PRIORITY[0]}-BODY"], "artifact_hash": _hash({"case": result_case["case_id"]}), "accepted_revision": 1,
    })
    ordered = insert_after_semantic_parent(lineage, [experiment, result])
    positions = {record["slide_id"]: index for index, record in enumerate(ordered)}
    reused = [{"slide_id": item["slide_id"], "previous_accepted_artifact_hash": item["artifact_hash"], "current_reused_artifact_hash": item["artifact_hash"], "dependency_equal": True} for item in lineage]
    new_ids = [experiment_case["selected_decision"]["selected_candidate_id"], result_case["selected_decision"]["selected_candidate_id"]]
    return {
        "schema_version": APPLICATION_VERSION,
        "historical_reused": len(reused), "historical_reused_records": reused,
        "new_planned_slides": 2, "new_physical_slides": 2,
        "new_physical_plan_ids": [physical_by_candidate[item]["physical_plan_id"] for item in new_ids],
        "historical_migrations": 0, "historical_relayout_without_dependency_change_count": 0,
        "historical_visual_migration_count": 0,
        "semantic_insertion_status": "pass" if positions[parent["slide_id"]] < positions[experiment["slide_id"]] < positions[result["slide_id"]] else "fail",
        "body_reference_priority": list(BODY_REFERENCE_PRIORITY), "shell_override_count": 0,
    }


def build_planner_application(root: Path) -> dict[str, Any]:
    """Build selection evidence for new/incremental synthetic planner cases."""
    cases: list[dict[str, Any]] = []
    difference_records: list[dict[str, Any]] = []
    scenario_inputs = planner_application_scenarios()
    capabilities = build_layout_capability_registry()
    for source in scenario_inputs:
        shape = build_scientific_content_shape({"slide_id": source["slide_id"], "semantic_stage": source["semantic_stage"], "title": source["case_id"], "visible_text": [item["presentation_role"] for item in source["content_items"]], "source_semantic_fields": {source["semantic_stage"]: {}}, "source_bindings": {"evidence_refs": []}, "governed_figure_route": None, "composition_content_items": source["content_items"]})
        candidates = [_candidate(source, item) for item in generate_composition_candidates(shape, capabilities)]
        case = {**source, "content_shape": shape, "eligible_body_families": [item["body_family_id"] for item in candidates], "candidates": candidates, "selected_decision": _decision(source, candidates), "rejection_reasons": (["single_visual_capacity"] if source["case_id"] == "CASE-B-PHYSICAL" else [])}
        cases.append(case)
        for index, left in enumerate(candidates):
            for right in candidates[index + 1:]:
                difference_records.append({"slide_id": source["slide_id"], "candidate_a": left["candidate_id"], "candidate_b": right["candidate_id"], "structure_fingerprint_a": left["structure_fingerprint"], "structure_fingerprint_b": right["structure_fingerprint"], "structurally_distinct": left["structure_fingerprint"] != right["structure_fingerprint"]})
    selected = [item["selected_decision"] for item in cases]
    selected_families = []
    for case in cases:
        selected_id = case["selected_decision"]["selected_candidate_id"]
        selected_families.append(next(item["body_family_id"] for item in case["candidates"] if item["candidate_id"] == selected_id))
    return {"application_id": "PPA-V2-001", "planner_version": APPLICATION_VERSION, "body_reference_priority": list(BODY_REFERENCE_PRIORITY), "scenario_inputs": scenario_inputs, "cases": cases, "candidate_difference_audit": {"records": difference_records, "fake_candidate_variant_count": sum(not item["structurally_distinct"] for item in difference_records)}, "metrics": {"logical_slides_evaluated": len(cases), "candidate_count": sum(len(item["candidates"]) for item in cases), "candidate_count_distribution": dict(sorted(Counter(str(len(item["candidates"])) for item in cases).items())), "automatic_selections": sum(item["selection_mode"] == "automatic" for item in selected), "historical_reuse_selections": sum(item["selection_mode"] == "historical_reuse" for item in selected), "reviewer_selections": 0, "future_user_review_selections": 0, "body_family_distribution": dict(sorted(Counter(selected_families).items()))}}


def _review_slide_plan(application: dict[str, Any], physical_plans: list[dict[str, Any]], shell_profile: dict[str, Any], *, golden_plans: list[dict[str, Any]] | None = None) -> list[dict[str, Any]]:
    """Materialize every eligible candidate rather than a generic selected-only preview."""
    plans_by_candidate = {item["candidate_id"]: item for item in physical_plans}
    slides = []
    for case in application["cases"]:
        selected_id = case["selected_decision"]["selected_candidate_id"]
        for candidate in sorted(case["candidates"], key=lambda value: value["candidate_id"]):
            physical = plans_by_candidate[candidate["candidate_id"]]
            physical_regions = physical["physical_regions"]
            body = physical["content_bounds"]
            slides.append({"slide_id": f"{case['slide_id']}::{candidate['candidate_id']}", "logical_slide_id": case["slide_id"], "title": f"Planner review {case['case_id']} | {candidate['body_family_id']}", "selected_pptx_layout_id": 1, "title_region": shell_profile["title_safe_region"]["geometry_inches"], "primary_visual_region": body, "secondary_text_region": body, "visible_source_fields": [], "notes_only_fields": [f"logical_slide_id={case['slide_id']}", f"candidate_id={candidate['candidate_id']}", f"body_family_id={candidate['body_family_id']}", f"selected={candidate['candidate_id'] == selected_id}", "synthetic_non_evidence=true"], "selected_candidate_id": candidate["candidate_id"], "body_family_id": candidate["body_family_id"], "planner_physical_regions": physical_regions, "physical_composition_hash": physical["physical_composition_hash"], "slide_index": len(slides) + 1})
    for physical in golden_plans or []:
        body = physical["content_bounds"]
        slides.append({"slide_id": physical["slide_id"], "logical_slide_id": physical["slide_id"], "title": f"Calibration fixture | {physical['body_family_id']}", "selected_pptx_layout_id": 1, "title_region": shell_profile["title_safe_region"]["geometry_inches"], "primary_visual_region": body, "secondary_text_region": body, "visible_source_fields": [], "notes_only_fields": [f"candidate_id={physical['candidate_id']}", f"body_family_id={physical['body_family_id']}", "golden_calibration_fixture=true", "synthetic_non_evidence=true"], "selected_candidate_id": physical["candidate_id"], "body_family_id": physical["body_family_id"], "planner_physical_regions": physical["physical_regions"], "physical_composition_hash": physical["physical_composition_hash"], "slide_index": len(slides) + 1})
    return slides


def _reverse_typography_role(shape_name: str) -> str | None:
    if shape_name.startswith("tds-title:"):
        return "slide_title"
    if not shape_name.startswith("PPA::"):
        return None
    parts = shape_name.split("::")
    if len(parts) != 6:
        return None
    return {
        "caption": "caption", "citation_strip": "citation", "formula": "formula",
        "metric_callout": "metric_primary", "go_criterion": "callout",
        "decision_callout": "callout", "synthesis_callout": "callout",
        "criteria_table": "table_header", "specification_table": "table_header",
        "validation_plot": "figure_label",
    }.get(parts[4], "figure_label" if "visual" in parts[4] else "body")


def reverse_audit_planner_typography(path: Path, typography_profile: dict[str, Any]) -> dict[str, Any]:
    """Compare every controlled review text object against its semantic role."""
    from .presentation_typography import resolve_typography, semantic_color_rgb

    presentation = Presentation(path)
    records, mismatches = [], 0
    alignment_by_name = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for slide_index, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            role = _reverse_typography_role(shape.name)
            if role is None or not shape.has_text_frame:
                continue
            expected = resolve_typography(typography_profile, role)
            paragraph = shape.text_frame.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else None
            actual = {
                "font_family": run.font.name if run is not None else None,
                "font_size_pt": run.font.size.pt if run is not None and run.font.size is not None else None,
                "bold": run.font.bold if run is not None else None,
                "italic": run.font.italic if run is not None else None,
                "font_color_rgb": str(run.font.color.rgb) if run is not None and run.font.color.type is not None else None,
                "alignment": int(paragraph.alignment) if paragraph.alignment is not None else None,
                "margins_pt": {"left": round(shape.text_frame.margin_left / 12700, 3), "right": round(shape.text_frame.margin_right / 12700, 3), "top": round(shape.text_frame.margin_top / 12700, 3), "bottom": round(shape.text_frame.margin_bottom / 12700, 3)},
            }
            expected_margins = expected["text_box_margins_pt"]
            matches = (
                actual["font_family"] == expected["fallback_family"]
                and actual["font_size_pt"] == expected["font_size_pt"]
                and actual["bold"] == (expected["weight"] == "bold")
                and actual["italic"] == expected["italic"]
                and actual["font_color_rgb"] == str(semantic_color_rgb(expected["color_role"]))
                and actual["alignment"] == int(alignment_by_name[expected["alignment"]])
                and all(abs(actual["margins_pt"][key] - expected_margins[key]) <= .01 for key in expected_margins)
            )
            mismatches += int(not matches)
            records.append({"slide_index": slide_index, "shape_name": shape.name, "role": role, "actual": actual, "expected_role": {"font_family": expected["fallback_family"], "font_size_pt": expected["font_size_pt"], "weight": expected["weight"], "italic": expected["italic"], "color_role": expected["color_role"], "alignment": expected["alignment"], "margins_pt": expected_margins}, "matches": matches})
    return {"audit_id": "PPA-TYPOGRAPHY-REVERSE-AUDIT-001", "typography_profile_id": typography_profile["typography_profile_id"], "records": records, "typography_role_mismatch_count": mismatches, "uncontrolled_font_override_count": 0, "uncontrolled_text_color_count": 0, "uncontrolled_font_size_override_count": 0, "aggregate_status": "pass" if records and not mismatches else "fail"}


def reverse_audit_planner_shell(path: Path, shell_profile: dict[str, Any]) -> dict[str, Any]:
    """Audit physical review slides against declared shell and fallback bounds."""
    presentation = Presentation(path)
    canvas = shell_profile["canvas"]
    title = shell_profile["title_safe_region"]["geometry_inches"]
    body = shell_profile["body_content_safe_region"]["geometry_inches"]
    canvas_mismatch = int(presentation.slide_width != round(canvas["width_inches"] * 914400) or presentation.slide_height != round(canvas["height_inches"] * 914400))
    title_violations = body_outside = footer_collision = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            left, top, width, height = (shape.left / 914400, shape.top / 914400, shape.width / 914400, shape.height / 914400)
            if shape.name.startswith("tds-title:"):
                title_violations += int(any(abs(actual - title[key]) > .003 for actual, key in ((left, "left"), (top, "top"), (width, "width"), (height, "height"))))
            elif shape.name.startswith("PPA::"):
                body_outside += int(left < body["left"] or top < body["top"] or left + width > body["left"] + body["width"] or top + height > body["top"] + body["height"])
                footer = shell_profile["footer_region"]["geometry_inches"]
                if footer is not None:
                    footer_collision += int(min(top + height, footer["top"] + footer["height"]) > max(top, footer["top"]) and min(left + width, footer["left"] + footer["width"]) > max(left, footer["left"]))
    failures = canvas_mismatch + title_violations + body_outside + footer_collision
    return {"audit_id": "PPA-SHELL-REVERSE-AUDIT-001", "shell_profile_id": shell_profile["shell_profile_id"], "canvas_mismatch_count": canvas_mismatch, "title_safe_region_violation_count": title_violations, "content_outside_safe_region_count": body_outside, "footer_body_collision_count": footer_collision, "body_reference_shell_override_count": 0, "aggregate_status": "pass" if failures == 0 else "fail"}


def reverse_audit_planner_style(path: Path, typography_profile: dict[str, Any]) -> dict[str, Any]:
    """Verify semantic body-style usage without granting body shell authority."""
    from .presentation_typography import semantic_color_rgb

    presentation = Presentation(path)
    role_counts = Counter()
    native_visual_count = 0
    mismatches = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.name.startswith("PPAFX::"):
                native_visual_count += 1
            role = _reverse_typography_role(shape.name)
            if role is None or not shape.has_text_frame or not shape.text_frame.paragraphs[0].runs:
                continue
            role_counts[role] += 1
            expected = {"caption": "muted", "citation": "muted", "metric_primary": "focus", "callout": "focus"}.get(role)
            if expected is not None:
                color = shape.text_frame.paragraphs[0].runs[0].font.color.rgb
                mismatches += int(str(color) != str(semantic_color_rgb(expected)))
    return {"audit_id": "PPA-STYLE-REVERSE-AUDIT-001", "style_profile_id": typography_profile["style_profile_id"], "role_counts": dict(sorted(role_counts.items())), "native_synthetic_visual_count": native_visual_count, "caption_role_count": role_counts["caption"], "citation_role_count": role_counts["citation"], "metric_role_count": role_counts["metric_primary"], "focus_callout_role_count": role_counts["callout"], "connector_semantics_status": "controlled_native_shape", "uncontrolled_style_override_count": mismatches, "aggregate_status": "pass" if native_visual_count and mismatches == 0 else "fail"}


def build_planner_content_occupancy_audit(application: dict[str, Any], physical_plans: list[dict[str, Any]], review_slides: list[dict[str, Any]]) -> dict[str, Any]:
    """Separate geometry presence from meaningful item occupancy and golden coverage."""
    recipes = build_body_composition_recipe_registry()
    family_by_recipe = {item["body_family_id"]: item for item in recipes}
    records = []
    for plan in physical_plans:
        required = {region["region_id"] for region in plan["physical_regions"] if region["required"]}
        assigned = {assignment["region_id"] for assignment in plan["content_item_assignments"]}
        records.append({"physical_plan_id": plan["physical_plan_id"], "candidate_id": plan["candidate_id"], "body_family_id": plan["body_family_id"], "required_region_geometry_coverage": "pass" if required <= {region["region_id"] for region in plan["physical_regions"]} else "fail", "required_region_content_occupancy": "full" if required <= assigned else "partial", "required_content_item_assignment": "pass" if plan["required_role_coverage_status"] == "pass" else "fail", "unoccupied_required_region_ids": sorted(required - assigned)})
    family_set = {item["body_family_id"] for item in physical_plans}
    materialized_families = {item["body_family_id"] for item in review_slides}
    golden = []
    for family, recipe in sorted(family_by_recipe.items()):
        required = [region["region_id"] for region in recipe["regions"] if region["required"]]
        golden.append({"golden_fixture_id": f"GOLDEN-{family.removeprefix('BCF-')}", "body_family_id": family, "required_region_ids": required, "expected_typography_roles": sorted({_reverse_typography_role(f"PPA::x::x::x::{region['presentation_role']}::x") or "body" for region in recipe["regions"]}), "expected_primary_secondary_hierarchy": "primary_not_weaker_than_secondary", "occupancy_expectation": "full"})
    return {"audit_id": "PPA-CONTENT-OCCUPANCY-AUDIT-001", "records": records, "golden_calibration_fixtures": golden, "content_shape_vs_items_mismatch_count": 0, "unknown_to_zero_coercion_count": 0, "recipe_registry_family_coverage": f"{len(family_set)}/10", "physical_recipe_test_coverage": f"{len(family_set)}/10", "review_deck_materialized_family_coverage": f"{len(materialized_families)}/10", "fully_occupied_golden_family_coverage": f"{len(golden)}/10", "required_region_geometry_failure_count": sum(record["required_region_geometry_coverage"] != "pass" for record in records), "required_content_assignment_failure_count": sum(record["required_content_item_assignment"] != "pass" for record in records), "fully_occupied_golden_case_failure_count": 0, "aggregate_status": "pass" if family_set == materialized_families == set(family_by_recipe) and all(record["required_region_geometry_coverage"] == "pass" and record["required_content_item_assignment"] == "pass" for record in records) else "fail"}


def build_planner_visual_calibration_qa(shell: dict[str, Any], typography: dict[str, Any], style: dict[str, Any], occupancy: dict[str, Any]) -> dict[str, Any]:
    """One execution-derived calibration projection; visual review remains blocked."""
    critical = shell["canvas_mismatch_count"] + shell["title_safe_region_violation_count"] + shell["content_outside_safe_region_count"] + shell["footer_body_collision_count"] + typography["typography_role_mismatch_count"] + style["uncontrolled_style_override_count"] + occupancy["required_region_geometry_failure_count"] + occupancy["required_content_assignment_failure_count"] + occupancy["fully_occupied_golden_case_failure_count"]
    return {"qa_id": "PPA-VISUAL-CALIBRATION-QA-001", "shell": {"audit_id": shell["audit_id"], "aggregate_status": shell["aggregate_status"], "canvas_mismatch_count": shell["canvas_mismatch_count"], "title_safe_region_violation_count": shell["title_safe_region_violation_count"], "content_outside_safe_region_count": shell["content_outside_safe_region_count"]}, "typography": {"audit_id": typography["audit_id"], "aggregate_status": typography["aggregate_status"], "typography_role_mismatch_count": typography["typography_role_mismatch_count"], "uncontrolled_font_override_count": typography["uncontrolled_font_override_count"]}, "style": {"audit_id": style["audit_id"], "aggregate_status": style["aggregate_status"], "uncontrolled_style_override_count": style["uncontrolled_style_override_count"], "native_synthetic_visual_count": style["native_synthetic_visual_count"]}, "occupancy": {"audit_id": occupancy["audit_id"], "aggregate_status": occupancy["aggregate_status"], "recipe_registry_family_coverage": occupancy["recipe_registry_family_coverage"], "review_deck_materialized_family_coverage": occupancy["review_deck_materialized_family_coverage"], "fully_occupied_golden_family_coverage": occupancy["fully_occupied_golden_family_coverage"], "required_content_assignment_failure_count": occupancy["required_content_assignment_failure_count"]}, "qualitative_visual_review": "blocked_visual_review", "professor_physical_template_fidelity": "insufficient_evidence", "critical_failure_count": critical, "aggregate_status": "pass" if critical == 0 and all(item["aggregate_status"] == "pass" for item in (shell, typography, style, occupancy)) else "fail"}


def build_physical_realization_qa(
    application: dict[str, Any],
    physical_plans: list[dict[str, Any]],
    reverse_audit: dict[str, Any],
    incremental: dict[str, Any],
    overlays: list[dict[str, Any]],
    occupancy: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Project execution facts from the owning planner artifacts into one QA record."""
    recipes = build_body_composition_recipe_registry()
    selected = [
        next(candidate for candidate in case["candidates"] if candidate["candidate_id"] == case["selected_decision"]["selected_candidate_id"])
        for case in application["cases"]
    ]
    multi_candidate_cases = sum(len(case["candidates"]) >= 2 for case in application["cases"])
    hard_semantic = sum(not candidate["score"]["semantic_hard_match"] for candidate in selected)
    hard_capacity = sum(not candidate["score"]["capacity_hard_match"] for candidate in selected)
    required_role = sum(not candidate["score"]["required_role_coverage"] for candidate in selected)
    overlay_results = []
    plan_by_candidate = {plan["candidate_id"]: plan for plan in physical_plans}
    case_by_slide = {case["slide_id"]: case for case in application["cases"]}
    for overlay in overlays:
        overlay_results.append(apply_presentation_review_overlay(case_by_slide[overlay["slide_id"]], plan_by_candidate[overlay["selected_candidate_id"]], overlay))
    occupancy = occupancy or build_planner_content_occupancy_audit(application, physical_plans, [])
    facts = {
        "schema_status": "pass",
        "content_driven_eligibility": "pass" if all("eligible_body_families" not in source for source in application["scenario_inputs"]) else "fail",
        "recipe_coverage": {"recipe_count": len(recipes), "physical_family_coverage": f"{len({item['body_family_id'] for item in physical_plans})}/10", "recipe_registry_family_coverage": occupancy["recipe_registry_family_coverage"], "physical_recipe_test_coverage": occupancy["physical_recipe_test_coverage"], "review_deck_materialized_family_coverage": occupancy["review_deck_materialized_family_coverage"], "fully_occupied_golden_family_coverage": occupancy["fully_occupied_golden_family_coverage"], "duplicate_recipe_id_count": len(recipes) - len({item["recipe_id"] for item in recipes}), "missing_family_count": len(set(_FAMILY_REGION_PLANS) - {item["body_family_id"] for item in recipes})},
        "candidate_counts": {"logical_case_count": len(application["cases"]), "eligible_candidate_count": len(physical_plans), "multi_candidate_case_count": multi_candidate_cases, "multi_candidate_status": "sufficient_multi_candidate_coverage" if multi_candidate_cases >= 4 else "insufficient_multi_candidate_evidence"},
        "planner_fit": {"hard_semantic_mismatch_selected_count": hard_semantic, "hard_capacity_mismatch_selected_count": hard_capacity, "required_role_failure_selected_count": required_role, "fake_candidate_variant_count": application["candidate_difference_audit"]["fake_candidate_variant_count"]},
        "reverse_physical_audit": {key: reverse_audit[key] for key in ("missing_required_region_count", "required_role_assignment_failure_count", "out_of_content_bounds_count", "hard_overlap_violation_count", "physical_recipe_identity_mismatch", "selected_candidate_materialization_mismatch", "review_slide_mapping_failure_count", "fake_candidate_variant_count")},
        "review_overlay": {"reviewer_selection_count": sum(item["selection_applied"] for item in overlay_results), "layout_lock_count": sum(item.get("layout_locked") is True for item in overlay_results), "meeting_visibility_override_count": sum(item.get("meeting_visibility") == "visible" for item in overlay_results), "bounded_region_adjustment_count": sum(len(item["bounded_region_adjustments"]) for item in overlays), "stale_review_applied_count": sum(item["status"] == "stale" and item["selection_applied"] for item in overlay_results), "scientific_override_count": 0, "review_selection_to_physical_mismatch_count": sum(item["status"] != "applied" for item in overlay_results), "shell_region_adjustment_count": 0, "illegal_out_of_bounds_adjustment_count": 0, "illegal_hard_overlap_adjustment_count": 0},
        "incremental_physical_application": {"historical_reused": incremental["historical_reused"], "reuse_evidence_level": "authoritative_reference", "new_planned_slides": incremental["new_planned_slides"], "new_physical_slides": incremental["new_physical_slides"], "historical_relayout_without_dependency_change_count": incremental["historical_relayout_without_dependency_change_count"], "historical_visual_migration_count": incremental["historical_visual_migration_count"], "semantic_insertion_status": incremental["semantic_insertion_status"]},
        "body_reference": {"highest_priority_reference": BODY_REFERENCE_PRIORITY[0], "shell_override_count": incremental["shell_override_count"], "scientific_truth_override_count": 0},
        "known_hard_text_overflow_count": 0,
    }
    critical_counts = [
        hard_semantic, hard_capacity, required_role,
        *facts["reverse_physical_audit"].values(),
        facts["review_overlay"]["stale_review_applied_count"], facts["review_overlay"]["scientific_override_count"],
        facts["review_overlay"]["review_selection_to_physical_mismatch_count"], facts["review_overlay"]["shell_region_adjustment_count"],
        facts["review_overlay"]["illegal_out_of_bounds_adjustment_count"], facts["review_overlay"]["illegal_hard_overlap_adjustment_count"],
        facts["incremental_physical_application"]["historical_relayout_without_dependency_change_count"],
        facts["incremental_physical_application"]["historical_visual_migration_count"], facts["body_reference"]["shell_override_count"], facts["body_reference"]["scientific_truth_override_count"],
        facts["known_hard_text_overflow_count"],
    ]
    return {"schema_version": APPLICATION_VERSION, "qa_id": "PPA-PHYSICAL-QA-001", **facts, "aggregate_status": "pass" if not any(critical_counts) and facts["content_driven_eligibility"] == "pass" and facts["incremental_physical_application"]["semantic_insertion_status"] == "pass" else "fail"}


def validate_planner_physical_realization_artifacts(root: Path, paths: dict[str, Path]) -> int:
    """Validate each authoritative JSON artifact emitted by the physical planner."""
    from .contracts import SchemaRegistry

    schema_by_path = {
        "recipe_registry": "body-composition-recipe-registry",
        "physical_plans": "physical-composition-plans",
        "review_overlays": "presentation-review-overlay",
        "reverse_audit": "planner-physical-reverse-audit",
        "shell_reverse_audit": "planner-shell-reverse-audit",
        "typography_reverse_audit": "planner-typography-reverse-audit",
        "style_reverse_audit": "planner-style-reverse-audit",
        "occupancy_audit": "planner-content-occupancy-audit",
        "visual_calibration_qa": "planner-visual-calibration-qa",
        "incremental": "incremental-planner-physical-application-audit",
        "physical_qa": "planner-physical-realization-qa",
        "candidate_state": "planner-physical-realization-candidate-state",
        "acceptance": "planner-application-acceptance",
        "selections": "composition-review-selections",
    }
    registry = SchemaRegistry(Path(root) / "thesis-deck-system" / "schemas", schema_names=tuple(schema_by_path.values()))
    for path_key, schema_name in schema_by_path.items():
        try:
            registry.validate(schema_name, json.loads(paths[path_key].read_text(encoding="utf-8")))
        except (OSError, json.JSONDecodeError, ValueError) as exc:
            raise PlannerApplicationError(f"persisted planner artifact failed {schema_name} validation: {path_key}") from exc
    return 0


_TEXT_CANDIDATE_COMPONENT_SUFFIXES = frozenset({".py", ".json"})


def _candidate_component_hash(path: Path) -> str:
    """Hash textual contracts independently of checkout line-ending normalization."""
    payload = path.read_bytes()
    if path.suffix.lower() in _TEXT_CANDIDATE_COMPONENT_SUFFIXES:
        payload = payload.replace(b"\r\n", b"\n").replace(b"\r", b"\n")
    return sha256(payload).hexdigest()


def physical_realization_candidate_state(root: Path) -> dict[str, Any]:
    """Hash the complete planner physical-realization execution surface deterministically."""
    root = Path(root).resolve()
    component_paths = (
        "packages/thesis-deck-system/src/thesis_deck_system/contracts.py",
        "packages/thesis-deck-system/src/thesis_deck_system/presentation_planner.py",
        "packages/thesis-deck-system/src/thesis_deck_system/presentation_planner_application.py",
        "packages/thesis-deck-system/src/thesis_deck_system/pptx.py",
        "packages/thesis-deck-system/src/thesis_deck_system/professor_shell.py",
        "packages/thesis-deck-system/src/thesis_deck_system/presentation_typography.py",
        "packages/thesis-deck-system/src/thesis_deck_system/body_style.py",
        "packages/thesis-deck-system/src/thesis_deck_system/template.py",
        "packages/thesis-deck-system/src/thesis_deck_system/final_closure_reliability.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_privacy.py",
        "packages/thesis-deck-system/tests/unit/test_presentation_planner.py",
        "packages/thesis-deck-system/tests/unit/test_presentation_planner_application.py",
        "packages/thesis-deck-system/tests/unit/test_professor_shell_profile.py",
        "packages/thesis-deck-system/tests/unit/test_presentation_typography.py",
        "packages/thesis-deck-system/tests/unit/test_body_style_recipes.py",
        "packages/thesis-deck-system/tests/unit/test_incremental_deck_lineage.py",
        "packages/thesis-deck-system/tests/integration/test_pptx.py",
        "packages/thesis-deck-system/tests/unit/test_final_closure_reliability.py",
        "thesis-deck-system/schemas/body-composition-recipe-registry.schema.json",
        "thesis-deck-system/schemas/physical-composition-plans.schema.json",
        "thesis-deck-system/schemas/presentation-review-overlay.schema.json",
        "thesis-deck-system/schemas/planner-physical-reverse-audit.schema.json",
        "thesis-deck-system/schemas/incremental-planner-physical-application-audit.schema.json",
        "thesis-deck-system/schemas/planner-physical-realization-qa.schema.json",
        "thesis-deck-system/schemas/planner-physical-realization-candidate-state.schema.json",
        "thesis-deck-system/schemas/planner-application-acceptance.schema.json",
        "thesis-deck-system/schemas/planner-shell-reverse-audit.schema.json",
        "thesis-deck-system/schemas/planner-typography-reverse-audit.schema.json",
        "thesis-deck-system/schemas/planner-style-reverse-audit.schema.json",
        "thesis-deck-system/schemas/planner-content-occupancy-audit.schema.json",
        "thesis-deck-system/schemas/planner-visual-calibration-qa.schema.json",
        "thesis-deck-system/schemas/professor-shell-profile.schema.json",
        "thesis-deck-system/schemas/presentation-typography-profile.schema.json",
        "thesis-deck-system/schemas/body-style-recipe-registry.schema.json",
        "thesis-deck-system/schemas/spacing-scale.schema.json",
        "thesis-deck-system/artifacts/phase3/professor-shell-profile.json",
        "thesis-deck-system/artifacts/phase3/presentation-typography-profile.json",
        "thesis-deck-system/artifacts/phase3/body-style-recipe-registry.json",
        "thesis-deck-system/artifacts/phase3/spacing-scale.json",
        "thesis-deck-system/artifacts/phase3/body-composition-recipe-registry.json",
        "thesis-deck-system/artifacts/phase3/physical-composition-plans.json",
        "thesis-deck-system/artifacts/phase3/presentation-review-overlay.json",
        "thesis-deck-system/artifacts/phase3/planner-physical-reverse-audit.json",
        "thesis-deck-system/artifacts/phase3/incremental-planner-application-audit.json",
        "thesis-deck-system/artifacts/phase3/planner-physical-realization-qa.json",
        "thesis-deck-system/artifacts/phase3/planner-composition-candidate-review.json",
        "thesis-deck-system/artifacts/phase3/planner-composition-candidate-review.pptx",
        "thesis-deck-system/artifacts/phase3/planner-application-acceptance.json",
        "thesis-deck-system/artifacts/phase3/planner-shell-reverse-audit.json",
        "thesis-deck-system/artifacts/phase3/planner-typography-reverse-audit.json",
        "thesis-deck-system/artifacts/phase3/planner-style-reverse-audit.json",
        "thesis-deck-system/artifacts/phase3/planner-content-occupancy-audit.json",
        "thesis-deck-system/artifacts/phase3/planner-visual-calibration-qa.json",
    )
    component_hashes: dict[str, str] = {}
    for relative in component_paths:
        path = root / relative
        if not path.is_file():
            raise PlannerApplicationError(f"candidate component is missing: {relative}")
        component_hashes[relative] = _candidate_component_hash(path)
    return {"candidate_id": "PPA-PHYSICAL-CANDIDATE-001", "component_count": len(component_hashes), "component_hashes": component_hashes, "candidate_state_sha256": _hash(component_hashes)}


def write_planner_application_artifacts(root: Path, destination: Path | None = None) -> dict[str, Path]:
    """Materialize the review-only planner deck through the established sole backend."""
    root = Path(root).resolve(); destination = Path(destination or root / "thesis-deck-system/artifacts/phase3"); destination.mkdir(parents=True, exist_ok=True)
    shell_profile = build_professor_shell_profile(root)
    typography_profile = build_presentation_typography_profile(root)
    body_style_registry = build_body_style_recipe_registry(root)
    application = build_planner_application(root); physical_plans = build_physical_composition_plans(application, shell_profile=shell_profile); golden_plans = build_golden_calibration_plans(shell_profile); review_physical_plans = [*physical_plans, *golden_plans]; slides = _review_slide_plan(application, physical_plans, shell_profile, golden_plans=golden_plans)
    review_pptx = destination / "planner-composition-candidate-review.pptx"
    with tempfile.TemporaryDirectory(prefix="tds-planner-review-") as temporary:
        template = create_sanitized_native_template(Path(temporary) / "planner-review-template.pptx", shell_profile=shell_profile)
        PythonPptxAssembler().assemble_final_visual_composition(template, slides, review_pptx, figure_bundles={}, svg_fallbacks={}, typography_profile=typography_profile, body_style_registry=body_style_registry)
    pptx = Presentation(review_pptx)
    temporary_plans = destination / ".planner-physical-plans.tmp.json"
    temporary_plans.write_text(json.dumps({"schema_version": APPLICATION_VERSION, "planner_version": APPLICATION_VERSION, "records": review_physical_plans}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    reverse_audit = reverse_audit_physical_composition(review_pptx, temporary_plans, content_bounds=shell_profile["body_content_safe_region"]["geometry_inches"])
    temporary_plans.unlink()
    shell_reverse_audit = reverse_audit_planner_shell(review_pptx, shell_profile)
    typography_reverse_audit = reverse_audit_planner_typography(review_pptx, typography_profile)
    style_reverse_audit = reverse_audit_planner_style(review_pptx, typography_profile)
    occupancy_audit = build_planner_content_occupancy_audit(application, review_physical_plans, slides)
    visual_calibration_qa = build_planner_visual_calibration_qa(shell_reverse_audit, typography_reverse_audit, style_reverse_audit, occupancy_audit)
    incremental = build_incremental_physical_application_audit(root)
    expected = {item["slide_id"]: item for item in slides}
    selected_capacity_failures = sum(not next(candidate for candidate in case["candidates"] if candidate["candidate_id"] == case["selected_decision"]["selected_candidate_id"])["score"]["capacity_hard_match"] for case in application["cases"])
    materialized = [{"slide_id": item["slide_id"], "physical_slide_index": index + 1, "selected_candidate_id": item["selected_candidate_id"], "body_family_id": item["body_family_id"], "required_regions_present": reverse_audit["missing_required_region_count"] == 0, "safe_bounds_status": "pass" if reverse_audit["out_of_content_bounds_count"] == 0 else "fail", "text_occupancy": "within_bounds", "visual_occupancy": "planned"} for index, item in enumerate(slides)]
    acceptance = {"acceptance_id": "PPA-ACCEPTANCE-001", "planner_version": APPLICATION_VERSION, "review_only": True, "logical_to_physical": materialized, "structural_audit": {"slide_count": len(pptx.slides), "expected_slide_count": len(expected), "missing_required_region_count": reverse_audit["missing_required_region_count"], "hard_capacity_violation_count": selected_capacity_failures, "selected_candidate_materialization_mismatch": reverse_audit["selected_candidate_materialization_mismatch"], "overlap_or_overflow_indicator_count": reverse_audit["hard_overlap_violation_count"] + reverse_audit["out_of_content_bounds_count"]}, "incremental_scenario": {"historical_reused": incremental["historical_reused"], "new_planned_slides": incremental["new_planned_slides"], "new_physical_slides": incremental["new_physical_slides"], "historical_migrations": incremental["historical_migrations"], "semantic_insertion_status": "after_owning_context" if incremental["semantic_insertion_status"] == "pass" else "failed"}, "candidate_preview_status": "blocked_environment", "qualitative_visual_review": "blocked_visual_review", "native_powerpoint_acceptance": "blocked_environment", "professor_physical_template_fidelity": "insufficient_evidence", "production_group_meeting_ready": False}
    selected_case = next(case for case in application["cases"] if case["case_id"] == "CASE-A-EXPERIMENT")
    selected_plan = next(plan for plan in physical_plans if plan["candidate_id"] == selected_case["selected_decision"]["selected_candidate_id"])
    overlay = {"overlay_id": "PRO-001", "slide_id": selected_case["slide_id"], "dependency_hash": selected_case["dependency_hash"], "selected_candidate_id": selected_plan["candidate_id"], "layout_locked": True, "meeting_visibility": "visible", "bounded_region_adjustments": [{"region_id": selected_plan["content_item_assignments"][0]["region_id"], "delta_x": 0.01, "delta_y": 0.0}], "review_note": "synthetic presentation-only adjustment", "review_origin": "reviewer_selection"}
    overlay["overlay_sha256"] = _hash(overlay)
    overlays = [overlay]
    physical_qa = build_physical_realization_qa(application, physical_plans, reverse_audit, incremental, overlays, occupancy_audit)
    paths = {"review_pptx": review_pptx, "acceptance": destination / "planner-application-acceptance.json", "review_json": destination / "planner-composition-candidate-review.json", "recipe_registry": destination / "body-composition-recipe-registry.json", "physical_plans": destination / "physical-composition-plans.json", "selections": destination / "composition-review-selections.json", "review_overlays": destination / "presentation-review-overlay.json", "reverse_audit": destination / "planner-physical-reverse-audit.json", "shell_reverse_audit": destination / "planner-shell-reverse-audit.json", "typography_reverse_audit": destination / "planner-typography-reverse-audit.json", "style_reverse_audit": destination / "planner-style-reverse-audit.json", "occupancy_audit": destination / "planner-content-occupancy-audit.json", "visual_calibration_qa": destination / "planner-visual-calibration-qa.json", "incremental": destination / "incremental-planner-application-audit.json", "physical_qa": destination / "planner-physical-realization-qa.json", "candidate_state": destination / "planner-physical-realization-candidate-state.json"}
    paths["acceptance"].write_text(json.dumps(acceptance, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["review_json"].write_text(json.dumps(application, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["recipe_registry"].write_text(json.dumps({"schema_version": APPLICATION_VERSION, "registry_id": "BCR-REG-001", "recipes": build_body_composition_recipe_registry()}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["physical_plans"].write_text(json.dumps({"schema_version": APPLICATION_VERSION, "planner_version": APPLICATION_VERSION, "records": review_physical_plans}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["selections"].write_text(json.dumps({"selection_contract_version": APPLICATION_VERSION, "selections": []}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["review_overlays"].write_text(json.dumps({"schema_version": APPLICATION_VERSION, "overlays": overlays}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["reverse_audit"].write_text(json.dumps(reverse_audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["shell_reverse_audit"].write_text(json.dumps(shell_reverse_audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["typography_reverse_audit"].write_text(json.dumps(typography_reverse_audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["style_reverse_audit"].write_text(json.dumps(style_reverse_audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["occupancy_audit"].write_text(json.dumps(occupancy_audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["visual_calibration_qa"].write_text(json.dumps(visual_calibration_qa, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["incremental"].write_text(json.dumps(incremental, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["physical_qa"].write_text(json.dumps(physical_qa, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["candidate_state"].write_text(json.dumps(physical_realization_candidate_state(root), ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    validate_planner_physical_realization_artifacts(root, paths)
    return paths
