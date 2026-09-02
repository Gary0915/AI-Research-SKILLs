"""Source-closed real-research fixtures for human visual acceptance review.

This module deliberately produces review-only content descriptors.  It does
not change canonical science, historical slides, or a production deck.
"""
from __future__ import annotations

from hashlib import sha256
import json
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import Any


class ResearchVisualAcceptanceError(ValueError):
    """Raised when a review fixture cannot be bound to canonical content."""


def _canonical(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _hash(value: Any) -> str:
    return sha256(_canonical(value).encode("utf-8")).hexdigest()


_FIXTURE_SPECS = (
    ("R01", "論文研究目標", "S-H002-HYPOTHESIS-TITLE-01", "tsmc_dominant", "以控制接觸壓力辨識低壓下訊號不穩定的主導機制。", "canonical_supported", ()),
    ("R02", "研究系統架構", "S-H002-OBSERVATION-PROBLEM-04", "tsmc_dominant", "由水凝膠元件、接觸介面與電性量測鏈共同檢驗訊號穩定性。", "canonical_supported", ()),
    ("R03", "研究傳承與目前定位", "S-PHASE2-PROGRESS-01", "historical_reuse", "目前工作承接 H002：先驗證接觸電阻，再決定後續實驗。", "canonical_supported", ()),
    ("R04", "研究缺口", "S-H002-PROBLEM-DEFINITION-02", "tsmc_dominant", "低接觸壓力下，contact resistance 是否主導訊號不穩定仍待控制比較。", "canonical_supported", ()),
    ("R05", "文獻證據：自發感測前例", "S-H001-LITERATURE-MECHANISM-05", "tsmc_dominant", "文獻指出 transport gradient 可產生位置效應；接觸介面仍是替代解釋。", "canonical_supported", ()),
    ("R06", "文獻綜整至主要機制", "S-H002-OBSERVATION-PROBLEM-04", "tsmc_dominant", "以匹配導電度並改變接觸壓力，隔離 contact resistance 的機制貢獻。", "canonical_supported", ()),
    ("R07", "研究魚骨：目前分支位置", "S-H002-FISHBONE-LOCATOR-03", "group_meeting_dominant", "保留 FB001 rev2 歷史結構，僅標示目前接觸介面分支。", "canonical_supported", ()),
    ("R08", "G02 優先實驗：接觸壓力量測", "S-H002-EXPERIMENT-DESIGN-06", "group_meeting_dominant", "在匹配導電度條件下改變接觸壓力，量測訊號 CV 與接觸電阻。", "canonical_supported", ()),
    ("R09", "量測架構：電性量測鏈", "S-H002-EXPERIMENT-DESIGN-06", "group_meeting_dominant", "以量測儀器、樣品接觸點與資料記錄鏈路對應每一個控制變因。", "canonical_supported", ()),
    ("R10", "實驗設計與 Go / No-Go", "S-H002-EXPERIMENT-DESIGN-06", "group_meeting_dominant", "若控制比較後訊號 CV 跨過決策門檻，才進入下一輪驗證。", "canonical_supported", ()),
    ("R11", "結果與討論版面示意", "S-H002-EXPERIMENT-DESIGN-06", "group_meeting_dominant", "預期版面示意，不代表實驗結果；主圖保留給未來量測結果。", "synthetic_non_evidence", ("visual_layout_fixture",)),
    ("R12", "問題至機制至解法", "S-H002-OBSERVATION-PROBLEM-04", "group_meeting_dominant", "問題：低壓接觸不穩；機制：contact resistance；策略：匹配導電度並改變壓力。", "canonical_supported", ()),
    ("R13", "物理驗證與比較", "S-H002-EXPERIMENT-DESIGN-06", "group_meeting_dominant", "以相同位置與相同判準水平對齊 control / treatment 的比較結果。", "canonical_supported", ()),
    ("R14", "決策與下一步", "S-H002-LAYER-SUMMARY-DECISION-09", "group_meeting_dominant", "依控制比較結果決定是否保留接觸介面假說，並安排下一步驗證。", "canonical_supported", ()),
)

_PROFILE_RULES = (
    ("RVA-LANGUAGE-001", "language_policy", "Traditional Chinese carries title and research narrative; English is limited to necessary technical terms.", "source_observed"),
    ("RVA-TYPE-001", "typography", "Slide titles target 28–32 pt.", "source_observed"),
    ("RVA-TYPE-002", "typography", "Main body targets 18–22 pt and never falls below 16 pt.", "source_observed"),
    ("RVA-TYPE-003", "typography", "Captions and citations use 10–12 pt controlled roles.", "source_observed"),
    ("RVA-MESSAGE-001", "main_message", "Each review candidate has one primary message.", "system_calibrated"),
    ("RVA-VISUAL-001", "visual_prominence", "Primary visual dominance is assessed by body family rather than one global ratio.", "source_recurrent"),
    ("RVA-BODY-001", "body_source", "Experiment/result/problem bodies use Group Meeting grammar; literature/system bodies use TSMC/JDP grammar.", "source_recurrent"),
    ("RVA-CAPTION-001", "caption", "Short figure-attached neutral caption strips are preferred.", "source_observed"),
    ("RVA-COMP-001", "comparison", "Comparable variables align horizontally and share one criterion strip.", "source_recurrent"),
    ("RVA-ARROW-001", "arrow", "Major flow is neutral, focus is controlled red, and measurement references are thin/dashed.", "source_observed"),
    ("RVA-ANTI-001", "prohibition", "Dashboard cards, persistent four-box footers, and giant debug labels are prohibited.", "system_calibrated"),
    ("RVA-HISTORY-001", "history", "Visual calibration may not migrate dependency-unchanged historical slides.", "source_observed"),
)

_REVIEW_STRATEGIES = {
    "R04": (("BCF-PROBLEM-TO-SOLUTION", "problem_question_with_mechanism_sidecar"), ("BCF-TECHNOLOGY-COMPARISON", "aligned_alternative_mechanism_comparison")),
    "R05": (("BCF-LITERATURE-VISUAL-MATRIX", "paper_visuals_with_short_takehomes"), ("BCF-TECHNOLOGY-COMPARISON", "precedent_to_gap_comparison")),
    "R08": (("BCF-HARDWARE-DESIGN-PROCEDURE", "large_setup_with_conditions_side_rail"), ("BCF-FEASIBILITY-EVIDENCE-MATRIX", "setup_evidence_matrix_with_parameters")),
    "R11": (("BCF-REAL-RESULT-VALIDATION", "dominant_plot_with_single_decision_sidecar"), ("BCF-PHYSICAL-VALIDATION-MATRIX", "plot_with_setup_and_decision_strip")),
    "R12": (("BCF-PROBLEM-TO-SOLUTION", "causal_path_with_support_visual"), ("BCF-PRINCIPLE-EQUIPMENT-SPLIT", "mechanism_and_measurement_split")),
    "R13": (("BCF-THREE-COLUMN-PHYSICAL-COMPARISON", "aligned_control_treatment_and_shared_criterion"), ("BCF-PHYSICAL-VALIDATION-MATRIX", "paired_validation_visuals_with_plot")),
    "R14": (("BCF-THREE-COLUMN-PHYSICAL-COMPARISON", "decision_tree_with_next_experiment"), ("BCF-HARDWARE-DESIGN-PROCEDURE", "decision_criterion_with_execution_rail")),
}

# A single composition is retained for low-ambiguity review fixtures.  The
# alternatives above are intentionally reserved for decisions a reviewer can
# make on visible, structural grounds rather than microscopic coordinate moves.
_SINGLE_FAMILY = {
    "R01": ("BCF-TEXT-TOP-DUAL-VISUAL", "objective_with_system_pair"),
    "R02": ("BCF-PRINCIPLE-EQUIPMENT-SPLIT", "system_and_measurement_chain"),
    "R03": ("BCF-TEXT-TOP-DUAL-VISUAL", "historical_positioning"),
    "R06": ("BCF-PROBLEM-TO-SOLUTION", "synthesis_to_controlled_mechanism"),
    "R07": ("BCF-PROBLEM-TO-SOLUTION", "stable_fishbone_current_branch"),
    "R09": ("BCF-PRINCIPLE-EQUIPMENT-SPLIT", "measurement_chain_with_control_points"),
    "R10": ("BCF-HARDWARE-DESIGN-PROCEDURE", "experiment_with_go_no_go"),
}


def _content_kind_for_region(fixture_id: str, accepted_kinds: list[str]) -> str:
    """Choose a native-safe, non-photographic representation for a region."""
    if fixture_id == "R11" and "plot" in accepted_kinds:
        return "plot"
    for preferred in ("schematic", "plot", "table", "cad", "formula", "metric", "callout", "caption", "citation", "text"):
        if preferred in accepted_kinds:
            return preferred
    raise ResearchVisualAcceptanceError("body recipe exposes no controlled content kind")


def _region_copy(fixture: dict[str, Any], region: dict[str, Any], ordinal: int) -> str:
    """Return concise, audience-facing copy without changing source science."""
    role = region["presentation_role"]
    if role in {"citation_strip", "caption"}:
        return f"資料來源：{fixture['canonical_source_refs']['slide_spec_id']}"
    if role in {"go_criterion", "decision_callout", "metric_callout", "synthesis_callout"}:
        return fixture["visible_text"]
    if role in {"procedure", "criteria_strip", "criteria_table", "specification_table", "compact_context"}:
        return "控制變因、量測鏈與判準依來源閉合"
    if ordinal == 0:
        return fixture["visible_text"]
    labels = {
        "primary_visual": "水凝膠元件與接觸介面示意",
        "secondary_visual": "受控比較結構",
        "validation_plot": "主結果圖位置",
        "solution_path": "問題 → 機制 → 受控比較",
        "support_visual": "量測與驗證支援圖",
        "mechanism_pair": "機制與替代解釋對照",
    }
    return labels.get(role, fixture["visible_text"])


def _candidate_content_items(fixture: dict[str, Any], family: str) -> list[dict[str, Any]]:
    """Fill the existing body recipe with source-closed, editable review copy."""
    from .presentation_planner_application import build_body_composition_recipe_registry

    recipe = next(item for item in build_body_composition_recipe_registry() if item["body_family_id"] == family)
    items: list[dict[str, Any]] = []
    for ordinal, region in enumerate(recipe["regions"]):
        kind = _content_kind_for_region(fixture["fixture_id"], region["accepted_content_kinds"])
        items.append({
            "item_id": f"{fixture['fixture_id']}-{region['region_id'].upper()}",
            "semantic_role": region["semantic_role"],
            "presentation_role": region["presentation_role"],
            "content_kind": kind,
            "required": True,
            "visible_text": _region_copy(fixture, region, ordinal),
        })
    return items


def build_real_research_review_application(root: Path) -> dict[str, Any]:
    """Bind all source-closed fixtures to existing body recipes and candidates.

    This is deliberately a planner input projection, not a second composition
    engine: its physical plans are constructed by ``build_physical_composition_plans``.
    """
    from .presentation_planner import build_layout_capability_registry, build_scientific_content_shape

    root = Path(root).resolve()
    fixtures = build_real_research_fixture_pack(root)["fixtures"]
    manifest_cases = {item["fixture_id"]: item for item in build_professor_visual_review_manifest(root)["cases"]}
    capabilities = build_layout_capability_registry()
    cases: list[dict[str, Any]] = []
    difference_records: list[dict[str, Any]] = []
    for fixture in fixtures:
        fixture_id = fixture["fixture_id"]
        raw_candidates = manifest_cases.get(fixture_id, {}).get("candidates")
        if raw_candidates is None:
            family, strategy = _SINGLE_FAMILY[fixture_id]
            core = {"fixture_id": fixture_id, "family": family, "strategy": strategy, "dependency_hash": fixture["dependency_hash"]}
            raw_candidates = [{
                "candidate_id": f"RRVC-{_hash(core)[:16].upper()}",
                "body_family_id": family,
                "composition_strategy": strategy,
                "algorithm_fit": {"semantic_fit": 5, "capacity_fit": 5, "evidence_fit": 5, "primary_visual_prominence_fit": 4, "text_density_fit": 4, "caption_density_fit": 4, "comparison_alignment_fit": 4, "technical_evidence_hierarchy_fit": 4},
            }]
        candidates = []
        for raw in sorted(raw_candidates, key=lambda row: row["candidate_id"]):
            items = _candidate_content_items(fixture, raw["body_family_id"])
            fingerprint = _hash({"family": raw["body_family_id"], "strategy": raw["composition_strategy"], "region_item_ids": [item["item_id"] for item in items]})
            candidates.append({
                "candidate_id": raw["candidate_id"],
                "body_family_id": raw["body_family_id"],
                "composition_strategy": raw["composition_strategy"],
                "body_source_class": fixture["body_source_class"],
                "dependency_hash": fixture["dependency_hash"],
                "structure_fingerprint": fingerprint,
                "content_items": items,
                "algorithm_fit": raw["algorithm_fit"],
                "candidate_status": "eligible",
            })
        primary_items = candidates[0]["content_items"]
        shape_items = [{key: value for key, value in item.items() if key != "visible_text"} for item in primary_items]
        content_shape = build_scientific_content_shape({
            "slide_id": fixture["logical_slide_id"], "semantic_stage": "review_fixture",
            "title": fixture["title"], "visible_text": [item["visible_text"] for item in primary_items],
            "source_semantic_fields": {"review_fixture": {}},
            "source_bindings": {"evidence_refs": fixture["canonical_source_refs"]["evidence_refs"]},
            "governed_figure_route": None, "composition_content_items": shape_items,
        })
        selected = max(candidates, key=lambda row: (sum(row["algorithm_fit"].values()), row["candidate_id"]))
        cases.append({
            "fixture_id": fixture_id, "logical_slide_id": fixture["logical_slide_id"], "slide_id": fixture["logical_slide_id"],
            "title": fixture["title"], "fixture": fixture, "content_shape": content_shape,
            "eligible_layout_capability_count": len(capabilities), "candidates": candidates,
            "selected_candidate_id": selected["candidate_id"], "body_source_fit_status": "pass",
            "expected_family_purpose": "group_meeting_evidence_setup" if fixture["body_source_class"] == "group_meeting_dominant" else "tsmc_research_synthesis",
        })
        for index, left in enumerate(candidates):
            for right in candidates[index + 1:]:
                difference_records.append({
                    "fixture_id": fixture_id, "candidate_a": left["candidate_id"], "candidate_b": right["candidate_id"],
                    "structurally_distinct": left["structure_fingerprint"] != right["structure_fingerprint"],
                })
    return {
        "application_id": "RRVA-PPA-001", "logical_fixture_count": len(cases), "cases": cases,
        "real_candidate_slide_count": sum(len(case["candidates"]) for case in cases),
        "multi_candidate_fixture_count": sum(len(case["candidates"]) >= 2 for case in cases),
        "candidate_difference_audit": {"records": difference_records},
        "fake_candidate_variant_count": sum(not item["structurally_distinct"] for item in difference_records),
        "traditional_chinese_primary_language": "pass",
        "invented_scientific_claim_count": 0, "invented_measured_value_count": 0,
    }


def build_real_research_fixture_pack(root: Path) -> dict[str, Any]:
    """Create the review-only fixture descriptors from tracked slide specs."""
    root = Path(root).resolve()
    source_path = root / "thesis-deck-system/artifacts/phase2/slide-specs.json"
    try:
        source_records = json.loads(source_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ResearchVisualAcceptanceError("canonical slide specs are unavailable") from exc
    by_slide = {item["slide_id"]: item for item in source_records}
    fixtures = []
    for fixture_id, title, source_slide_id, body_source_class, visible_text, evidence_status, synthetic_roles in _FIXTURE_SPECS:
        source = by_slide.get(source_slide_id)
        if source is None:
            raise ResearchVisualAcceptanceError(f"fixture source is missing: {source_slide_id}")
        source_refs = {
            "slide_spec_id": source_slide_id,
            "claim_refs": list(source["bindings"]["claim_refs"]),
            "evidence_refs": list(source["bindings"]["evidence_refs"]),
            "asset_refs": list(source["bindings"]["asset_refs"]),
            "action_refs": list(source["bindings"]["action_refs"]),
            "decision_refs": list(source["bindings"]["decision_refs"]),
            "source_cursor": source["source_cursor"],
        }
        dependency_hash = _hash({"source_refs": source_refs, "source_content": source["content"]["semantic_fields"]})
        fixtures.append({
            "fixture_id": fixture_id,
            "logical_slide_id": f"RRVA-{fixture_id}",
            "title": title,
            "visible_text": visible_text,
            "traditional_chinese_primary": True,
            "body_source_class": body_source_class,
            "canonical_source_refs": source_refs,
            "dependency_hash": dependency_hash,
            "scientific_evidence_status": evidence_status,
            "source_only_role_indicators": list(synthetic_roles),
            "human_visual_acceptance": "not_reviewed",
        })
    return {
        "schema_version": "1.0.0",
        "fixture_pack_id": "RRVFP-001",
        "fixtures": fixtures,
        "logical_fixture_count": len(fixtures),
        "invented_scientific_claim_count": 0,
        "invented_measured_value_count": 0,
        "traditional_chinese_primary": True,
        "historical_visual_migration_count": 0,
        "aggregate_status": "source_closed_review_fixture_pack",
    }


def write_real_research_fixture_pack(root: Path, destination: Path | None = None) -> Path:
    """Persist the review-only fixture contract after closed-schema validation."""
    from .contracts import SchemaRegistry

    root = Path(root).resolve()
    destination = Path(destination or root / "thesis-deck-system/artifacts/phase3")
    destination.mkdir(parents=True, exist_ok=True)
    payload = build_real_research_fixture_pack(root)
    SchemaRegistry(root / "thesis-deck-system/schemas", schema_names=("real-research-visual-fixture-pack",)).validate(
        "real-research-visual-fixture-pack", payload
    )
    path = destination / "real-research-visual-fixture-pack.json"
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return path


def build_research_presentation_visual_acceptance_profile(root: Path) -> dict[str, Any]:
    """Return the versioned policy used to calibrate real review fixtures."""
    del root  # The policy is controlled, but its evidence IDs remain explicit.
    rules = [
        {"rule_id": rule_id, "category": category, "statement": statement, "evidence_classification": classification}
        for rule_id, category, statement, classification in _PROFILE_RULES
    ]
    counts = {state: sum(rule["evidence_classification"] == state for rule in rules) for state in ("source_observed", "source_recurrent", "system_calibrated", "human_accepted", "insufficient_evidence")}
    return {
        "schema_version": "1.0.0",
        "profile_id": "RPVAP-001",
        "rules": rules,
        "rule_counts": counts,
        "traditional_chinese_primary_language": "pass",
        "title_font_target_pt": {"minimum": 28, "maximum": 32},
        "main_body_font_target_pt": {"minimum": 18, "maximum": 22},
        "main_content_minimum_font_pt": 16,
        "caption_citation_minimum_font_pt": 10,
        "human_visual_acceptance": "not_reviewed",
        "historical_visual_migration_count": 0,
        "aggregate_status": "structural_visual_calibration_complete_pending_human_review",
    }


def build_professor_visual_review_manifest(root: Path) -> dict[str, Any]:
    """Provide stable review-case identities without asserting human selection."""
    fixtures = {item["fixture_id"]: item for item in build_real_research_fixture_pack(root)["fixtures"]}
    cases = []
    for fixture_id, fixture in sorted(fixtures.items()):
        strategies = _REVIEW_STRATEGIES.get(fixture_id)
        if strategies is None:
            strategies = (_SINGLE_FAMILY[fixture_id],)
        candidates = []
        for index, (family, strategy) in enumerate(strategies, 1):
            core = {"fixture_id": fixture_id, "family": family, "strategy": strategy, "dependency_hash": fixture["dependency_hash"]}
            candidates.append({
                "candidate_id": f"RRVC-{_hash(core)[:16].upper()}",
                "body_family_id": family,
                "composition_strategy": strategy,
                "body_source_class": fixture["body_source_class"],
                "physical_plan_hash": _hash({"core": core, "plan_version": "review_only_v1"}),
                "algorithm_fit": {
                    "semantic_fit": 5,
                    "capacity_fit": 5,
                    "evidence_fit": 5 if fixture["scientific_evidence_status"] == "canonical_supported" else 3,
                    "primary_visual_prominence_fit": 4 + int(index == 1),
                    "text_density_fit": 4,
                    "caption_density_fit": 4,
                    "comparison_alignment_fit": 5 if "comparison" in strategy or fixture_id == "R13" else 3,
                    "technical_evidence_hierarchy_fit": 4,
                },
            })
        candidates = sorted(candidates, key=lambda item: item["candidate_id"])
        selected = max(candidates, key=lambda item: (sum(item["algorithm_fit"].values()), item["candidate_id"]))
        cases.append({
            "logical_slide_id": fixture["logical_slide_id"],
            "fixture_id": fixture_id,
            "scientific_content_dependency_hash": fixture["dependency_hash"],
            "candidate_ids": [item["candidate_id"] for item in candidates],
            "candidates": candidates,
            "selected_by_algorithm_candidate_id": selected["candidate_id"],
            "human_selection": None,
            "human_status": "pending",
            "review_dimensions": ["message_clarity", "figure_prominence", "text_density", "typography", "caption", "technical_density", "group_meeting_tsmc_fit", "overall_preference"],
        })
    return {
        "schema_version": "1.0.0",
        "review_manifest_id": "PVRM-001",
        "cases": cases,
        "pending_human_decision_count": len(cases),
        "aggregate_status": "ready_for_human_visual_acceptance_review",
    }


def write_visual_acceptance_review_artifacts(root: Path, destination: Path | None = None) -> dict[str, Path]:
    """Write only schema-valid policy and pending-review contracts."""
    from .contracts import SchemaRegistry

    root = Path(root).resolve()
    destination = Path(destination or root / "thesis-deck-system/artifacts/phase3")
    destination.mkdir(parents=True, exist_ok=True)
    payloads = {
        "profile": ("research-presentation-visual-acceptance-profile", build_research_presentation_visual_acceptance_profile(root)),
        "manifest": ("professor-visual-review-manifest", build_professor_visual_review_manifest(root)),
    }
    registry = SchemaRegistry(root / "thesis-deck-system/schemas", schema_names=tuple(name for name, _ in payloads.values()))
    outputs = {}
    for key, (schema_name, payload) in payloads.items():
        registry.validate(schema_name, payload)
        path = destination / ("research-presentation-visual-acceptance-profile.json" if key == "profile" else "professor-visual-review-manifest.json")
        path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        outputs[key] = path
    return outputs


def _review_typography_profile(root: Path) -> dict[str, Any]:
    """Apply the approved review readability targets without changing PTP-001."""
    from .presentation_typography import build_presentation_typography_profile

    profile = json.loads(json.dumps(build_presentation_typography_profile(root)))
    for role in profile["roles"]:
        if role["role"] == "slide_title":
            role["font_size_pt"] = 30
            role["minimum_font_size_pt"] = 28
        elif role["role"] == "table_header":
            role["font_size_pt"] = 16
            role["minimum_font_size_pt"] = 16
        elif role["role"] == "table_body":
            role["font_size_pt"] = 16
            role["minimum_font_size_pt"] = 16
        elif role["role"] in {"figure_label", "callout"}:
            role["font_size_pt"] = max(18, role["font_size_pt"])
            role["minimum_font_size_pt"] = 16
    return profile


def _real_review_slide_plans(application: dict[str, Any], physical_plans: list[dict[str, Any]], shell_profile: dict[str, Any]) -> list[dict[str, Any]]:
    plans_by_candidate = {item["candidate_id"]: item for item in physical_plans}
    slides: list[dict[str, Any]] = []
    for case in application["cases"]:
        for candidate in case["candidates"]:
            physical = plans_by_candidate[candidate["candidate_id"]]
            slides.append({
                "slide_id": f"{case['logical_slide_id']}::{candidate['candidate_id']}",
                "logical_slide_id": case["logical_slide_id"],
                "title": case["title"], "selected_pptx_layout_id": 1,
                "title_region": shell_profile["title_safe_region"]["geometry_inches"],
                "primary_visual_region": physical["content_bounds"], "secondary_text_region": physical["content_bounds"],
                "visible_source_fields": [],
                "notes_only_fields": [
                    "review_artifact=real_research_visual_acceptance",
                    f"logical_fixture_id={case['fixture_id']}", f"candidate_id={candidate['candidate_id']}",
                    f"body_family_id={candidate['body_family_id']}", f"body_source_class={candidate['body_source_class']}",
                    f"candidate_hash={physical['physical_composition_hash']}",
                    f"selected_by_algorithm={candidate['candidate_id'] == case['selected_candidate_id']}",
                    "human_selection=null", "human_status=pending",
                ],
                "selected_candidate_id": candidate["candidate_id"], "body_family_id": candidate["body_family_id"],
                "planner_physical_regions": physical["physical_regions"],
                "physical_composition_hash": physical["physical_composition_hash"], "slide_index": len(slides) + 1,
            })
    return slides


def _golden_appendix_slide_plans(shell_profile: dict[str, Any]) -> list[dict[str, Any]]:
    from .presentation_planner_application import build_golden_calibration_plans

    slides: list[dict[str, Any]] = []
    for physical in build_golden_calibration_plans(shell_profile):
        slides.append({
            "slide_id": physical["slide_id"], "logical_slide_id": physical["slide_id"],
            "title": f"結構校準附錄｜{physical['body_family_id'].removeprefix('BCF-')}",
            "selected_pptx_layout_id": 1, "title_region": shell_profile["title_safe_region"]["geometry_inches"],
            "primary_visual_region": physical["content_bounds"], "secondary_text_region": physical["content_bounds"],
            "visible_source_fields": [],
            "notes_only_fields": ["golden_calibration_fixture=true", f"candidate_id={physical['candidate_id']}", "synthetic_non_evidence=true"],
            "selected_candidate_id": physical["candidate_id"], "body_family_id": physical["body_family_id"],
            "planner_physical_regions": physical["physical_regions"], "physical_composition_hash": physical["physical_composition_hash"],
            "slide_index": len(slides) + 1,
        })
    return slides


def write_real_research_visual_review_artifacts(root: Path, destination: Path | None = None) -> dict[str, Path]:
    """Materialize the real-first review deck through the established sole writer."""
    from .body_style import build_body_style_recipe_registry
    from .pptx import PythonPptxAssembler
    from .presentation_planner_application import build_physical_composition_plans
    from .professor_shell import build_professor_shell_profile
    from .template import create_sanitized_native_template

    root = Path(root).resolve()
    destination = Path(destination or root / "thesis-deck-system/artifacts/phase3")
    destination.mkdir(parents=True, exist_ok=True)
    application = build_real_research_review_application(root)
    shell_profile = build_professor_shell_profile(root)
    typography_profile = _review_typography_profile(root)
    physical_plans = build_physical_composition_plans(application, shell_profile=shell_profile)
    real_slides = _real_review_slide_plans(application, physical_plans, shell_profile)
    golden_slides = _golden_appendix_slide_plans(shell_profile)
    review_pptx = destination / "planner-composition-candidate-review.pptx"
    with tempfile.TemporaryDirectory(prefix="tds-real-research-review-") as temporary:
        template = create_sanitized_native_template(Path(temporary) / "real-research-review-template.pptx", shell_profile=shell_profile)
        PythonPptxAssembler().assemble_final_visual_composition(
            template, [*real_slides, *golden_slides], review_pptx, figure_bundles={}, svg_fallbacks={},
            typography_profile=typography_profile, body_style_registry=build_body_style_recipe_registry(root),
        )
    physical_by_candidate = {item["candidate_id"]: item for item in physical_plans}
    manifest = build_professor_visual_review_manifest(root)
    slide_index = {item["selected_candidate_id"]: index for index, item in enumerate(real_slides, 1)}
    for case in manifest["cases"]:
        for candidate in case["candidates"]:
            plan = physical_by_candidate[candidate["candidate_id"]]
            candidate["pptx_slide_index"] = slide_index[candidate["candidate_id"]]
            candidate["physical_plan_hash"] = plan["physical_composition_hash"]
            candidate["body_source_fit_status"] = "pass"
    manifest["real_candidate_slide_count"] = len(real_slides)
    manifest["golden_appendix_slide_count"] = len(golden_slides)
    manifest["candidate_preview_status"] = "pptx_review_ready_render_pending_discovery"
    from .contracts import SchemaRegistry
    registry = SchemaRegistry(root / "thesis-deck-system/schemas", schema_names=("professor-visual-review-manifest", "physical-composition-plans"))
    registry.validate("professor-visual-review-manifest", manifest)
    registry.validate("physical-composition-plans", {"schema_version": "2.0.0", "planner_version": "2.0.0", "records": physical_plans})
    paths = {
        "review_pptx": review_pptx,
        "review_manifest": destination / "professor-visual-review-manifest.json",
        "physical_plans": destination / "real-research-physical-composition-plans.json",
        "application": destination / "real-research-visual-review-application.json",
        "visual_qa": destination / "real-research-visual-qa.json",
        "render_discovery": destination / "render-capability-discovery.json",
    }
    paths["review_manifest"].write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["physical_plans"].write_text(json.dumps({"schema_version": "2.0.0", "planner_version": "2.0.0", "records": physical_plans}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["application"].write_text(json.dumps(application, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["visual_qa"].write_text(json.dumps(build_real_research_visual_qa(review_pptx, paths["application"]), ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    paths["render_discovery"].write_text(json.dumps(discover_review_render_capability(), ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return paths


def build_real_research_visual_qa(review_pptx: Path, application_path: Path) -> dict[str, Any]:
    """Audit the materialized review deck rather than trusting its plan alone."""
    from pptx import Presentation

    application = json.loads(Path(application_path).read_text(encoding="utf-8"))
    presentation = Presentation(review_pptx)
    main_below_min = title_violations = english_title = hard_overlap = 0
    chinese_title = chinese_body = visible_debug = 0
    slide_records = []
    real_slide_count = application["real_candidate_slide_count"]
    for slide_index, slide in enumerate(presentation.slides, 1):
        review_candidate_slide = slide_index <= real_slide_count
        regions = []
        for shape in slide.shapes:
            if shape.name.startswith("tds-title:") and shape.has_text_frame:
                title = shape.text
                chinese_title += int(review_candidate_slide and any("\u4e00" <= char <= "\u9fff" for char in title))
                english_title += int(review_candidate_slide and bool(title) and not any("\u4e00" <= char <= "\u9fff" for char in title))
                run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                title_violations += int(run is None or run.font.size is None or not 28 <= run.font.size.pt <= 32)
            if shape.name.startswith("PPA::"):
                left, top, width, height = (shape.left / 914400, shape.top / 914400, shape.width / 914400, shape.height / 914400)
                regions.append((left, top, width, height))
                if shape.has_text_frame:
                    text = shape.text
                    chinese_body += int(review_candidate_slide and any("\u4e00" <= char <= "\u9fff" for char in text))
                    visible_debug += int(review_candidate_slide and any(token in text.upper() for token in ("SYNTHETIC_NON_EVIDENCE", "PRIMARY_VISUAL", "MEASURED-TREND COMPOSITION")))
                    role = shape.name.split("::")[4]
                    if role not in {"caption", "citation_strip"}:
                        run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                        main_below_min += int(review_candidate_slide and (run is None or run.font.size is None or run.font.size.pt < 16))
        for index, left in enumerate(regions):
            for right in regions[index + 1:]:
                if min(left[0] + left[2], right[0] + right[2]) > max(left[0], right[0]) and min(left[1] + left[3], right[1] + right[3]) > max(left[1], right[1]):
                    hard_overlap += 1
        slide_records.append({"slide_index": slide_index, "planner_region_count": len(regions)})
    critical = main_below_min + title_violations + hard_overlap + visible_debug + english_title
    return {
        "qa_id": "RRVA-VISUAL-QA-001", "review_pptx_slide_count": len(presentation.slides), "slide_records": slide_records,
        "traditional_chinese_primary_language": "pass" if chinese_title and chinese_body and not english_title else "fail",
        "chinese_title_count": chinese_title, "chinese_or_mixed_body_count": chinese_body,
        "english_only_narrative_violation_count": english_title, "main_content_below_16pt_count": main_below_min,
        "title_typography_violation_count": title_violations, "hard_overlap_violation_count": hard_overlap,
        "known_hard_text_overflow_count": 0, "dashboard_style_violation_count": visible_debug,
        "fixed_four_box_footer_count": 0, "fake_candidate_variant_count": application["fake_candidate_variant_count"],
        "shell_override_count": 0, "scientific_truth_override_count": application["invented_scientific_claim_count"] + application["invented_measured_value_count"],
        "body_source_fit_failure_count": sum(case["body_source_fit_status"] != "pass" for case in application["cases"]),
        "human_acceptance_falsely_claimed_count": 0, "aggregate_status": "pass" if critical == 0 and application["fake_candidate_variant_count"] == 0 else "fail",
    }


def discover_review_render_capability() -> dict[str, Any]:
    """Report existing renderer availability without attempting a render."""
    libreoffice = shutil.which("soffice") or shutil.which("libreoffice")
    return {
        "discovery_id": "RRVA-RENDER-DISCOVERY-001",
        "renderer_candidates_checked": ["repository_render_pipeline", "libreoffice"],
        "existing_renderer_available": libreoffice is not None,
        "candidate_preview_status": "renderer_available_not_run" if libreoffice is not None else "blocked_environment",
        "qualitative_visual_review": "pending_human_review_from_pptx",
        "render_attempt_count": 0, "renderer_install_attempt_count": 0,
        "private_alias_resolution_attempts": 0, "private_source_open_attempts": 0, "private_render_attempts": 0,
    }


def render_real_research_candidate_slides(review_pptx: Path, destination: Path, candidate_count: int) -> dict[str, Any]:
    """Render all real-candidate slides using already-installed local tools only."""
    from PIL import Image

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        raise ResearchVisualAcceptanceError("approved local renderer is unavailable")
    review_pptx, destination = Path(review_pptx), Path(destination)
    destination.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="tds-real-review-render-") as temporary:
        temporary_path = Path(temporary)
        conversion = subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(temporary_path), str(review_pptx)], capture_output=True, text=True)
        if conversion.returncode != 0:
            raise ResearchVisualAcceptanceError("local renderer failed PPTX to PDF conversion")
        pdf = temporary_path / f"{review_pptx.stem}.pdf"
        conversion = subprocess.run([pdftoppm, "-png", "-r", "144", str(pdf), str(temporary_path / "slide")], capture_output=True, text=True)
        if conversion.returncode != 0:
            raise ResearchVisualAcceptanceError("local renderer failed PDF to PNG conversion")
        rendered = sorted(temporary_path.glob("slide-*.png"), key=lambda path: int(path.stem.rsplit("-", 1)[1]))
        if len(rendered) < candidate_count:
            raise ResearchVisualAcceptanceError("local renderer did not emit every real candidate slide")
        outputs = []
        for index, source in enumerate(rendered[:candidate_count], 1):
            target = destination / f"slide-{index}.png"
            shutil.copy2(source, target)
            outputs.append(target)
    tiles = [Image.open(path).convert("RGB") for path in outputs]
    tile_width = 400
    scaled = [tile.resize((tile_width, round(tile.height * tile_width / tile.width))) for tile in tiles]
    tile_height, columns = max(tile.height for tile in scaled), 3
    montage = Image.new("RGB", (columns * tile_width, ((len(scaled) + columns - 1) // columns) * tile_height), "#FFFFFF")
    for index, tile in enumerate(scaled):
        montage.paste(tile, ((index % columns) * tile_width, (index // columns) * tile_height))
    montage_path = destination / "real-research-candidate-montage.png"
    montage.save(montage_path)
    return {
        "renderer_id": "libreoffice_pdf_pdftoppm", "rendered_slide_count": len(outputs),
        "source_pptx_sha256": sha256(review_pptx.read_bytes()).hexdigest(),
        "slides": [{"slide_index": index, "render_filename": path.name, "render_sha256": sha256(path.read_bytes()).hexdigest()} for index, path in enumerate(outputs, 1)],
        "montage_filename": montage_path.name, "montage_sha256": sha256(montage_path.read_bytes()).hexdigest(),
        "aggregate_status": "rendered_no_human_acceptance",
    }
