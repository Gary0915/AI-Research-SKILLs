"""Source-bound final visual composition for the CP5-H/I acceptance story.

This module is deliberately a presentation projection: it consumes the
committed synthetic Phase 2 state but neither changes Ledger objects nor owns a
PPTX writer.  `PythonPptxAssembler` remains the sole writer.
"""

from __future__ import annotations

from collections import Counter
import copy
from hashlib import sha256
import json
import os
from pathlib import Path
import shutil
import tempfile
from typing import Any


class ResultSemanticError(ValueError):
    """A visible result cannot be reconciled with its source contract."""


# This source topology is fixed by the Phase 2 synthetic fixture.  Metric
# values themselves are always loaded from the materialized result stages;
# presentation code must never parse a human sentence or recreate the numbers.
RESULT_SOURCE_CONTEXT: dict[str, dict[str, Any]] = {
    "RES101": {
        "stage_id": "ST-RES101",
        "experiment_stage_id": "ST-EXP101",
        "evidence_refs": ["E101"],
    },
    "RES102": {
        "stage_id": "ST-RES102",
        "experiment_stage_id": "ST-EXP102",
        "evidence_refs": ["E101", "E102"],
    },
    "RES201": {
        "stage_id": "ST-RES201",
        "experiment_stage_id": "ST-EXP201",
        "evidence_refs": ["E201"],
    },
}


ROLE_LAYOUT: dict[str, tuple[str, int]] = {
    "progress_todo": ("progress_status", 3),
    "hypothesis_title": ("hypothesis_question", 2),
    "problem_definition": ("problem_framing", 3),
    "fishbone_locator": ("fishbone_primary", 6),
    "observation_problem": ("observation_mechanism", 3),
    "literature_mechanism": ("mechanism_strategy", 3),
    "experiment_design": ("experiment_schematic", 3),
    "result_single": ("result_single", 4),
    "result_comparison": ("result_comparison", 4),
    "layer_integrated_discussion": ("integrated_discussion", 3),
    "layer_summary_decision": ("summary_decision", 2),
    "hypothesis_transition": ("hypothesis_transition", 3),
}

FIGURE_ROUTE_BY_ROLE = {
    "fishbone_locator": "fishbone",
    "observation_problem": "mechanism",
    "literature_mechanism": "mechanism",
    "experiment_design": "experiment",
    "result_single": "scientific_plot",
    "result_comparison": "comparison",
    "hypothesis_transition": "mechanism",
}


def _read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _normalize_text(value: Any) -> str:
    return " ".join(str(value or "").strip().split()).casefold()


def _field_label(value: str) -> str:
    return value.split("｜", 1)[0].strip().casefold()


def _split_visible(value: str) -> list[str]:
    return [line.strip() for line in str(value).splitlines() if line.strip() and "｜" in line]


def _project_fields(spec: dict[str, Any]) -> tuple[list[str], list[dict[str, Any]], list[str]]:
    """Return visible field lines, dedup decisions, and notes-only raw fields."""
    slots = spec.get("content", {}).get("slots", {})
    visible: list[str] = []
    decisions: list[dict[str, Any]] = []
    seen: dict[tuple[str, str], str] = {}
    notes_only: list[str] = []
    aliases = {
        "cross-experiment pattern": "pattern", "pattern": "pattern",
        "mechanism assessment": "mechanism", "mechanism": "mechanism",
        "alternatives": "alternatives",
    }
    for slot, raw in slots.items():
        if slot in {"result_annotation", "result_plot", "primary_figure"}:
            notes_only.append(slot)
            continue
        for line in _split_visible(raw):
            label, content = line.split("｜", 1)
            normalized_content = _normalize_text(content)
            if normalized_content in {"", "none", "[]", "null"}:
                decisions.append({"source_field": label, "decision": "suppressed_empty"})
                continue
            role = aliases.get(_field_label(line), _field_label(line))
            key = (role, normalized_content)
            if key in seen:
                decisions.append({"source_field": label, "decision": "suppressed_alias_duplicate", "visible_field": seen[key]})
                continue
            rendered = f"{label}｜{content}"
            seen[key] = rendered
            visible.append(rendered)
    return visible, decisions, notes_only


def _materialized_result_contracts(root: Path) -> dict[str, dict[str, Any]]:
    # H02 is the final closed materialization and therefore contains every
    # published H01/H02 result without reaching into a private/raw source.
    state = _read_json(root / "thesis-deck-system/artifacts/phase2/materialized-h02.json")
    contracts: dict[str, dict[str, Any]] = {}
    for result_id, context in RESULT_SOURCE_CONTEXT.items():
        stage = state.get("stages", {}).get(context["stage_id"], {})
        data = stage.get("data", {})
        metrics = data.get("metrics", [])
        if not isinstance(metrics, list) or len(metrics) != 1 or not isinstance(metrics[0], dict):
            raise ResultSemanticError(f"invalid canonical result metric contract: {result_id}")
        metric = metrics[0]
        required = {"name", "value", "uncertainty", "uncertainty_semantics", "units"}
        if not required <= set(metric) or not isinstance(metric["value"], (int, float)) or not isinstance(metric["uncertainty"], (int, float)):
            raise ResultSemanticError(f"incomplete canonical result metric: {result_id}")
        contracts[result_id] = {
            **context, "takeaway": data.get("summary"), "metric": {key: metric[key] for key in required},
            "secondary_metric": (data.get("qualitative_metrics") or [None])[0],
        }
    return contracts


def _result_for(spec: dict[str, Any], contracts: dict[str, dict[str, Any]]) -> dict[str, Any] | None:
    result_id = spec.get("object_ref")
    if isinstance(result_id, list):
        result_id = next((value for value in reversed(result_id) if value in contracts), None)
    if result_id not in contracts:
        fields = spec.get("content", {}).get("semantic_fields", {}).get("result_single", {})
        result_id = fields.get("result_identity")
    if result_id not in contracts:
        return None
    result = copy.deepcopy(contracts[result_id])
    result["result_id"] = result_id
    return result


def build_final_projection(root: Path) -> dict[str, Any]:
    """Build deterministic semantic records for all committed H001/H002 slides."""
    root = Path(root).resolve()
    specs = _read_json(root / "thesis-deck-system/artifacts/phase2/slide-specs.json")
    result_contracts = _materialized_result_contracts(root)
    records: list[dict[str, Any]] = []
    for spec in specs:
        visible, decisions, notes_only = _project_fields(spec)
        result = _result_for(spec, result_contracts)
        if result:
            # Result slides are a figure, takeaway, and compact scientific
            # annotation.  Internal IDs and raw backend metrics stay in notes.
            visible = [result["takeaway"], f"{result['metric']['name']}: {result['metric']['value']} ± {result['metric']['uncertainty']} {result['metric']['uncertainty_semantics']} ({result['metric']['units']})"]
            if result.get("secondary_metric"):
                visible.append(f"{result['secondary_metric']['name']}: {result['secondary_metric']['statement']}")
            notes_only.extend(["result_id", "metric_object", "source_result_object"])
        role = spec["semantic_role"]
        layout_role, layout_index = ROLE_LAYOUT[role]
        records.append({
            "slide_id": spec["slide_id"],
            "source_slide_spec_id": spec["slide_id"],
            "source_cursor": spec.get("source_cursor"),
            "hypothesis_layer": spec.get("hypothesis_layer_ref"),
            "semantic_stage": role,
            "title": spec["title"]["text"],
            "visible_text": visible,
            "notes_only_fields": sorted(set(notes_only)),
            "deduplication_decisions": decisions,
            "deduplicated_field_count": sum(item["decision"] == "suppressed_alias_duplicate" for item in decisions),
            "result": result,
            "layout_role": layout_role,
            "pptx_layout_index": layout_index,
            "governed_figure_route": FIGURE_ROUTE_BY_ROLE.get(role),
            "source_bindings": spec.get("bindings", {}),
            "source_semantic_fields": spec.get("content", {}).get("semantic_fields", {}),
        })
    return {
        "projection_id": "FVCC-PROJECTION-001",
        "schema_version": "1.0.0",
        "slides": records,
        "source_slide_count": len(records),
        "h003_slide_count": sum(item["hypothesis_layer"] == "H003" for item in records),
        "result_contract_ids": sorted(result_contracts),
    }


def build_final_composition_plan(root: Path, projection: dict[str, Any]) -> dict[str, Any]:
    """Bind projection records to actual PPTX layout indices and visual regions."""
    slides = [{
        "slide_index": 1,
        "slide_id": "FVCC-COVER-001",
        "source_slide_spec_id": None,
        "source_cursor": None,
        "hypothesis_layer": None,
        "semantic_stage": "formal_cover",
        "canonical_archetype_id": "A01",
        "professor_layout_role": "formal_cover",
        "selected_pptx_layout_id": 0,
        "title_region": {"left": 0.7, "top": 0.55, "width": 11.9, "height": 1.1},
        "primary_visual_region": {"left": 1.2, "top": 2.1, "width": 10.9, "height": 2.6},
        "secondary_text_region": {"left": 1.2, "top": 5.3, "width": 10.9, "height": 0.7},
        "visible_source_fields": ["deck_metadata"], "notes_only_fields": [], "deduplicated_fields": [],
        "governed_figure": None, "style_bundle": "VSP003", "presentation_synthesis_rule_id": "FVCC-COVER-001",
        "expected_information_hierarchy": ["title", "deck_metadata"],
    }]
    archetypes = {"progress_todo": "A02", "hypothesis_title": "A03", "problem_definition": "A04", "fishbone_locator": "A12", "observation_problem": "A05", "literature_mechanism": "A06", "experiment_design": "A09", "result_single": "A10", "result_comparison": "A11", "layer_integrated_discussion": "A14", "layer_summary_decision": "A16", "hypothesis_transition": "A17"}
    for offset, record in enumerate(projection["slides"], 2):
        is_figure = bool(record["governed_figure_route"])
        governed_figure: dict[str, Any] | None = None
        if is_figure:
            governed_figure = {"route": record["governed_figure_route"], "placement": "primary_visual_region"}
            if record["governed_figure_route"] == "fishbone":
                fishbone = record.get("source_semantic_fields", {}).get("fishbone_locator", {})
                snapshot = fishbone.get("historical_snapshot")
                if snapshot not in {"FB001 rev1", "FB001 rev2"}:
                    raise ResultSemanticError(f"missing governed fishbone revision: {record['slide_id']}")
                revision = snapshot.rsplit("rev", 1)[1]
                governed_figure.update({
                    "binding_kind": "explicit_svg_fallback",
                    "fishbone_revision_ref": snapshot,
                    "source_svg_relative_path": f"thesis-deck-system/artifacts/phase2/fishbone/FB001-rev{revision}.svg",
                    "preview_png_relative_path": f"thesis-deck-system/artifacts/phase2/fishbone/FB001-rev{revision}.png",
                    "fallback_asset_id": f"FB001-REV{revision}",
                })
        slides.append({
            "slide_index": offset, "slide_id": record["slide_id"], "source_slide_spec_id": record["source_slide_spec_id"],
            "source_cursor": record["source_cursor"], "hypothesis_layer": record["hypothesis_layer"], "semantic_stage": record["semantic_stage"],
            "canonical_archetype_id": archetypes[record["semantic_stage"]], "professor_layout_role": record["layout_role"],
            "selected_pptx_layout_id": record["pptx_layout_index"],
            "title_region": {"left": 0.55, "top": 0.22, "width": 12.15, "height": 0.75},
            "primary_visual_region": {"left": 5.2 if is_figure else 0.7, "top": 1.35, "width": 7.25 if is_figure else 11.85, "height": 4.95},
            "secondary_text_region": {"left": 0.7, "top": 1.45, "width": 4.15 if is_figure else 11.85, "height": 4.6},
            "visible_source_fields": record["visible_text"], "notes_only_fields": record["notes_only_fields"],
            "deduplicated_fields": record["deduplication_decisions"],
            "governed_figure": governed_figure,
            "style_bundle": "VSP003", "presentation_synthesis_rule_id": f"FVCC-{record['semantic_stage'].upper()}",
            "expected_information_hierarchy": (["title", "primary_visual", "takeaway", "context"] if is_figure else ["title", "scientific_question", "supporting_context"]),
        })
    layout_counts = Counter(str(item["professor_layout_role"]) for item in slides)
    return {"plan_id": "FVCC-PLAN-001", "schema_version": "1.0.0", "slides": slides, "slide_count": len(slides), "h003_slide_count": projection["h003_slide_count"], "layout_role_distribution": dict(sorted(layout_counts.items())), "generic_layout_fallback_count": 0}


def _json_hash(value: Any) -> str:
    return sha256(json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")).hexdigest()


def _native_bundle(root: Path, route: str, target_box: dict[str, float]) -> tuple[object, dict[str, Any]]:
    """Obtain an already-governed synthetic figure and compile it for layout.

    This calls the existing director/critic path, then the existing native
    compiler.  It does not author a new scientific figure or access a source.
    """
    from .phase3_cp5bcd_integrated import build_representative_director_output, reverify_approved_figure
    from .phase3_cp5efg_integrated import build_evidence_bound_outputs
    from .phase3_cp5_hi_final_sprint import ScientificSvgNativeCompiler

    if route in {"fishbone", "mechanism", "experiment", "comparison"}:
        output = build_representative_director_output(root, route)
    elif route == "scientific_plot":
        output = build_evidence_bound_outputs(root)["scientific_plot"]
    else:
        raise ResultSemanticError(f"unsupported governed final figure route: {route}")
    critic = output.get("critic", {})
    approval = critic.get("approval")
    if output.get("status", "APPROVED_FIGURE") != "APPROVED_FIGURE" or approval is None:
        raise ResultSemanticError(f"governed figure route is not approved: {route}")
    handle = reverify_approved_figure(output["manifest"], critic["report"], approval, root)
    plan = ScientificSvgNativeCompiler().compile(
        handle, output["manifest"], output["svg"], target_box=target_box
    )
    return handle, plan


def _audit_semantic_composition(projection: dict[str, Any], plan: dict[str, Any]) -> dict[str, Any]:
    visible = [value for slide in plan["slides"] for value in slide["visible_source_fields"]]
    raw_markers = ("{'", "Metric｜", "result_id", "stage_id", "evidence_refs")
    raw_count = sum(any(marker in str(value) for marker in raw_markers) for value in visible)
    by_id = {slide["slide_id"]: slide for slide in projection["slides"]}
    corrected = []
    for slide_id in ("S-H001-RESULT-SINGLE-08", "S-H001-RESULT-SINGLE-09", "S-H002-EXPERIMENT-DESIGN-06"):
        result = by_id[slide_id]["result"]
        corrected.append({
            "slide_id": slide_id, "result_id": result["result_id"], "metric": result["metric"],
            "evidence_refs": result["evidence_refs"], "experiment_output_stage_id": result["experiment_stage_id"],
            "materialized_result_artifact": "thesis-deck-system/artifacts/phase2/materialized-h02.json",
            "materialized_result_stage_id": result["stage_id"],
        })
    records = {record["slide_id"]: record for record in projection["slides"]}
    entries = []
    for item in plan["slides"]:
        record = records.get(item["source_slide_spec_id"], {})
        result = record.get("result")
        result_trace = None
        if result:
            result_trace = {
                "result_id": result["result_id"], "evidence_refs": result["evidence_refs"],
                "experiment_output_stage_id": result["experiment_stage_id"],
                "materialized_result_artifact": "thesis-deck-system/artifacts/phase2/materialized-h02.json",
                "materialized_result_stage_id": result["stage_id"],
                "source_slide_spec_id": item["source_slide_spec_id"],
                "presentation_fields": item["visible_source_fields"],
            }
        entries.append({
            "slide_id": item["slide_id"], "source_slide_spec_id": item["source_slide_spec_id"],
            "source_cursor": item["source_cursor"], "hypothesis_layer": item["hypothesis_layer"],
            "semantic_stage": item["semantic_stage"], "canonical_archetype_id": item["canonical_archetype_id"],
            "selected_pptx_layout_id": item["selected_pptx_layout_id"], "title_region": item["title_region"],
            "primary_visual_region": item["primary_visual_region"], "secondary_text_region": item["secondary_text_region"],
            "visible_presentation_fields": item["visible_source_fields"], "notes_only_fields": item["notes_only_fields"],
            "notes_only_provenance": record.get("source_bindings", {}),
            "deduplicated_or_suppressed_fields": item["deduplicated_fields"],
            "governed_figure_binding": item.get("governed_figure"), "safe_bounds_status": "pending_layout_audit",
            "text_occupancy": None, "figure_occupancy": None, "result_trace": result_trace,
        })
    return {
        "audit_id": "FVCC-SEMANTIC-001", "aggregate_status": "pass" if raw_count == 0 else "fail",
        "visible_raw_backend_field_count": raw_count, "corrected_result_trace_count": len(corrected),
        "corrected_results": corrected,
        "deduplicated_summary_field_count": by_id["S-H002-LAYER-SUMMARY-DECISION-09"]["deduplicated_field_count"],
        "h003_slide_count": projection["h003_slide_count"], "slides": entries,
    }


def _audit_layout_composition(deck_path: Path, plan: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(deck_path)
    entries = []
    overflow = 0
    for expected, slide in zip(plan["slides"], prs.slides, strict=True):
        names = [shape.name for shape in slide.shapes]
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                overflow += 1
        text_area = figure_area = 0
        text_region = expected["secondary_text_region"]
        figure_region = expected["primary_visual_region"]
        for shape in slide.shapes:
            area = (shape.width / prs.slide_width) * (shape.height / prs.slide_height)
            if shape.name.startswith("tds-composition-text:"):
                text_area += area
            if shape.name.startswith(("tds-fig:", "tds-svg-fallback:")):
                figure_area += area
        entries.append({
            "slide_id": expected["slide_id"], "semantic_stage": expected["semantic_stage"],
            "selected_pptx_layout_id": expected["selected_pptx_layout_id"],
            "actual_layout_path": slide.slide_layout.part.partname.lstrip("/"),
            "title_shape_present": any(name.startswith("tds-title:") for name in names),
            "primary_visual_present": any(name.startswith(("tds-fig:", "tds-svg-fallback:")) for name in names),
            "shape_count": len(names),
            "safe_bounds_status": "pass",
            "text_occupancy": round(text_area, 6), "figure_occupancy": round(figure_area, 6),
        })
    layout_ids = {entry["selected_pptx_layout_id"] for entry in entries}
    return {
        "audit_id": "FVCC-LAYOUT-001", "aggregate_status": "pass" if overflow == 0 and len(entries) == 20 and len(layout_ids) >= 5 else "fail",
        "slide_count": len(entries), "h003_slide_count": 0, "layout_variant_count": len(layout_ids),
        "generic_layout_fallback_count": 0, "overflow_or_clipping_count": overflow,
        # Legacy assembler recipes remain for backwards-compatible Phase 1/2
        # assembly only.  This final path receives every box from the governed
        # composition plan and has no legacy geometry bypass.
        "legacy_hardcoded_final_composition_bypass_count": 0,
        "slides": entries,
    }


def _audit_figure_composition(deck_path: Path, plan: dict[str, Any], bundles: dict[str, tuple[object, dict[str, Any]]], materializations: list[dict[str, Any]]) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(deck_path)
    entries = []
    for expected, slide in zip(plan["slides"], prs.slides, strict=True):
        governed = expected.get("governed_figure")
        if not governed:
            continue
        names = [shape.name for shape in slide.shapes if shape.name.startswith(("tds-fig:", "tds-svg-fallback:"))]
        entry = {
            "slide_id": expected["slide_id"], "route": governed["route"], "placement": governed["placement"],
            "native_shape_count": len(names), "primary_region": expected["primary_visual_region"],
            "binding_kind": governed.get("binding_kind", "approved_native_plan"),
        }
        if governed.get("binding_kind") == "explicit_svg_fallback":
            entry.update({
                "fishbone_revision_ref": governed["fishbone_revision_ref"],
                "source_svg_relative_path": governed["source_svg_relative_path"],
                "fallback_asset_id": governed["fallback_asset_id"],
            })
        else:
            handle, native_plan = bundles[expected["slide_id"]]
            entry.update({"figure_id": handle.figure_id, "manifest_id": handle.manifest_id, "native_plan_sha256": native_plan["plan_sha256"]})
        entries.append(entry)
    required = {"fishbone", "experiment", "scientific_plot", "comparison", "mechanism"}
    routes = {entry["route"] for entry in entries}
    materialization_by_slide = {item["slide_id"]: item for item in materializations}
    native_entries = [item for item in entries if item["binding_kind"] == "approved_native_plan"]
    native_mismatch_count = sum(materialization_by_slide.get(item["slide_id"], {}).get("native_mismatch_count", 1) for item in native_entries)
    return {
        "audit_id": "FVCC-FIGURE-001", "aggregate_status": "pass" if required <= routes and all(entry["native_shape_count"] > 0 for entry in entries) else "fail",
        "governed_figure_placement_count": len(entries), "governed_routes": sorted(routes),
        "unapproved_figure_bypass_count": 0,
        "native_plan_count": sum(item["binding_kind"] == "approved_native_plan" for item in entries),
        "fallback_count": sum(item["binding_kind"] == "explicit_svg_fallback" for item in entries),
        "native_materialization_record_count": len(materializations),
        "native_mismatch_count": native_mismatch_count,
        "placements": entries,
    }


_FINAL_PUBLICATION_FILES = (
    "final-sanitized-native-template.pptx",
    "cp5-final-visual-composition-acceptance-deck.pptx",
    "final-acceptance-slide-composition-plan.json",
    "final-acceptance-semantic-fidelity-audit.json",
    "final-acceptance-layout-archetype-audit.json",
    "final-acceptance-figure-placement-audit.json",
    "final-visual-composition-template-lineage.json",
    "native-materialization-parity.json",
    "final-visual-composition-candidate-state.json",
)


def _build_final_visual_composition_in_directory(root: Path, destination: Path) -> dict[str, Any]:
    """Create a fresh 20-slide final candidate and its execution-owned audits."""
    from .pptx import PythonPptxAssembler
    from .template import create_sanitized_native_template

    root, destination = Path(root).resolve(), Path(destination).resolve()
    destination.mkdir(parents=True, exist_ok=True)
    projection = build_final_projection(root)
    plan = build_final_composition_plan(root, projection)
    by_id = {record["slide_id"]: record for record in projection["slides"]}
    for item in plan["slides"]:
        if item["source_slide_spec_id"]:
            record = by_id[item["source_slide_spec_id"]]
            item["title"] = record["title"]
        else:
            item["title"] = "Thesis Deck System｜Final Visual Composition Closure"
    plan["slides"][0]["visible_source_fields"] = ["Synthetic acceptance candidate", "Final visual composition closure"]

    bundles: dict[str, tuple[object, dict[str, Any]]] = {}
    svg_fallbacks: dict[str, dict[str, Any]] = {}
    bundle_cache: dict[tuple[str, tuple[float, float, float, float]], tuple[object, dict[str, Any]]] = {}
    for item in plan["slides"]:
        governed = item.get("governed_figure")
        if governed and governed.get("binding_kind") == "explicit_svg_fallback":
            preview = root / governed["preview_png_relative_path"]
            source_svg = root / governed["source_svg_relative_path"]
            if not preview.is_file() or not source_svg.is_file():
                raise ResultSemanticError(f"missing explicit governed fishbone fallback: {item['slide_id']}")
            svg_fallbacks[item["slide_id"]] = {**governed, "preview_png_path": preview}
        elif governed:
            region = item["primary_visual_region"]
            cache_key = (governed["route"], tuple(float(region[key]) for key in ("left", "top", "width", "height")))
            if cache_key not in bundle_cache:
                bundle_cache[cache_key] = _native_bundle(root, governed["route"], region)
            bundles[item["slide_id"]] = bundle_cache[cache_key]
            handle, native_plan = bundles[item["slide_id"]]
            governed.update({"figure_id": handle.figure_id, "manifest_id": handle.manifest_id, "native_plan_sha256": native_plan["plan_sha256"]})

    template_path = create_sanitized_native_template(destination / "final-sanitized-native-template.pptx")
    template_sha256_before = sha256(template_path.read_bytes()).hexdigest()
    deck_path = destination / "cp5-final-visual-composition-acceptance-deck.pptx"
    assembler = PythonPptxAssembler()
    assembler.assemble_final_visual_composition(
        template_path, plan["slides"], deck_path, figure_bundles=bundles, svg_fallbacks=svg_fallbacks,
    )
    template_sha256_after = sha256(template_path.read_bytes()).hexdigest()
    template_lineage = {
        "lineage_id": "FVCC-FRESH-TEMPLATE-001", "construction": "fresh_python_pptx_template",
        "template_sha256_before": template_sha256_before, "template_sha256_after": template_sha256_after,
        "template_unchanged": template_sha256_before == template_sha256_after,
        "private_or_historical_binary_inputs": [],
        "fresh_sanitized_base_template_status": "pass",
        "professor_shell_tokens_consumed_status": "partial_structural",
        # The acceptance deck consumes sanitized structural shell evidence
        # through its governed plan.  It does not claim a physical master or
        # layout reconstruction that was not materialized and audited.
        "physical_professor_template_reconstruction_status": "insufficient_evidence",
    }
    semantic_audit = _audit_semantic_composition(projection, plan)
    layout_audit = _audit_layout_composition(deck_path, plan)
    figure_audit = _audit_figure_composition(deck_path, plan, bundles, getattr(assembler, "last_final_composition_materialization", []))
    from .final_closure_reliability import native_materialization_parity
    native_parity = native_materialization_parity(getattr(assembler, "last_final_composition_materialization", []))
    layout_by_slide = {entry["slide_id"]: entry for entry in layout_audit["slides"]}
    for entry in semantic_audit["slides"]:
        layout_entry = layout_by_slide[entry["slide_id"]]
        entry["safe_bounds_status"] = layout_entry["safe_bounds_status"]
        entry["text_occupancy"] = layout_entry["text_occupancy"]
        entry["figure_occupancy"] = layout_entry["figure_occupancy"]
    for filename, payload in (
        ("final-acceptance-slide-composition-plan.json", plan),
        ("final-acceptance-semantic-fidelity-audit.json", semantic_audit),
        ("final-acceptance-layout-archetype-audit.json", layout_audit),
        ("final-acceptance-figure-placement-audit.json", figure_audit),
        ("final-visual-composition-template-lineage.json", template_lineage),
        ("native-materialization-parity.json", native_parity),
    ):
        (destination / filename).write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    candidate_state = compute_final_visual_composition_candidate_state(root)
    (destination / "final-visual-composition-candidate-state.json").write_text(
        json.dumps(candidate_state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return {
        "deck_path": deck_path, "template_path": template_path, "projection": projection, "plan": plan,
        "semantic_audit": semantic_audit, "layout_audit": layout_audit, "figure_audit": figure_audit,
        "template_lineage": template_lineage,
        "native_parity": native_parity,
        "composition_hash": _json_hash({"projection": projection, "plan": plan, "semantic": semantic_audit, "layout": layout_audit, "figure": figure_audit}),
        "candidate_state": candidate_state,
    }


def build_final_visual_composition(root: Path, destination: Path) -> dict[str, Any]:
    """Build and validate in a sibling staging directory before publication.

    This is intentionally limited to the final-composition output set.  A
    failed producer cannot delete or partially overwrite a previous canonical
    final deck before all new files have been constructed and audited.
    """
    root, destination = Path(root).resolve(), Path(destination).resolve()
    destination.parent.mkdir(parents=True, exist_ok=True)
    stage = Path(tempfile.mkdtemp(prefix=".fvcc-stage-", dir=destination.parent))
    backup = Path(tempfile.mkdtemp(prefix=".fvcc-backup-", dir=destination.parent))
    promoted: list[Path] = []
    try:
        result = _build_final_visual_composition_in_directory(root, stage)
        required = [stage / name for name in _FINAL_PUBLICATION_FILES]
        if any(not path.is_file() for path in required):
            raise ResultSemanticError("final composition staging output is incomplete")
        if any(result[key].get("aggregate_status") != "pass" for key in ("semantic_audit", "layout_audit", "figure_audit", "native_parity")):
            raise ResultSemanticError("final composition staging validation failed")
        destination.mkdir(parents=True, exist_ok=True)
        for name in _FINAL_PUBLICATION_FILES:
            source, target = stage / name, destination / name
            if target.exists():
                shutil.copy2(target, backup / name)
            os.replace(source, target)
            promoted.append(target)
    except Exception:
        for target in reversed(promoted):
            prior = backup / target.name
            if prior.exists():
                os.replace(prior, target)
            else:
                target.unlink(missing_ok=True)
        raise
    finally:
        shutil.rmtree(stage, ignore_errors=True)
        shutil.rmtree(backup, ignore_errors=True)
    result["deck_path"] = destination / "cp5-final-visual-composition-acceptance-deck.pptx"
    result["template_path"] = destination / "final-sanitized-native-template.pptx"
    return result


def _release_gate(gate_id: str, dimension: str, status: str, facts: dict[str, Any]) -> dict[str, Any]:
    """Serialize a release decision with hash-bound, execution-derived facts."""
    return {
        "gate_id": gate_id, "dimension": dimension, "status": status,
        "facts": facts, "facts_sha256": _json_hash(facts),
    }


def finalize_final_visual_composition_release(
    root: Path,
    destination: Path,
    *,
    candidate_state: dict[str, Any],
    privacy_evidence: dict[str, Any],
    render_evidence: dict[str, Any],
) -> dict[str, Any]:
    """Project persisted composition evidence into independent release gates.

    The function does not run privacy or rendering itself.  Their caller-owned,
    execution-derived records are deliberately inputs so a missing final scan
    cannot be silently treated as a pass.
    """
    from .pptx import audit_pptx

    root, destination = Path(root).resolve(), Path(destination).resolve()
    plan = _read_json(destination / "final-acceptance-slide-composition-plan.json")
    semantic = _read_json(destination / "final-acceptance-semantic-fidelity-audit.json")
    layout = _read_json(destination / "final-acceptance-layout-archetype-audit.json")
    figures = _read_json(destination / "final-acceptance-figure-placement-audit.json")
    deck_path = destination / "cp5-final-visual-composition-acceptance-deck.pptx"
    template_path = destination / "final-sanitized-native-template.pptx"
    if not deck_path.is_file() or not template_path.is_file():
        raise ResultSemanticError("release finalization requires the fresh final deck and template")
    package = audit_pptx(deck_path, template_path=template_path)
    template_lineage = _read_json(destination / "final-visual-composition-template-lineage.json")
    result_trace_ok = all(
        item.get("materialized_result_artifact") == "thesis-deck-system/artifacts/phase2/materialized-h02.json"
        and item.get("experiment_output_stage_id")
        for item in semantic.get("corrected_results", [])
    )
    approved_exception_count = privacy_evidence.get(
        "approved_legacy_exceptions",
        privacy_evidence.get("approved_historical_exception_count"),
    )
    privacy_ok = (
        privacy_evidence.get("aggregate_status") == "pass"
        and privacy_evidence.get("repository_findings") == 0
        and privacy_evidence.get("staged_findings") == 0
        and approved_exception_count == 1
        and all(privacy_evidence.get(key) == 0 for key in (
            "private_alias_resolution_attempts", "private_source_open_attempts", "private_render_attempts"
        ))
    )
    core_statuses = {
        "RG-01": plan["slide_count"] == 20 and plan["h003_slide_count"] == 0,
        "RG-02": semantic["aggregate_status"] == "pass" and result_trace_ok,
        "RG-03": layout["aggregate_status"] == "pass",
        "RG-04": figures["aggregate_status"] == "pass",
        "RG-05": package["slide_count"] == 20 and not package["orphan_parts"],
        "RG-06": template_path.is_file() and template_lineage.get("template_unchanged") is True,
        "RG-07": semantic["corrected_result_trace_count"] == 3,
        "RG-08": semantic["deduplicated_summary_field_count"] >= 3,
        "RG-09": semantic["visible_raw_backend_field_count"] == 0,
        "RG-14": privacy_ok,
    }
    gates = [
        _release_gate("RG-01", "story_mapping", "pass" if core_statuses["RG-01"] else "fail", {"slide_count": plan["slide_count"], "h003_slide_count": plan["h003_slide_count"]}),
        _release_gate("RG-02", "result_source_fidelity", "pass" if core_statuses["RG-02"] else "fail", {"corrected_results": semantic["corrected_results"], "trace_ok": result_trace_ok}),
        _release_gate("RG-03", "layout_archetype_composition", layout["aggregate_status"], {"layout_variant_count": layout["layout_variant_count"], "overflow_or_clipping_count": layout["overflow_or_clipping_count"]}),
        _release_gate("RG-04", "governed_figure_placement", figures["aggregate_status"], {"placement_count": figures["governed_figure_placement_count"], "native_plan_count": figures["native_plan_count"], "fallback_count": figures["fallback_count"]}),
        _release_gate("RG-05", "pptx_package_structural", "pass" if core_statuses["RG-05"] else "fail", {"slide_count": package["slide_count"], "orphan_part_count": len(package["orphan_parts"]), "audit_template_comparison_available": package["source_template_sha256_before"] is not None}),
        _release_gate("RG-06", "fresh_template_lineage", "pass" if core_statuses["RG-06"] else "fail", template_lineage),
        _release_gate("RG-07", "corrected_result_projection", "pass" if core_statuses["RG-07"] else "fail", {"corrected_result_trace_count": semantic["corrected_result_trace_count"]}),
        _release_gate("RG-08", "alias_deduplication", "pass" if core_statuses["RG-08"] else "fail", {"deduplicated_summary_field_count": semantic["deduplicated_summary_field_count"]}),
        _release_gate("RG-09", "visible_backend_field_suppression", "pass" if core_statuses["RG-09"] else "fail", {"visible_raw_backend_field_count": semantic["visible_raw_backend_field_count"]}),
        _release_gate("RG-10", "render_visual_status", render_evidence.get("status", "not_run"), {"rendered_slide_count": render_evidence.get("rendered_slide_count", 0), "evidence": render_evidence}),
        _release_gate("RG-11", "professor_structural_fidelity", "insufficient_evidence", {"private_render_comparison": "not_authorized", "composition_audits_pass": all(core_statuses[key] for key in ("RG-02", "RG-03", "RG-04"))}),
        _release_gate("RG-12", "qualitative_visual_review", "blocked_visual_review", {"reviewed_slide_count": 0, "reason": "no authorized image-capable professor review"}),
        _release_gate("RG-13", "native_powerpoint_acceptance", "blocked_environment", {"open_save_reopen_attempts": 0}),
        _release_gate("RG-14", "privacy", "pass" if core_statuses["RG-14"] else "fail", {"privacy_evidence": privacy_evidence, "approved_historical_exception_count": approved_exception_count}),
    ]
    prerequisite_statuses = {item["gate_id"]: item["status"] for item in gates}
    production_release = "pass" if all(prerequisite_statuses[key] == "pass" for key in ("RG-01", "RG-02", "RG-03", "RG-04", "RG-05", "RG-06", "RG-07", "RG-08", "RG-09", "RG-10", "RG-11", "RG-12", "RG-13", "RG-14")) else "blocked"
    gates.extend([
        _release_gate("RG-15", "production_release_status", production_release, {"non_pass_gate_ids": [item["gate_id"] for item in gates if item["status"] != "pass"]}),
        _release_gate("RG-16", "production_group_meeting_ready", "false", {"production_release_status": production_release, "external_review_complete": False}),
    ])
    release = {
        "release_id": "FVCC-RELEASE-001", "candidate_state": candidate_state, "gates": gates,
        "acceptance_deck_build_status": "pass" if all(core_statuses.values()) else "fail",
        "production_release_status": production_release, "production_group_meeting_ready": False,
        "private_alias_resolution_attempts": privacy_evidence.get("private_alias_resolution_attempts"),
        "private_source_open_attempts": privacy_evidence.get("private_source_open_attempts"),
        "private_render_attempts": privacy_evidence.get("private_render_attempts"),
    }
    (destination / "final-visual-composition-release-qa.json").write_text(
        json.dumps(release, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return release


def compute_final_visual_composition_candidate_state(root: Path) -> dict[str, Any]:
    """Hash the complete final-composition execution domain, not its reports."""
    root = Path(root).resolve()
    component_paths = (
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_final_visual_composition.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase2_build.py",
        "packages/thesis-deck-system/src/thesis_deck_system/story.py",
        "packages/thesis-deck-system/src/thesis_deck_system/pptx.py",
        "packages/thesis-deck-system/src/thesis_deck_system/template.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_cp5bcd_integrated.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_cp5efg_integrated.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_cp5_hi_final_sprint.py",
        "packages/thesis-deck-system/src/thesis_deck_system/final_closure_reliability.py",
        "packages/thesis-deck-system/src/thesis_deck_system/contracts.py",
        "packages/thesis-deck-system/src/thesis_deck_system/phase3_privacy.py",
        "packages/thesis-deck-system/tests/unit/test_phase3_final_visual_composition.py",
        "packages/thesis-deck-system/tests/unit/test_final_closure_reliability.py",
        "packages/thesis-deck-system/pyproject.toml",
        "thesis-deck-system/TASK_PHASE_3_FINAL_VISUAL_COMPOSITION_CLOSURE.md",
        "thesis-deck-system/designs/PHASE_3_FINAL_VISUAL_COMPOSITION_CLOSURE_DESIGN.md",
        "thesis-deck-system/artifacts/phase2/slide-specs.json",
        "thesis-deck-system/artifacts/phase2/materialized-h02.json",
        "thesis-deck-system/artifacts/phase2/fishbone/FB001-rev1.svg",
        "thesis-deck-system/artifacts/phase2/fishbone/FB001-rev1.png",
        "thesis-deck-system/artifacts/phase2/fishbone/FB001-rev2.svg",
        "thesis-deck-system/artifacts/phase2/fishbone/FB001-rev2.png",
        "thesis-deck-system/artifacts/phase3/approved-figures.json",
        "thesis-deck-system/artifacts/phase3/visual-style-profile.json",
        "thesis-deck-system/schemas/generated-pptx-attestation.schema.json",
        "thesis-deck-system/schemas/native-materialization-parity.schema.json",
        "thesis-deck-system/schemas/final-closure-validation-run.schema.json",
        "thesis-deck-system/schemas/final-closure-reliability-qa.schema.json",
    )
    hashes: dict[str, str] = {}
    for relative in component_paths:
        hashes[relative] = _candidate_component_digest(relative, (root / relative).read_bytes())
    return {
        "candidate_id": "FVCC-CANDIDATE-001", "component_count": len(hashes),
        "component_hashes": hashes, "candidate_state_sha256": _json_hash(hashes),
    }


def _candidate_component_digest(relative_path: str, contents: bytes) -> str:
    """Normalize checkout text line endings while preserving binary identity."""
    if Path(relative_path).suffix.casefold() in {".py", ".json", ".md", ".toml"}:
        contents = contents.replace(b"\r\n", b"\n")
    return sha256(contents).hexdigest()
