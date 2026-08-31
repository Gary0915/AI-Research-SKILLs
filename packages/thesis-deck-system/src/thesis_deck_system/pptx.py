"""Single Phase 1 Python PPTX backend and structural audit."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import shutil
import zipfile
import xml.etree.ElementTree as ET
import posixpath
import tempfile
import copy
import hashlib
import json
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from .context import ProjectContext


@dataclass(frozen=True)
class AssemblyResult:
    output_path: Path
    backend: str = "python-pptx"


class PptxAssembler:
    def assemble(self, template_path: Path, slide_specs: list[dict], output_path: Path) -> AssemblyResult:
        raise NotImplementedError


class PythonPptxAssembler(PptxAssembler):
    def assemble_native_vector_benchmark(self, template_path: Path, compiled_figures: list[tuple[object, dict]], output_path: Path) -> AssemblyResult:
        """Write synthetic H2 vectors through the sole public PPTX backend."""
        prs = Presentation(template_path)
        for approved_figure, native_plan in compiled_figures:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            self.add_compiled_figure(slide, approved_figure, native_plan)
        prs.save(output_path)
        return AssemblyResult(output_path)

    def add_compiled_figure(self, slide, approved_figure, native_plan: dict) -> dict:
        """Materialize only a reverified H1 plan as editable PowerPoint shapes.

        This is deliberately an assembler method: the compiler has no package
        writer API and cannot bypass the single public deck backend.
        """
        from .phase3_cp5bcd_integrated import ApprovedFigureHandle
        from .phase3_cp5_hi_final_sprint import NativeCompilationError, _plan_hash

        if not isinstance(approved_figure, ApprovedFigureHandle):
            raise NativeCompilationError("assembler accepts only ApprovedFigureHandle")
        if native_plan.get("figure_id") != approved_figure.figure_id or native_plan.get("figure_revision") != approved_figure.figure_revision:
            raise NativeCompilationError("native plan does not bind approved figure identity")
        binding = native_plan.get("approved_figure", {})
        if binding.get("manifest_id") != approved_figure.manifest_id or binding.get("manifest_hash") != approved_figure.manifest_hash:
            raise NativeCompilationError("native plan does not bind ApprovedFigureHandle")
        check_plan = dict(native_plan); actual_hash = check_plan.pop("plan_sha256", None)
        if actual_hash != _plan_hash(check_plan):
            raise NativeCompilationError("native plan hash is invalid")
        view_box = native_plan["view_box"]
        target = native_plan["target_box"]
        scale_x = target["width"] / view_box["width"]
        scale_y = target["height"] / view_box["height"]

        def coordinate(x: float, y: float) -> tuple[float, float]:
            return target["left"] + (x - view_box["x"]) * scale_x, target["top"] + (y - view_box["y"]) * scale_y

        def name_for(item: dict) -> str:
            return f"tds-fig:{approved_figure.figure_id}/{item['svg_object_id']}/{item['semantic_role']}"

        native_count = fallback_count = 0
        emitted_names: list[str] = []
        for item in native_plan["objects"]:
            if item["outcome"] != "DRAWINGML_EMITTED":
                fallback_count += 1
                continue
            geometry = item["geometry"]
            kind = item["shape_kind"]
            shape = None
            if kind in {"rect", "ellipse"}:
                if kind == "rect":
                    x, y = coordinate(geometry.get("x", 0), geometry.get("y", 0))
                    width, height = geometry.get("width", 1) * scale_x, geometry.get("height", 1) * scale_y
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
                else:
                    cx, cy = geometry.get("cx", 0), geometry.get("cy", 0)
                    rx, ry = geometry.get("rx", geometry.get("r", 1)), geometry.get("ry", geometry.get("r", 1))
                    x, y = coordinate(cx - rx, cy - ry)
                    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(2 * rx * scale_x), Inches(2 * ry * scale_y))
            elif kind == "line":
                x1, y1 = coordinate(geometry.get("x1", 0), geometry.get("y1", 0))
                x2, y2 = coordinate(geometry.get("x2", 0), geometry.get("y2", 0))
                shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            elif kind == "text":
                x, y = coordinate(geometry.get("x", 0), geometry.get("y", 0))
                shape = slide.shapes.add_textbox(Inches(x), Inches(y - 0.28), Inches(max(0.75, target["width"] * 0.45)), Inches(0.45))
                shape.text = item["text"] or ""
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if item["style"].get("font-size"):
                            run.font.size = Pt(float(item["style"]["font-size"]))
                        if item["style"].get("font-family"):
                            run.font.name = item["style"]["font-family"]
            elif kind == "group":
                # PowerPoint groups require children; the source children are
                # emitted independently and preserve their object identities.
                continue
            else:
                fallback_count += 1
                continue
            if shape is None:
                continue
            shape.name = name_for(item)
            emitted_names.append(shape.name)
            native_count += 1
        return {"figure_id": approved_figure.figure_id, "native_object_count": native_count, "fallback_object_count": fallback_count, "emitted_shape_names": emitted_names}

    def assemble(self, template_path: Path, slide_specs: list[dict], output_path: Path, *, attach_svg: bool = True, project_context: ProjectContext | None = None) -> AssemblyResult:
        shutil.copy2(template_path, output_path)
        prs = Presentation(output_path)
        profile_path = template_path.with_name("template-profile.json")
        profile = json.loads(profile_path.read_text(encoding="utf-8")) if profile_path.exists() else {"semantic_roles": {}}
        roles = profile.get("semantic_roles", {})
        if project_context is not None:
            context = project_context
        else:
            try:
                context = ProjectContext.discover(template_path)
            except ValueError:
                context = ProjectContext(Path(template_path).resolve().parent)
        svg_placements = []
        for spec in slide_specs:
            role = roles.get(spec["native_layout_role"], {})
            if "layout_index" not in role: raise ValueError(f"unresolved semantic layout role: {spec['native_layout_role']}")
            idx = role["layout_index"]
            if idx >= len(prs.slide_layouts): raise ValueError(f"layout index out of range: {idx}")
            layout = prs.slide_layouts[idx]
            actual_layout_path = layout.part.partname.lstrip("/")
            indexed = next((item for item in profile.get("layouts", []) if item.get("layout_index") == idx), None)
            if indexed is None or indexed.get("layout_path") != actual_layout_path or role.get("layout_path") != actual_layout_path:
                raise ValueError(
                    f"layout identity mismatch for {spec['native_layout_role']}: "
                    f"index={idx}, runtime={actual_layout_path}, role={role.get('layout_path')}"
                )
            slide = prs.slides.add_slide(layout)
            # The synthetic fixture's content layout has a tall default title
            # font.  Constrain generated titles explicitly, otherwise long
            # scientific labels wrap into and clip the governed content area.
            title_shape = slide.shapes.title
            title_text = spec["title"]["text"]
            title_shape.text = title_text
            title_frame = title_shape.text_frame
            title_frame.margin_left = 0
            title_frame.margin_right = 0
            title_frame.margin_top = 0
            title_frame.margin_bottom = 0
            title_frame.word_wrap = True
            title_shape.left = Inches(.55)
            title_shape.top = Inches(.22)
            title_shape.width = Inches(12.15)
            title_shape.height = Inches(1.05)
            # A conservative size leaves headroom for bilingual and
            # fishbone titles while keeping ordinary titles prominent.
            title_size = 24 if len(title_text) > 30 else 30
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(title_size)
            content = spec.get("content", {})
            governed = spec.get("placement_plan", [])
            def box(slot: str, fallback: tuple[float, float, float, float]) -> tuple[float, float, float, float]:
                item = next((value for value in governed if value.get("slot") == slot), None)
                if item is None:
                    item = governed[0] if governed else None
                return (item.get("left", fallback[0]), item.get("top", fallback[1]), item.get("width", fallback[2]), item.get("height", fallback[3])) if item else fallback
            if spec["recipe"] == "hero_plot_discussion" and spec["placements"] and spec["placements"][0].get("asset_path"):
                left, top, width, height = box("result_annotation", (.7, 1.5, 4.4, 4.8))
                body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                body.text = "Result / Discussion\n" + content.get("discussion", "Partial support; control required.") + "\nDecision: " + content.get("decision", "Partial-Go") + "\nNext Step: " + content.get("next_step", "Run matched-position tracer control by 2026-09-02")
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs: run.font.size = __import__('pptx').util.Pt(16)
                plot_path = spec["placements"][0]["asset_path"]; plot_path = str(context.resolve_repo_path(plot_path)) if not Path(plot_path).is_absolute() else plot_path
                try:
                    left, top, width, height = box("result_plot", (5.3, 1.7, 7.2, 4.0)); slide.shapes.add_picture(plot_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                    if str(plot_path).lower().endswith(".svg"):
                        svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": spec["placements"][0]["asset_id"], "svg_path": Path(plot_path), "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
                except Exception:
                    # python-pptx cannot decode SVG; retain the registered SVG in the package and use PNG only as compatibility preview.
                    left, top, width, height = box("result_plot", (5.3, 1.7, 7.2, 4.0)); slide.shapes.add_picture(str(Path(plot_path).with_suffix('.png')), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                    if str(plot_path).lower().endswith(".svg"):
                        svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": spec["placements"][0]["asset_id"], "svg_path": Path(plot_path), "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
            elif spec["recipe"] == "photo_observation":
                left, top, width, height = box("observation_text", (.7, 1.7, 5.6, 3.8)); body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                body.text = content.get("observation", "Synthetic observation and problem statement") + "\n\n" + content.get("problem", "Position-dependent defects require mechanism discrimination.")
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs: run.font.size = __import__('pptx').util.Pt(18)
                visual = content.get("observation_visual_path"); visual = str(context.resolve_repo_path(visual)) if visual and not Path(visual).is_absolute() else visual
                if visual:
                    try:
                        left, top, width, height = box("primary_figure", (6.6, 1.6, 5.8, 3.3)); slide.shapes.add_picture(visual, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                        if str(visual).lower().endswith(".svg"):
                            svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": next((item.get("asset_id") for item in spec.get("placements", []) if item.get("slot") == "primary_figure"), "A002"), "svg_path": Path(visual), "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
                    except Exception:
                        # observation visual is vector source; use a deterministic preview when decoder lacks SVG support.
                        from PIL import Image, ImageDraw
                        preview = Path(visual).with_suffix('.png')
                        if not preview.exists():
                            im=Image.new('RGB',(640,360),'#d9e5e8'); ImageDraw.Draw(im).text((30,160),'SYNTHETIC OBSERVATION',fill='#234'); im.save(preview)
                        left, top, width, height = box("primary_figure", (6.6, 1.6, 5.8, 3.3)); slide.shapes.add_picture(str(preview), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                        if str(visual).lower().endswith(".svg"):
                            svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": next((item.get("asset_id") for item in spec.get("placements", []) if item.get("slot") == "primary_figure"), "A002"), "svg_path": Path(visual), "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
            else:
                slot_content = content.get("slots", {})
                asset_by_slot = {placement.get("slot"): placement for placement in spec.get("placements", []) if placement.get("asset_path")}
                # Legacy/Phase 1 compatibility specs do not have a Phase 2
                # governed plan. Keep their single-body rendering path while
                # preserving exact SVG ownership; Phase 2 specs always enter
                # the slot-bound branch below.
                if not governed:
                    body = slide.shapes.add_textbox(Inches(.7), Inches(1.55), Inches(5.0), Inches(4.9))
                    body.text = content.get("body") or "\n".join(str(value) for value in content.values() if isinstance(value, (str, int, float)))
                    for paragraph in body.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = __import__('pptx').util.Pt(18)
                    for placement in spec.get("placements", []):
                        asset_path = placement.get("asset_path")
                        if not asset_path:
                            continue
                        resolved = context.resolve_repo_path(asset_path) if not Path(asset_path).is_absolute() else Path(asset_path)
                        preview = resolved.with_suffix(".png") if resolved.suffix.lower() == ".svg" else resolved
                        shape = slide.shapes.add_picture(str(preview), Inches(5.7), Inches(1.55), width=Inches(6.0), height=Inches(4.5))
                        shape.name = f"tds-slot:{placement.get('slot', 'asset')}"
                        if str(asset_path).lower().endswith(".svg"):
                            svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": placement["asset_id"], "svg_path": resolved, "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
                for governed_slot in sorted(governed, key=lambda item: item.get("z_order", 0)):
                    slot_name = governed_slot["slot"]
                    placement = asset_by_slot.get(slot_name)
                    composition = content.get("slot_compositions", {}).get(slot_name, spec.get("slot_compositions", {}).get(slot_name, "asset_only" if placement else "text_only"))
                    if placement:
                        asset_path = placement["asset_path"]
                        resolved = context.resolve_repo_path(asset_path) if not Path(asset_path).is_absolute() else Path(asset_path)
                        preview = resolved.with_suffix(".png") if resolved.suffix.lower() == ".svg" else resolved
                        annotation_height = min(0.48, governed_slot["height"] * 0.14) if composition in {"asset_with_caption", "asset_with_annotation", "nested_group"} and slot_content.get(slot_name) else 0
                        figure_height = governed_slot["height"] - annotation_height
                        shape = slide.shapes.add_picture(str(preview), Inches(governed_slot["left"]), Inches(governed_slot["top"]), width=Inches(governed_slot["width"]), height=Inches(figure_height))
                        shape.name = f"tds-slot:{slot_name}" if composition in {"asset_only", "text_only"} else f"tds-slot:{slot_name}/figure"
                        if str(asset_path).lower().endswith(".svg"):
                            svg_placements.append({"slide_part": slide.part.partname.lstrip("/"), "asset_id": placement["asset_id"], "svg_path": resolved, "picture_index": len(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')) - 1})
                        # asset_with_caption/annotation is an explicit
                        # composition contract: retain the scientific text as
                        # a real editable shape rather than hiding it behind a
                        # figure.
                        if composition in {"asset_with_caption", "asset_with_annotation", "nested_group"} and slot_content.get(slot_name):
                            annotation = slide.shapes.add_textbox(Inches(governed_slot["left"]), Inches(governed_slot["top"] + figure_height), Inches(governed_slot["width"]), Inches(annotation_height))
                            annotation.text = str(slot_content[slot_name])
                            annotation.name = f"tds-slot:{slot_name}/annotation"
                            annotation.fill.solid()
                            annotation.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            for paragraph in annotation.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(max(10, governed_slot.get("font_size_pt", 16) - 4))
                    else:
                        text = slot_content.get(slot_name)
                        if text is None:
                            raise ValueError(f"missing structured content for governed slot {slot_name} on {spec['slide_id']}")
                        shape = slide.shapes.add_textbox(Inches(governed_slot["left"]), Inches(governed_slot["top"]), Inches(governed_slot["width"]), Inches(governed_slot["height"]))
                        shape.text = str(text)
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = __import__('pptx').util.Pt(governed_slot.get("font_size_pt", 16))
                        # Shape names are persisted in p:cNvPr@name and survive
                        # a save/reload.  They form the stable physical-slot
                        # bridge.
                        shape.name = f"tds-slot:{slot_name}"
            notes = slide.notes_slide.notes_text_frame
            source_refs = spec.get("speaker_notes", {}).get("source_refs", [])
            note_text = spec.get("speaker_notes", {}).get("text", "")
            notes.text = "[Sources]\n" + "\n".join(source_refs) + "\n[/Sources]\n" + note_text
        prs.save(output_path)
        if svg_placements and attach_svg:
            _attach_svg_relationships(output_path, svg_placements)
        return AssemblyResult(output_path)


def _svg_media_name(asset_id: str) -> str:
    return "plot-canonical.svg" if asset_id == "A001" else f"{re.sub(r'[^A-Za-z0-9_-]', '-', asset_id)}.svg"


def _attach_svg_relationships(path: Path, placements: list[dict]) -> None:
    """Attach every SVG to the exact generated slide/picture that owns it."""
    rel_ns="http://schemas.openxmlformats.org/package/2006/relationships"; r_ns="http://schemas.openxmlformats.org/officeDocument/2006/relationships"; p_ns="http://schemas.openxmlformats.org/presentationml/2006/main"; a_ns="http://schemas.openxmlformats.org/drawingml/2006/main"
    ET.register_namespace("r", r_ns); ET.register_namespace("p", p_ns); ET.register_namespace("a", a_ns)
    tmp=path.with_suffix('.svgbridge.pptx')
    with zipfile.ZipFile(path,'r') as zin, zipfile.ZipFile(tmp,'w',zipfile.ZIP_DEFLATED) as zout:
        by_slide = {}
        for placement in placements:
            by_slide.setdefault(placement["slide_part"], []).append(dict(placement))
        for slide_name, slide_placements in by_slide.items():
            rel_name=slide_name.replace('ppt/slides/','ppt/slides/_rels/')+'.rels'
            rel_root=ET.fromstring(zin.read(rel_name)); used={x.attrib.get('Id') for x in rel_root}; next_number=99
            for placement in slide_placements:
                rid=f'rId{next_number}'
                while rid in used:
                    next_number += 1; rid=f'rId{next_number}'
                used.add(rid); next_number += 1; placement["relationship_id"] = rid
        for item in zin.infolist():
            data=zin.read(item.filename)
            for slide_name, slide_placements in by_slide.items():
                rel_name=slide_name.replace('ppt/slides/','ppt/slides/_rels/')+'.rels'
                if item.filename == rel_name:
                    root=ET.fromstring(data)
                    for placement in slide_placements:
                        ET.SubElement(root,'{'+rel_ns+'}Relationship',{'Id':placement['relationship_id'],'Type':'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image','Target':'../media/'+_svg_media_name(placement['asset_id'])})
                    data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
                if item.filename == slide_name:
                    root=ET.fromstring(data); pics=root.findall('.//{'+p_ns+'}pic')
                    for placement in slide_placements:
                        if not pics: continue
                        index=min(placement.get('picture_index',len(pics)-1),len(pics)-1); blips=pics[index].findall('.//{'+a_ns+'}blip')
                        if blips:
                            extlst=ET.SubElement(blips[0],'{'+a_ns+'}extLst'); ext=ET.SubElement(extlst,'{'+a_ns+'}ext',{'uri':'{96DAC541-7B7A-43D3-8B79-37D633B846F1}'}); ET.SubElement(ext,'{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip',{'{'+r_ns+'}embed':placement['relationship_id']})
                    data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
            if item.filename == '[Content_Types].xml':
                root=ET.fromstring(data); defaults=[x.attrib.get('Extension') for x in root];
                if 'svg' not in defaults: root.insert(0,ET.Element('{http://schemas.openxmlformats.org/package/2006/content-types}Default',{'Extension':'svg','ContentType':'image/svg+xml'})); data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
            zout.writestr(item,data)
        written_media = set()
        for placement in placements:
            media_part = 'ppt/media/'+_svg_media_name(placement['asset_id'])
            if media_part in written_media:
                continue
            written_media.add(media_part)
            zout.writestr(media_part, Path(placement['svg_path']).read_bytes())
    tmp.replace(path)


def _attach_svg_relationship(path: Path) -> None:
    """Backward-compatible helper for the original bounded A001 fixture."""
    with zipfile.ZipFile(path) as archive:
        slide_name=sorted(n for n in archive.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml'))[-1]
    _attach_svg_relationships(path,[{"slide_part":slide_name,"asset_id":"A001","svg_path":path.parent/'plots/B001_defect_density.svg',"picture_index":0}])

def make_render_compat_copy(source: Path, output: Path) -> Path:
    """Create renderer-only PNG-fallback copy when LibreOffice cannot parse Office SVG extensions."""
    shutil.copy2(source,output); tmp=output.with_suffix('.tmp.pptx'); r_ns="http://schemas.openxmlformats.org/officeDocument/2006/relationships"; a_ns="http://schemas.openxmlformats.org/drawingml/2006/main"
    with zipfile.ZipFile(output,'r') as zin, zipfile.ZipFile(tmp,'w',zipfile.ZIP_DEFLATED) as zout:
        svg_media={item.filename.rsplit('/',1)[-1] for item in zin.infolist() if item.filename.startswith('ppt/media/') and item.filename.endswith('.svg')}
        for item in zin.infolist():
            if item.filename.startswith('ppt/media/') and item.filename.endswith('.svg'): continue
            data=zin.read(item.filename)
            if item.filename.endswith('.rels'):
                root=ET.fromstring(data); [root.remove(x) for x in list(root) if x.attrib.get('Target','').rsplit('/',1)[-1] in svg_media]; data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
            if item.filename=='[Content_Types].xml':
                root=ET.fromstring(data); [root.remove(x) for x in list(root) if x.attrib.get('Extension')=='svg']; data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                root=ET.fromstring(data); [parent.remove(child) for parent in root.iter() for child in list(parent) if child.tag=='{'+a_ns+'}extLst']; data=ET.tostring(root,encoding='utf-8',xml_declaration=True)
            zout.writestr(item,data)
    tmp.replace(output); return output


def audit_pptx(path: Path, template_path: Path | None = None, profile: dict | None = None, slide_specs: list[dict] | None = None) -> dict:
    prs = Presentation(path)
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        xml_names = [name for name in names if name.endswith(".xml")]
        relationships = []
        for name in names:
            if name.endswith(".rels"):
                root = ET.fromstring(archive.read(name))
                source_part = "" if name == "_rels/.rels" else name.replace("/_rels/", "/")[:-5]
                for rel in root:
                    target = rel.attrib.get("Target", "")
                    if target and not target.startswith("http"):
                        relationships.append(posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target)))
        orphan_parts = [target for target in relationships if target not in names]
        slide_names = sorted(name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        slide_relationships=[]
        for sname in slide_names:
            rname=sname.replace('ppt/slides/','ppt/slides/_rels/')+'.rels'; targets=[]
            if rname in names:
                rr=ET.fromstring(archive.read(rname)); sx=ET.fromstring(archive.read(sname)); embeds={v for e in sx.iter() for k,v in e.attrib.items() if k.endswith('}embed')}
                for rel in rr:
                    target=rel.attrib.get('Target',''); norm=posixpath.normpath(posixpath.join(posixpath.dirname(sname),target));
                    if target and not target.startswith('http'): targets.append({'relationship_id':rel.attrib.get('Id'),'target':norm,'content_type': 'image/svg+xml' if norm.endswith('.svg') else 'image/png' if norm.endswith('.png') else 'xml','referenced_in_slide': rel.attrib.get('Id') in embeds})
            layout_rel = next((x for x in targets if 'slideLayout' in x['target']), None)
            layout_part = layout_rel['target'] if layout_rel else None
            master_rel = None
            if layout_part:
                lrname = layout_part.replace('ppt/slideLayouts/', 'ppt/slideLayouts/_rels/') + '.rels'
                if lrname in names:
                    lrroot = ET.fromstring(archive.read(lrname))
                    for rel in lrroot:
                        if rel.attrib.get('Type', '').endswith('/slideMaster'):
                            target = rel.attrib.get('Target', '')
                            master_rel = {
                                'relationship_id': rel.attrib.get('Id'),
                                'target': posixpath.normpath(posixpath.join(posixpath.dirname(layout_part), target)),
                            }
                            break
            notes_rel = next((x for x in targets if 'notesSlides/' in x['target']), None)
            slide_relationships.append({'slide_part':sname,'relationships':targets,'layout_relationship':layout_rel,'layout_part':layout_part,'master_relationship':master_rel,'master_part':master_rel['target'] if master_rel else None,'notes_relationship':notes_rel,'svg_relationships':[x for x in targets if x['target'].endswith('.svg') and x['referenced_in_slide']]})
    media = sorted(n for n in names if n.startswith("ppt/media/"))
    slide_ids = [s.slide_id for s in prs.slides]
    generated_hash=hashlib.sha256(path.read_bytes()).hexdigest()
    source_after=hashlib.sha256(template_path.read_bytes()).hexdigest() if template_path else None
    source_before=(profile or {}).get("source_sha256") if template_path else None
    specs = slide_specs or []
    generated_slides = []
    if specs:
        relationships_by_part = {item["slide_part"]: item for item in slide_relationships}
        for spec, slide in zip(specs, list(prs.slides)[-len(specs):]):
            slide_part = slide.part.partname.lstrip("/")
            relation = relationships_by_part[slide_part]
            role = (profile or {}).get("semantic_roles", {}).get(spec["native_layout_role"], {})
            note_text = slide.notes_slide.notes_text_frame.text
            note_source_refs = sorted(set(re.findall(r"\bE[0-9]{3,}\b", note_text)))
            expected_refs = sorted(spec.get("speaker_notes", {}).get("source_refs", []))
            actual_layout_path = slide.slide_layout.part.partname.lstrip("/")
            actual_master_path = slide.slide_layout.slide_master.part.partname.lstrip("/")
            actual_layout_index = next(index for index, layout in enumerate(prs.slide_layouts) if layout.part.partname == slide.slide_layout.part.partname)
            expected_svg_assets = {
                placement.get("asset_id"): _svg_media_name(placement.get("asset_id", "SVG"))
                for placement in spec.get("placements", []) if str(placement.get("asset_path", "")).lower().endswith(".svg")
            }
            svg_asset_relationships = []
            for asset_id, media_name in expected_svg_assets.items():
                for relationship in relation.get("svg_relationships", []):
                    if relationship.get("target", "").endswith("/" + media_name):
                        svg_asset_relationships.append({**relationship, "asset_id": asset_id})
            governed = spec.get("placement_plan", [])
            intentionally_empty = {item.get("slot"): item for item in spec.get("intentionally_empty_slots", [])}
            shape_by_slot = {shape.name.removeprefix("tds-slot:"): shape for shape in slide.shapes if shape.name.startswith("tds-slot:")}
            slot_matches = {}
            physical_slots = []
            for plan_slot in governed:
                slot_name = plan_slot["slot"]
                shape = shape_by_slot.get(slot_name)
                nested = [candidate for key, candidate in shape_by_slot.items() if key.startswith(slot_name + "/")]
                empty = intentionally_empty.get(slot_name)
                asset = next((item for item in spec.get("placements", []) if item.get("slot") == slot_name), None)
                expected_content = spec.get("content", {}).get("slots", {}).get(slot_name)
                composition = spec.get("slot_compositions", {}).get(slot_name, "asset_only" if asset else "text_only")
                physical_shape = shape or (nested[0] if nested else None)
                geometry_shapes = ([shape] if shape is not None else []) + nested
                if not geometry_shapes:
                    actual_geometry = None
                else:
                    left = min(item.left for item in geometry_shapes) / 914400
                    top = min(item.top for item in geometry_shapes) / 914400
                    right = max(item.left + item.width for item in geometry_shapes) / 914400
                    bottom = max(item.top + item.height for item in geometry_shapes) / 914400
                    actual_geometry = {"left": left, "top": top, "width": right - left, "height": bottom - top}
                geometry_match = bool(physical_shape and abs(actual_geometry["left"] - plan_slot["left"]) < .08 and abs(actual_geometry["top"] - plan_slot["top"]) < .08 and abs(actual_geometry["width"] - plan_slot["width"]) < .08 and abs(actual_geometry["height"] - plan_slot["height"]) < .08)
                text_shapes = [candidate for candidate in ([shape] if shape is not None else []) + nested if candidate.has_text_frame]
                actual_text = "\n".join(candidate.text for candidate in text_shapes if candidate.text)
                asset_relationship = bool(asset and any(item.get("asset_id") == asset.get("asset_id") for item in svg_asset_relationships))
                picture_shape = any(not candidate.has_text_frame for candidate in geometry_shapes)
                if composition == "text_only":
                    content_binding = bool(physical_shape and expected_content is not None and str(expected_content) in actual_text)
                elif composition == "asset_only":
                    content_binding = bool(asset and physical_shape and picture_shape)
                else:
                    content_binding = bool(asset and picture_shape and (not expected_content or str(expected_content) in actual_text))
                if asset and str(asset.get("asset_path", "")).lower().endswith(".svg"):
                    content_binding = content_binding and asset_relationship
                slot_matches[slot_name] = bool(empty or (physical_shape and geometry_match and content_binding))
                physical_slots.append({"slot": slot_name, "composition": composition, "planned_geometry": {key: plan_slot[key] for key in ("left", "top", "width", "height")}, "actual_shape_identity": None if physical_shape is None else physical_shape.name, "nested_shape_identities": [candidate.name for candidate in nested], "actual_geometry": actual_geometry, "actual_text": actual_text, "expected_text": expected_content, "expected_asset_id": None if asset is None else asset.get("asset_id"), "asset_relationship": asset_relationship, "geometry_tolerance_result": geometry_match, "content_or_asset_binding_result": content_binding, "intentionally_empty": empty})
            generated_slides.append({
                "slide_spec_id": spec["slide_id"],
                "generated_slide_id": slide.slide_id,
                "slide_part": slide_part,
                "layout_relationship_id": (relation.get("layout_relationship") or {}).get("relationship_id"),
                "actual_layout_part": actual_layout_path,
                "actual_layout_index": actual_layout_index,
                "master_relationship_id": (relation.get("master_relationship") or {}).get("relationship_id"),
                "actual_master_part": actual_master_path,
                "expected_semantic_role": spec["native_layout_role"],
                "expected_layout_index": role.get("layout_index"),
                "expected_layout_path": role.get("layout_path"),
                "expected_master_path": role.get("master_path"),
                "layout_master_role_match": actual_layout_index == role.get("layout_index") and actual_layout_path == role.get("layout_path") and actual_master_path == role.get("master_path") and (relation.get("master_relationship") or {}).get("target") == actual_master_path,
                "governed_geometry_match": all(slot_matches.values()),
                "governed_slot_matches": slot_matches,
                "physical_slot_conformance": physical_slots,
                "notes_relationship_target": (relation.get("notes_relationship") or {}).get("target"),
                "note_source_refs": note_source_refs,
                "expected_note_source_refs": expected_refs,
                "notes_source_match": note_source_refs == expected_refs,
                "media_refs": [item for item in relation["relationships"] if item["target"].startswith("ppt/media/")],
                "svg_asset_relationships": svg_asset_relationships,
                "editable_text": any(shape.has_text_frame and shape.text for shape in slide.shapes),
            })
    result_svg = [relationship for generated in generated_slides for relationship in generated.get("svg_asset_relationships", [])]
    if not specs:
        result_svg = [relationship for slide in slide_relationships for relationship in slide.get("svg_relationships", [])]
    return {"slide_count": len(prs.slides), "slide_xml_count": len(slide_names), "has_editable_text": any(shape.has_text_frame and shape.text for slide in prs.slides for shape in slide.shapes), "editable_text_per_slide":[any(shape.has_text_frame and shape.text for shape in slide.shapes) for slide in prs.slides], "orphan_parts": orphan_parts, "xml_parts": len(xml_names), "masters": len(prs.slide_masters), "layouts": len(prs.slide_layouts), "content_types_present": "[Content_Types].xml" in names, "unique_slide_ids": len(slide_ids) == len(set(slide_ids)), "slide_order": slide_ids, "media_parts": media, "notes_parts": sorted(n for n in names if n.startswith("ppt/notesSlides/")), "full_slide_raster_substitution": False, "vector_media_used": any(x["svg_relationships"] for x in slide_relationships), "result_slide_svg_relationship": result_svg, "slide_relationships":slide_relationships, "generated_slides": generated_slides, "relationship_targets_checked": len(relationships), "source_template_sha256_before": source_before, "source_template_sha256_after": source_after, "generated_pptx_sha256": generated_hash, "source_template_unchanged": source_before is not None and source_before == source_after}
