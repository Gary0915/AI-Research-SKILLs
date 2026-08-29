"""Synthetic-first RED tests for Phase 3 Checkpoint 2."""

from __future__ import annotations

import hashlib
from pathlib import Path
import subprocess

import pytest
from pptx import Presentation

from thesis_deck_system.contracts import SchemaRegistry
from thesis_deck_system.phase3_checkpoint2 import (
    Checkpoint2PolicyViolation,
    Checkpoint2Run,
    LocalPrivateAliasResolver,
    build_checkpoint2,
    sanitize_body_descriptor,
    sanitize_shell_descriptor,
    validate_checkpoint2_qa,
)
from thesis_deck_system.phase3_checkpoint2 import (
    _classify_structural_family,
    _color_role,
    _compose_transform,
    _connector_semantics,
    _metric_observation,
)
from thesis_deck_system.phase3_contracts import (
    canonical_observation_catalogs,
    validate_observation_visual_binding,
)
from thesis_deck_system.phase3_privacy import RepositoryPrivacyScanner
from thesis_deck_system.phase3_privacy import LEGACY_EXCEPTION_PATH, LEGACY_EXCEPTION_BLOB_SHA


ROOT = Path(__file__).resolve().parents[4]
SCHEMAS = ROOT / "thesis-deck-system" / "schemas"
SHA = "a" * 64
ALIASES = (
    "private://template_primary_1",
    "private://layout_exemplar_2",
    "private://template_primary_3",
)


def _synthetic_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
    slide.shapes.add_shape(1, 914400, 2286000, 2743200, 1371600)
    presentation.save(path)
    return path


def _evidence(evidence_id: str, kind: str, *, verified: str = "verified") -> dict:
    return {
        "schema_version": "1.0.0", "evidence_id": evidence_id, "kind": kind,
        "title": "Synthetic evidence", "provenance": "synthetic_test_only",
        "source": {"source_id": "S001", "uri": "fixtures/synthetic.json", "sha256": SHA},
        "claim_support_refs": [], "claim_contradict_refs": [], "scope": {},
        "verification": {"status": verified},
    }


def _plot(evidence_id: str) -> dict:
    return {
        "schema_version": "3.0.0", "figure_output_id": "FOM001", "figure_id": "FIG001",
        "figure_type": "scientific_plot", "primary_artifact_kind": "svg_vector",
        "renderer": "synthetic_renderer", "source_spec_sha256": SHA,
        "provenance_refs": [evidence_id], "style_profile_ref": "VSP001",
        "evidence_status": "empirical", "primary_artifact": {
            "path": "artifacts/phase3/synthetic.svg", "sha256": SHA,
            "data_provenance_refs": [evidence_id],
        }, "output_part_lineage": ["generated"],
    }


def _binding(evidence_id: str = "E001") -> dict:
    return {
        "observation_id": "OBS001", "empirical_evidence_required": True,
        "observation_evidence_ref": evidence_id, "observation_output_ref": "FOM001",
        "evidence_refs": [evidence_id], "auxiliary_visuals": [],
    }


def _alias_mapping(tmp_path: Path) -> dict[str, Path]:
    return {alias: _synthetic_pptx(tmp_path / f"fixture-{index}.pptx") for index, alias in enumerate(ALIASES, 1)}


def test_repository_wide_scan_blocks_ordinary_tracked_source_canary(tmp_path: Path):
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    (tmp_path / "ordinary.py").write_text("value = 'D:/SYNTHETIC_PRIVATE_ROOT/secret.pptx'", encoding="utf-8")
    subprocess.run(["git", "add", "ordinary.py"], cwd=tmp_path, check=True)
    findings = RepositoryPrivacyScanner(private_root_signatures=["SYNTHETIC_PRIVATE_ROOT"]).scan_repository(tmp_path)
    assert {item.classification for item in findings} >= {"absolute_path"}
    assert all("SYNTHETIC_PRIVATE_ROOT" not in str(item) for item in findings)


@pytest.mark.parametrize("value", ["D:/SYNTHETIC_ROOT/file.pptx", "D:\\\\SYNTHETIC_ROOT\\file.pptx", "\\\\server\\share\\file.pptx", "/mnt/d/SYNTHETIC_ROOT/file.pptx"])
def test_repository_scan_detects_absolute_path_forms_without_configured_signature(tmp_path: Path, value: str):
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    (tmp_path / "ordinary.py").write_text(f"value = {value!r}", encoding="utf-8")
    subprocess.run(["git", "add", "ordinary.py"], cwd=tmp_path, check=True)
    findings = RepositoryPrivacyScanner().scan_repository(tmp_path)
    assert any(item.classification == "absolute_path" for item in findings)


def _legacy_basename_tokens() -> list[str]:
    text = subprocess.check_output(["git", "show", f"HEAD:{LEGACY_EXCEPTION_PATH}"], cwd=ROOT, text=True, encoding="utf-8")
    section = text.split("## D3-2 — Private exemplar roles remain asymmetric", 1)[1].split("\n## ", 1)[0]
    return sorted(set(__import__("re").findall(r"[A-Za-z0-9_ -]+\.pptx", section, __import__("re").I)))


def test_exact_reviewed_legacy_occurrence_is_recorded_not_silently_suppressed():
    tokens = _legacy_basename_tokens()
    findings, exceptions = RepositoryPrivacyScanner(forbidden_basenames=tokens).scan_repository_with_legacy_exception(ROOT, forbidden_basenames=tokens)
    assert not [item for item in findings if item.location == LEGACY_EXCEPTION_PATH]
    assert exceptions == [{"exception_id": "CP2-PRE-1-LEGACY-D3-2", "repository_relative_path": LEGACY_EXCEPTION_PATH, "reviewed_blob_sha": LEGACY_EXCEPTION_BLOB_SHA, "privacy_rule_id": "forbidden_private_basename", "status": "applied_legacy_exception"}]
    assert all(token not in str(exception) for token in tokens for exception in exceptions)


def test_changed_legacy_blob_invalidates_exception(tmp_path: Path):
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    (tmp_path / "thesis-deck-system" / "reviews").mkdir(parents=True)
    target = tmp_path / LEGACY_EXCEPTION_PATH
    target.write_text("## D3-2 — Private exemplar roles remain asymmetric\nlegacy.pptx\n", encoding="utf-8")
    subprocess.run(["git", "add", "."], cwd=tmp_path, check=True)
    subprocess.run(["git", "-c", "user.email=test@example.invalid", "-c", "user.name=test", "commit", "-qm", "fixture"], cwd=tmp_path, check=True)
    target.write_text("## D3-2 — Private exemplar roles remain asymmetric\nlegacy.pptx\nchanged\n", encoding="utf-8")
    findings, exceptions = RepositoryPrivacyScanner(forbidden_basenames=["legacy.pptx"]).scan_repository_with_legacy_exception(tmp_path, forbidden_basenames=["legacy.pptx"])
    assert not exceptions
    assert any(item.classification == "forbidden_private_basename" for item in findings)


def test_legacy_file_with_absolute_path_or_extra_basename_is_not_exempted(tmp_path: Path):
    subprocess.run(["git", "init", "-q"], cwd=tmp_path, check=True)
    (tmp_path / "thesis-deck-system" / "reviews").mkdir(parents=True)
    target = tmp_path / LEGACY_EXCEPTION_PATH
    target.write_text("## D3-2 — Private exemplar roles remain asymmetric\nlegacy.pptx\nD:/PRIVATE_ROOT_CANARY/file.pptx\n", encoding="utf-8")
    subprocess.run(["git", "add", "."], cwd=tmp_path, check=True)
    subprocess.run(["git", "-c", "user.email=test@example.invalid", "-c", "user.name=test", "commit", "-qm", "fixture"], cwd=tmp_path, check=True)
    findings, exceptions = RepositoryPrivacyScanner(private_root_signatures=["PRIVATE_ROOT_CANARY"], forbidden_basenames=["legacy.pptx", "extra.pptx"]).scan_repository_with_legacy_exception(tmp_path, forbidden_basenames=["legacy.pptx", "extra.pptx"])
    assert not exceptions
    assert findings


@pytest.mark.parametrize("kind", ["synthetic_measurement", "synthetic_observation", "simulation_output"])
def test_production_observation_rejects_synthetic_and_simulation_evidence(kind: str):
    registry = SchemaRegistry(SCHEMAS, include_phase3=True)
    catalog = canonical_observation_catalogs(registry, [_evidence("E001", kind)], [_plot("E001")])
    assert "P3-OBSERVATION-PRODUCTION-EMPIRICAL-POLICY" in validate_observation_visual_binding(
        _binding(), catalog=catalog, evidence_policy="production"
    )


def test_fixture_mode_preserves_synthetic_observation_test_support():
    registry = SchemaRegistry(SCHEMAS, include_phase3=True)
    catalog = canonical_observation_catalogs(registry, [_evidence("E001", "synthetic_measurement", verified="synthetic_test_only")], [_plot("E001")])
    assert validate_observation_visual_binding(_binding(), catalog=catalog, evidence_policy="fixture") == []


def test_alias_resolver_accepts_only_three_stable_aliases_and_local_mapping(tmp_path: Path):
    resolver = LocalPrivateAliasResolver(_alias_mapping(tmp_path), private_root=tmp_path / "raw")
    assert resolver.resolve(ALIASES[0]).alias_uri == ALIASES[0]
    with pytest.raises(Checkpoint2PolicyViolation):
        resolver.resolve("D:/arbitrary.pptx")
    with pytest.raises(Checkpoint2PolicyViolation):
        resolver.resolve("private://unrecognized")


def test_source_session_is_recorded_before_open_and_derives_hash_slide_count(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    resolver = LocalPrivateAliasResolver(_alias_mapping(tmp_path), private_root=tmp_path / "raw", execution=run)
    session = resolver.resolve(ALIASES[0]).open_read_only()
    profile = session.profile_structurally("shell")
    assert run.evidence.authorized_source_sessions == 0
    assert run.evidence.source_sessions[ALIASES[0]]["profiling_status"] == "pass"
    assert run.evidence.source_sessions[ALIASES[0]]["sanitizer_handoff"] == "pending"
    assert profile["source_sha256"] == hashlib.sha256((tmp_path / "fixture-1.pptx").read_bytes()).hexdigest()
    assert profile["slide_count"] == 1


def test_malformed_package_is_blocked_before_structural_profile(tmp_path: Path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not an ooxml zip")
    resolver = LocalPrivateAliasResolver({ALIASES[0]: bad, ALIASES[1]: bad, ALIASES[2]: bad}, private_root=tmp_path / "raw")
    with pytest.raises(Checkpoint2PolicyViolation):
        resolver.resolve(ALIASES[0]).open_read_only()


def test_profiler_is_read_only_and_structural_without_render(tmp_path: Path):
    mapping = _alias_mapping(tmp_path)
    before = hashlib.sha256(mapping[ALIASES[1]].read_bytes()).hexdigest()
    profile = LocalPrivateAliasResolver(mapping, private_root=tmp_path / "raw").resolve(ALIASES[1]).open_read_only().profile_structurally("body")
    after = hashlib.sha256(mapping[ALIASES[1]].read_bytes()).hexdigest()
    assert before == after
    assert profile["slide_count"] == 1
    assert profile["render_count"] == 0
    assert all("text" not in str(value).casefold() for value in profile.get("forbidden_exports", {}).values())


def test_sanitizers_fail_closed_and_block_exemplar_two_shell_contamination():
    shell = {"alias_uri": ALIASES[0], "source_sha256": SHA, "profile_id": "SHELL001", "slide_size": {"width": 13.333, "height": 7.5}, "master_count": 1, "layout_count": 1, "shell_primitives": [], "slide_count": 1}
    assert sanitize_shell_descriptor(shell)["profile_id"] == "SHELL001"
    with pytest.raises(Checkpoint2PolicyViolation):
        sanitize_shell_descriptor({**shell, "private_text": "forbidden"})
    body = {"alias_uri": ALIASES[1], "source_sha256": SHA, "profile_id": "BODY001", "slide_size": {"width": 13.333, "height": 7.5}, "slide_count": 1, "candidate_families": [], "body_measurements": []}
    assert sanitize_body_descriptor(body)["profile_id"] == "BODY001"
    with pytest.raises(Checkpoint2PolicyViolation):
        sanitize_body_descriptor({**body, "shell_primitives": []})


def test_private_unapproved_provider_creates_no_render_and_blocks_review(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    assert run.private_render_review({"image_capable": True, "approved_for_private_exemplars": False}) == "blocked_visual_review"
    assert run.evidence.private_renders_created == 0


def test_approved_synthetic_provider_records_render_delete_lifecycle(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    assert run.private_render_review({"image_capable": True, "approved_for_private_exemplars": True, "private_content_allowed": True, "hash_binding_supported": True, "egress_mode": "local_only", "retention_class": "ephemeral", "supported_input_forms": ["local_private_handle"]}) == "blocked_visual_review"
    assert (run.evidence.private_renders_created, run.evidence.private_renders_deleted, run.evidence.private_renders_retained) == (0, 0, 0)


def test_checkpoint_two_qa_rejects_literal_status_or_retained_render(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    run.evidence.record_pre_open_gate("CP2-PRE-1", "pass")
    run.evidence.record_pre_open_gate("CP2-PRE-2", "pass")
    run.evidence.private_renders_created = 1
    run.evidence.private_renders_retained = 1
    record = run.qa_record()
    assert record["aggregate_status"] == "fail"
    record["aggregate_status"] = "pass"
    assert "CP2-QA-AGGREGATE-NONDERIVED" in validate_checkpoint2_qa(record)


def test_checkpoint_two_build_produces_only_sanitized_descriptors(tmp_path: Path):
    artifact_root = tmp_path / "artifacts"
    qa = build_checkpoint2(repository_root=ROOT, local_aliases=_alias_mapping(tmp_path), private_root=tmp_path / "raw", artifact_root=artifact_root)
    assert qa["aggregate_status"] == "pass"
    assert {path.name for path in artifact_root.iterdir()} == {
        "sanitized-exemplar-manifest.json", "sanitized-shell-structural-descriptors.json",
        "sanitized-body-structural-descriptors.json", "checkpoint-2-qa.json",
    }


def test_shell_profile_contains_measured_topology_regions_and_basis(tmp_path: Path):
    profile = LocalPrivateAliasResolver(_alias_mapping(tmp_path), private_root=tmp_path / "raw").resolve(ALIASES[0]).open_read_only().profile_structurally("shell")
    assert profile["measurement_basis"]
    assert profile["layout_master_topology"]
    assert profile["slide_layout_topology"]
    assert profile["shell_regions"]
    assert profile["safe_content_bounds"]["basis"] in {"measured", "derived", "not_observable_structurally"}
    assert profile["safe_content_bounds"]["source_scope"] in {"slide_master", "slide_layout", "not_observable_structurally"}
    assert profile["safe_content_bounds"]["value"] is None or profile["safe_content_bounds"]["value"] != {"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0, "basis": "derived"}


def test_body_profile_contains_structural_metrics_and_classification(tmp_path: Path):
    profile = LocalPrivateAliasResolver(_alias_mapping(tmp_path), private_root=tmp_path / "raw").resolve(ALIASES[1]).open_read_only().profile_structurally("body")
    measurement = profile["body_measurements"][0]
    assert measurement["objects"]
    assert "panel_count" in measurement["metrics"] and "comparison_symmetry" in measurement["metrics"]
    assert "measurement_basis" in measurement
    assert profile["candidate_families"][0]["confidence"] in {"structurally_supported", "provisional", "insufficient_structural_evidence"}


def test_nested_sanitizer_rejects_unknown_nested_keys_and_free_text():
    shell = {"alias_uri": ALIASES[0], "source_sha256": SHA, "profile_id": "SHELL001", "slide_size": {"width": 13.333, "height": 7.5}, "master_count": 1, "layout_count": 1, "shell_primitives": [{"primitive_id": "S001", "primitive_class": "header", "geometry": {"x": 0.1, "y": 0.1, "w": 0.2, "h": 0.1, "basis": "measured"}, "unexpected": "private words"}], "slide_count": 1}
    with pytest.raises(Checkpoint2PolicyViolation):
        sanitize_shell_descriptor(shell)


def test_source_session_lifecycle_records_started_and_failed(tmp_path: Path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a package")
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    resolver = LocalPrivateAliasResolver({ALIASES[0]: bad}, private_root=tmp_path / "raw", execution=run)
    with pytest.raises(Checkpoint2PolicyViolation):
        resolver.resolve(ALIASES[0]).open_read_only()
    session = run.evidence.source_sessions[ALIASES[0]]
    assert session["started"] is True
    assert session["closed"] is True
    assert session["outcome"] == "failed"
    assert run.evidence.failed_source_sessions == 1


def test_capability_only_provider_cannot_fabricate_private_review(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    status = run.private_render_review({"image_capable": True, "approved_for_private_exemplars": True, "private_content_allowed": True, "hash_binding_supported": True, "egress_mode": "local_only", "retention_class": "ephemeral", "supported_input_forms": ["local_private_handle"]})
    assert status in {"blocked_visual_review", "not_run"}
    assert (run.evidence.private_renders_created, run.evidence.private_renders_deleted, run.evidence.private_renders_retained) == (0, 0, 0)


def test_qa_requires_descriptor_quality_owning_checks():
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=Path("."))
    run.evidence.private_root_status = "pass"
    run.evidence.privacy_scan_status = "pass"
    run.evidence.source_sessions = {alias: {"closed": True, "outcome": "success"} for alias in ALIASES}
    record = run.qa_record()
    assert record["aggregate_status"] == "fail"
    assert "descriptor_quality_checks" in record["execution_evidence"]


def test_qa_rejects_inconsistent_session_counters():
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=Path("."))
    run.evidence.private_root_status = "pass"
    run.evidence.privacy_scan_status = "pass"
    run.evidence.source_sessions = {alias: {"started": True, "closed": True, "outcome": "success"} for alias in ALIASES}
    run.evidence.descriptor_quality_checks = [{"check_id": f"CP2-DQ-{i}", "status": "pass"} for i in range(7)]
    record = run.qa_record()
    record["execution_evidence"]["successful_closed_sessions"] = 0
    record["execution_evidence_sha256"] = hashlib.sha256(__import__("json").dumps(record["execution_evidence"], sort_keys=True, separators=(",", ":")).encode()).hexdigest()
    record["aggregate_status"] = "pass"
    assert "CP2-QA-AGGREGATE-NONDERIVED" in validate_checkpoint2_qa(record)


def test_group_transform_composes_translation_scaling_and_nested_coordinates():
    parent = {"off": [100.0, 200.0], "ext": [400.0, 300.0], "ch_off": [10.0, 20.0], "ch_ext": [200.0, 100.0], "flip_h": False, "flip_v": False}
    child = {"off": [60.0, 45.0], "ext": [40.0, 20.0], "ch_off": [0.0, 0.0], "ch_ext": [1.0, 1.0], "flip_h": False, "flip_v": False}
    absolute = _compose_transform(parent, child)
    assert absolute == pytest.approx((200.0, 275.0, 80.0, 60.0), abs=1e-6)


@pytest.mark.parametrize(
    ("flip_h", "flip_v", "head", "tail", "expected"),
    [
        (False, False, "triangle", "none", ("directed", [0.1, 0.2], [0.4, 0.6])),
        (False, False, "none", "triangle", ("directed", [0.4, 0.6], [0.1, 0.2])),
        (True, False, "triangle", "none", ("directed", [0.4, 0.2], [0.1, 0.6])),
        (False, False, "none", "none", ("plain", [0.1, 0.2], [0.4, 0.6])),
    ],
)
def test_connector_semantics_uses_arrowheads_and_flip_not_bbox_order(flip_h, flip_v, head, tail, expected):
    value = _connector_semantics(0.1, 0.2, 0.3, 0.4, flip_h=flip_h, flip_v=flip_v, head_arrow=head, tail_arrow=tail)
    assert (value["directedness"], value["start"], value["end"]) == expected


def test_metric_observation_never_serializes_unmeasured_placeholder_zero():
    unavailable = _metric_observation(None, basis="not_observable_structurally", evidence_ids=[])
    assert unavailable == {"value": None, "basis": "not_observable_structurally", "evidence_state": "unavailable", "supporting_object_ids": []}
    derived = _metric_observation(0.5, basis="derived", evidence_ids=["O001", "O002"])
    assert derived["value"] == 0.5
    assert derived["evidence_state"] == "derived"


def test_family_classifier_requires_specific_signature_not_counts_alone():
    unrelated = _classify_structural_family(
        objects=[{"object_id": f"O{i:03d}", "object_class": "picture", "geometry": {"x": i * 0.2, "y": 0.1, "w": 0.1, "h": 0.1}} for i in range(4)],
        connectors=[], metrics={}, groups=[],
    )
    assert unrelated["confidence"] != "structurally_supported"
    flowchart = _classify_structural_family(
        objects=[{"object_id": "O001", "object_class": "native_shape", "geometry": {"x": 0.1, "y": 0.1, "w": 0.1, "h": 0.1}}],
        connectors=[{"object_id": f"O{i:03d}", "directedness": "directed"} for i in range(2, 5)], metrics={}, groups=[],
    )
    assert flowchart["family"] != "fishbone_research_map" or flowchart["confidence"] != "structurally_supported"


def test_shell_measurements_are_master_layout_scoped_and_expose_placeholders(tmp_path: Path):
    profile = LocalPrivateAliasResolver(_alias_mapping(tmp_path), private_root=tmp_path / "raw").resolve(ALIASES[0]).open_read_only().profile_structurally("shell")
    assert all(item["source_scope"] in {"slide_master", "slide_layout", "theme", "not_observable_structurally"} for item in profile["shell_primitives"])
    assert all(item["source_scope"] in {"slide_master", "slide_layout", "theme", "not_observable_structurally"} for item in profile["style_roles"])
    assert all(item["source_scope"] in {"slide_master", "slide_layout", "theme", "not_observable_structurally"} for item in profile["typography_roles"])
    assert "placeholder_measurements" in profile


def test_metric_observation_rejects_arbitrary_relation_text():
    body = {"alias_uri": ALIASES[1], "source_sha256": SHA, "profile_id": "BODY001", "slide_size": {"width": 13.333, "height": 7.5, "basis": "measured"}, "slide_count": 1, "candidate_families": [{"family": "other_insufficient_structural_evidence", "confidence": "insufficient_structural_evidence", "evidence_basis": []}], "body_measurements": [{"slide_id": "SL001", "measurement_basis": "measured", "objects": [], "connectors": [], "groups": [], "panels": [], "metrics": {key: {"value": None, "basis": "not_observable_structurally", "evidence_state": "unavailable", "supporting_object_ids": []} for key in ("text_area_ratio", "figure_area_ratio", "dominant_figure_ratio", "figure_text_ratio", "annotation_density", "whitespace_fraction", "comparison_symmetry", "matrix_rows", "matrix_columns", "panel_count", "caption_candidate_count", "callout_candidate_count", "photo_schematic_relation")}, "style_roles": []}]}
    body["body_measurements"][0]["metrics"]["photo_schematic_relation"] = {"value": "private free text", "basis": "measured", "evidence_state": "measured", "supporting_object_ids": []}
    with pytest.raises(Checkpoint2PolicyViolation):
        sanitize_body_descriptor(body)


def test_scheme_color_is_retained_as_theme_role_and_unknown_is_not_none():
    from xml.etree import ElementTree as ET
    node = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:schemeClr val="accent1"/></a:solidFill></a:sp>')
    assert _color_role(node, True) == "theme:accent1"
    unknown = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:schemeClr val="mystery"/></a:solidFill></a:sp>')
    assert _color_role(unknown, True) == "unknown"


def test_shell_regions_report_distinct_container_recurrence_without_overlap():
    from thesis_deck_system.phase3_checkpoint2 import _shell_regions
    slides = [
        {"source_scope": "slide_layout", "source_container_id": "L001", "objects": [
            {"object_class": "text", "geometry": {"x": .1, "y": .05, "w": .3, "h": .08}, "source_scope": "slide_layout", "source_container_id": "L001", "placeholder_type": "title"},
            {"object_class": "text", "geometry": {"x": .1, "y": .05, "w": .3, "h": .08}, "source_scope": "slide_layout", "source_container_id": "L001", "placeholder_type": "header"},
        ]},
        {"source_scope": "slide_layout", "source_container_id": "L002", "objects": [
            {"object_class": "text", "geometry": {"x": .1, "y": .05, "w": .3, "h": .08}, "source_scope": "slide_layout", "source_container_id": "L002", "placeholder_type": "title"},
        ]},
    ]
    regions = _shell_regions(slides, 100, 100)
    by_role = {item["role"]: item for item in regions}
    support = by_role["title"]["support_by_scope"][0]
    assert support["occurrence_count"] == 2
    assert support["source_container_count"] == 2
    assert support["eligible_container_count"] == 2
    assert support["coverage_ratio"] == pytest.approx(1.0)
    assert "L001" in support["supporting_source_ids"]
    assert "header" not in by_role or by_role["header"]["occurrence_count"] == 0


def test_shell_region_placeholder_semantics_override_geometry_fallback():
    from thesis_deck_system.phase3_checkpoint2 import _shell_regions
    slides = [{"source_scope": "slide_master", "source_container_id": "M001", "objects": [
        {"object_class": "text", "geometry": {"x": .1, "y": .03, "w": .3, "h": .05}, "source_scope": "slide_master", "source_container_id": "M001", "placeholder_type": "title"}
    ]}]
    regions = _shell_regions(slides, 100, 100)
    assert {item["role"] for item in regions} == {"title"}
    assert regions[0]["role_evidence"] == "placeholder_semantic"


def test_color_evidence_preserves_distinct_rgb_and_theme_palette():
    from thesis_deck_system.phase3_checkpoint2 import _color_evidence
    from xml.etree import ElementTree as ET
    a = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:srgbClr val="112233"/></a:solidFill></a:sp>')
    b = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:srgbClr val="445566"/></a:solidFill></a:sp>')
    theme = {"accent1": "AABBCC"}
    assert _color_evidence(a, True, theme)["direct_rgb"] != _color_evidence(b, True, theme)["direct_rgb"]
    themed = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:schemeClr val="accent1"/></a:solidFill></a:sp>')
    evidence = _color_evidence(themed, True, theme)
    assert evidence["theme_token"] == "accent1"
    assert evidence["resolved_rgb"] == "AABBCC"


def test_unsupported_color_transform_is_explicit():
    from thesis_deck_system.phase3_checkpoint2 import _color_evidence
    from xml.etree import ElementTree as ET
    node = ET.fromstring('<a:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:srgbClr val="112233"><a:gamma/></a:srgbClr></a:solidFill></a:sp>')
    assert _color_evidence(node, True, {})["transform_status"] == "unsupported"


def test_safe_non_default_font_survives_font_policy():
    from thesis_deck_system.phase3_checkpoint2 import _font_family
    assert _font_family("Yu Gothic") == "Yu Gothic"


def test_rotated_geometry_is_explicitly_unsupported():
    from thesis_deck_system.phase3_checkpoint2 import _rotation_state
    assert _rotation_state({"rot": 600000}) == {"rotation_status": "unsupported", "rotation_deg": 1.0, "geometry_eligible": False}
    assert _rotation_state({"rot": 0})["geometry_eligible"] is True


def test_rotated_shape_cannot_be_measured_for_family_confidence():
    family = _classify_structural_family(
        objects=[{"object_id": "O001", "object_class": "picture", "geometry": {"x": .1, "y": .1, "w": .2, "h": .2}, "geometry_eligible": False}],
        connectors=[], metrics={}, groups=[]
    )
    assert family["confidence"] == "insufficient_structural_evidence"


def test_body_descriptor_preserves_typography_observation_without_text():
    from thesis_deck_system.phase3_checkpoint2 import sanitize_body_descriptor
    body = {
        "alias_uri": ALIASES[1], "source_sha256": SHA, "profile_id": "P3-BODY-001",
        "slide_size": {"width": 13.333, "height": 7.5, "basis": "measured"}, "slide_count": 1,
        "candidate_families": [{"family": "other_insufficient_structural_evidence", "confidence": "insufficient_structural_evidence", "evidence_basis": ["insufficient"]}],
        "body_measurements": [{"slide_id": "SL001", "measurement_basis": "measured", "objects": [], "connectors": [], "groups": [], "panels": [], "metrics": {key: {"value": None, "basis": "not_observable_structurally", "evidence_state": "unavailable", "supporting_object_ids": []} for key in ("text_area_ratio", "figure_area_ratio", "dominant_figure_ratio", "figure_text_ratio", "annotation_density", "whitespace_fraction", "comparison_symmetry", "matrix_rows", "matrix_columns", "panel_count", "caption_candidate_count", "callout_candidate_count", "photo_schematic_relation")}, "style_roles": [], "typography_observations": [{"observation_id": "T001", "role": "annotation", "role_confidence": "provisional", "family": "Yu Gothic", "theme_font_role": None, "script_role": "latin", "font_evidence_state": "explicit_font", "size_pt": 11.0, "weight": "regular", "style": "normal", "basis": "measured", "source_scope": "slide_body", "supporting_object_id": "O001"}]}],
        "theme_profiles": [], "slide_theme_topology": []
    }
    result = sanitize_body_descriptor(body)
    assert result["body_measurements"][0]["typography_observations"][0]["family"] == "Yu Gothic"


def test_date_time_placeholder_never_becomes_navigation_and_keeps_semantic_role():
    from thesis_deck_system.phase3_checkpoint2 import _shell_regions
    slides = [{"source_scope": "slide_layout", "source_container_id": "L001", "objects": [
        {"object_class": "text", "geometry": {"x": .1, "y": .9, "w": .2, "h": .04}, "source_scope": "slide_layout", "source_container_id": "L001", "placeholder_type": "dt"}
    ]}]
    regions = _shell_regions(slides, 100, 100)
    assert {region["role"] for region in regions} == {"date_time"}
    assert "navigation" not in {region["role"] for region in regions}


def test_shell_region_support_is_scope_aware_and_coverage_is_not_cross_scope():
    from thesis_deck_system.phase3_checkpoint2 import _shell_regions
    slides = [
        {"source_scope": "slide_master", "source_container_id": "M001", "objects": [
            {"object_class": "text", "geometry": {"x": .1, "y": .02, "w": .3, "h": .05}, "source_scope": "slide_master", "source_container_id": "M001", "placeholder_type": "title"}
        ]},
        {"source_scope": "slide_layout", "source_container_id": "L001", "objects": [
            {"object_class": "text", "geometry": {"x": .1, "y": .02, "w": .3, "h": .05}, "source_scope": "slide_layout", "source_container_id": "L001", "placeholder_type": "title"}
        ]},
        {"source_scope": "slide_layout", "source_container_id": "L002", "objects": []},
    ]
    title = next(region for region in _shell_regions(slides, 100, 100) if region["role"] == "title")
    by_scope = {item["source_scope"]: item for item in title["support_by_scope"]}
    assert by_scope["slide_master"]["coverage_ratio"] == pytest.approx(1.0)
    assert by_scope["slide_layout"]["coverage_ratio"] == pytest.approx(0.5)
    assert title["support_by_scope"] == sorted(title["support_by_scope"], key=lambda value: value["source_scope"])


def test_navigation_requires_explicit_independent_role_evidence():
    from thesis_deck_system.phase3_checkpoint2 import _shell_regions
    slides = [{"source_scope": "slide_layout", "source_container_id": "L001", "objects": [
        {"object_class": "text", "geometry": {"x": .1, "y": .9, "w": .2, "h": .04}, "source_scope": "slide_layout", "source_container_id": "L001", "placeholder_type": "dt"}
    ]}]
    assert "navigation" not in {region["role"] for region in _shell_regions(slides, 100, 100)}


def test_theme_color_resolution_is_bound_to_the_observations_master_theme():
    from thesis_deck_system.phase3_checkpoint2 import resolve_theme_color
    themes = {
        "T001": {"palette": {"accent1": "AAAAAA"}},
        "T002": {"palette": {"accent1": "BBBBBB"}},
    }
    master_theme = {"M001": "T001", "M002": "T002"}
    assert resolve_theme_color("accent1", master_id="M001", master_theme=master_theme, theme_profiles=themes) == "AAAAAA"
    assert resolve_theme_color("accent1", master_id="M002", master_theme=master_theme, theme_profiles=themes) == "BBBBBB"


@pytest.mark.parametrize(
    ("token", "expected_role", "expected_script"),
    [("+mj-lt", "major", "latin"), ("+mn-lt", "minor", "latin"), ("+mj-ea", "major", "east_asian"), ("+mn-ea", "minor", "east_asian"), ("+mj-cs", "major", "complex_script")],
)
def test_theme_font_tokens_are_typed_as_role_and_script(token, expected_role, expected_script):
    from thesis_deck_system.phase3_checkpoint2 import _theme_font_token
    assert _theme_font_token(token) == (expected_role, expected_script)


def test_east_asian_unicode_font_is_safe_but_private_like_font_is_rejected():
    from thesis_deck_system.phase3_checkpoint2 import _font_family
    assert _font_family("微軟正黑體") == "微軟正黑體"
    assert _font_family("D:/private-font") == "unknown"
    assert _font_family("https://unsafe-font") == "unknown"


def _typography_slide(font_nodes: str) -> object:
    """Synthetic text shape with deliberately multi-script run properties."""
    from xml.etree import ElementTree as ET

    return ET.fromstring(
        """<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:cSld><p:spTree><p:sp><p:nvSpPr/>
            <p:spPr><a:xfrm><a:off x="100" y="300"/><a:ext cx="600" cy="100"/></a:xfrm></p:spPr>
            <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr sz="1200">"""
        + font_nodes
        + """</a:rPr><a:t>synthetic</a:t></a:r></a:p></p:txBody>
          </p:sp></p:spTree></p:cSld></p:sld>"""
    )


def test_one_run_with_latin_and_east_asian_fonts_emits_two_script_observations():
    from thesis_deck_system.phase3_checkpoint2 import ReadOnlyPrivateSourceSession

    profile = ReadOnlyPrivateSourceSession._slide_profile(
        _typography_slide('<a:latin typeface="Arial"/><a:ea typeface="微軟正黑體"/>'),
        1000,
        1000,
        1,
        source_scope="slide_body",
    )
    observations = profile["typography_observations"]
    assert [(item["script_role"], item["family"]) for item in observations] == [
        ("latin", "Arial"),
        ("east_asian", "微軟正黑體"),
    ]
    assert len({item["observation_id"] for item in observations}) == 2
    assert {item["supporting_object_id"] for item in observations} == {"O001"}


def test_one_run_with_all_three_script_fonts_emits_three_observations():
    from thesis_deck_system.phase3_checkpoint2 import ReadOnlyPrivateSourceSession

    profile = ReadOnlyPrivateSourceSession._slide_profile(
        _typography_slide('<a:latin typeface="Arial"/><a:ea typeface="微軟正黑體"/><a:cs typeface="Calibri"/>'),
        1000,
        1000,
        1,
        source_scope="slide_body",
    )
    assert [item["script_role"] for item in profile["typography_observations"]] == [
        "latin",
        "east_asian",
        "complex_script",
    ]


def test_theme_font_tokens_resolve_independently_for_each_script():
    from thesis_deck_system.phase3_checkpoint2 import ReadOnlyPrivateSourceSession

    profile = ReadOnlyPrivateSourceSession._slide_profile(
        _typography_slide('<a:latin typeface="+mn-lt"/><a:ea typeface="+mj-ea"/><a:cs typeface="+mn-cs"/>'),
        1000,
        1000,
        1,
        source_scope="slide_body",
        theme_font_scheme={
            "major": {"east_asian": "微軟正黑體"},
            "minor": {"latin": "Arial", "complex_script": "Calibri"},
        },
    )
    assert [(item["script_role"], item["theme_font_role"], item["family"], item["font_evidence_state"]) for item in profile["typography_observations"]] == [
        ("latin", "minor", "Arial", "theme_font_resolved"),
        ("east_asian", "major", "微軟正黑體", "theme_font_resolved"),
        ("complex_script", "minor", "Calibri", "theme_font_resolved"),
    ]


def test_supplemental_theme_cjk_font_mappings_are_sanitized_as_controlled_metadata():
    from thesis_deck_system.phase3_checkpoint2 import _sanitize_theme_profile

    result = _sanitize_theme_profile({
        "theme_profile_id": "T001", "usage_state": "unreferenced", "authority_state": "reference_only", "supporting_master_ids": [], "supporting_slide_ids": [], "palette": [],
        "font_scheme": {"major": {}, "minor": {}},
        "supplemental_fonts": [
            {"theme_font_role": "major", "script_code": code, "family": "微軟正黑體", "authority_state": "reference_only"}
            for code in ("Hans", "Hant", "Jpan", "Hang")
        ],
    })
    assert [item["script_code"] for item in result["supplemental_fonts"]] == ["Hang", "Hans", "Hant", "Jpan"]


@pytest.mark.parametrize(
    "supplemental",
    [
        {"theme_font_role": "major", "script_code": "Hant", "family": "D:/unsafe", "authority_state": "reference_only"},
        {"theme_font_role": "major", "script_code": "UnknownScript", "family": "Arial", "authority_state": "reference_only"},
    ],
)
def test_unsafe_or_unapproved_supplemental_theme_font_is_rejected(supplemental: dict):
    from thesis_deck_system.phase3_checkpoint2 import _sanitize_theme_profile

    with pytest.raises(Checkpoint2PolicyViolation):
        _sanitize_theme_profile({"theme_profile_id": "T001", "usage_state": "unreferenced", "authority_state": "reference_only", "supporting_master_ids": [], "supporting_slide_ids": [], "palette": [], "font_scheme": {"major": {}, "minor": {}}, "supplemental_fonts": [supplemental]})


def test_supplemental_mapping_does_not_fabricate_east_asian_run_resolution():
    from thesis_deck_system.phase3_checkpoint2 import ReadOnlyPrivateSourceSession

    profile = ReadOnlyPrivateSourceSession._slide_profile(
        _typography_slide(''), 1000, 1000, 1, source_scope="slide_body",
        theme_font_scheme={"major": {}, "minor": {}, "supplemental_fonts": {"Hant": "微軟正黑體"}},
    )
    assert [(item["script_role"], item["font_evidence_state"], item["family"]) for item in profile["typography_observations"]] == [
        ("unspecified", "inherited_unresolved", "unknown"),
    ]


def test_descriptor_local_theme_ids_do_not_collide_across_descriptors():
    from thesis_deck_system.phase3_checkpoint2 import resolve_descriptor_theme_color

    descriptors = {
        "P3-TEMPLATE-PRIMARY-1": {"T001": {"palette": {"accent1": "AAAAAA"}}},
        "P3-TEMPLATE-PRIMARY-3": {"T001": {"palette": {"accent1": "BBBBBB"}}},
    }
    assert resolve_descriptor_theme_color("accent1", profile_id="P3-TEMPLATE-PRIMARY-1", theme_profile_id="T001", descriptor_theme_profiles=descriptors) == "AAAAAA"
    assert resolve_descriptor_theme_color("accent1", profile_id="P3-TEMPLATE-PRIMARY-3", theme_profile_id="T001", descriptor_theme_profiles=descriptors) == "BBBBBB"


def test_typography_resolution_counts_reconcile_each_persisted_observation():
    from thesis_deck_system.phase3_checkpoint2 import typography_resolution_counts

    observations = [
        {"script_role": "latin", "font_evidence_state": "explicit_font"},
        {"script_role": "east_asian", "font_evidence_state": "theme_font_resolved"},
        {"script_role": "complex_script", "font_evidence_state": "inherited_unresolved"},
    ]
    counts = typography_resolution_counts(observations)
    assert counts["latin"]["explicit_font"] == 1
    assert counts["east_asian"]["theme_font_resolved"] == 1
    assert counts["complex_script"]["inherited_unresolved"] == 1
    assert sum(sum(states.values()) for states in counts.values()) == len(observations)


def test_absent_script_nodes_emit_unspecified_inherited_typography_not_latin():
    from thesis_deck_system.phase3_checkpoint2 import ReadOnlyPrivateSourceSession

    profile = ReadOnlyPrivateSourceSession._slide_profile(
        _typography_slide(""), 1000, 1000, 1, source_scope="slide_body"
    )
    assert [(item["script_role"], item["font_evidence_state"], item["family"])
            for item in profile["typography_observations"]] == [
        ("unspecified", "inherited_unresolved", "unknown"),
    ]


def test_typography_counts_reconcile_unspecified_observations_separately():
    from thesis_deck_system.phase3_checkpoint2 import typography_resolution_counts

    counts = typography_resolution_counts([
        {"script_role": "latin", "font_evidence_state": "explicit_font"},
        {"script_role": "unspecified", "font_evidence_state": "inherited_unresolved"},
    ])
    assert counts["latin"]["explicit_font"] == 1
    assert counts["unspecified"]["inherited_unresolved"] == 1
    assert sum(sum(states.values()) for states in counts.values()) == 2


def test_theme_usage_is_derived_from_topology_and_orphans_are_reference_only():
    from thesis_deck_system.phase3_checkpoint2 import classify_theme_usage

    profiles = [
        {"theme_profile_id": "T001", "palette": {}, "font_scheme": {"major": {}, "minor": {}}, "supplemental_fonts": []},
        {"theme_profile_id": "T002", "palette": {}, "font_scheme": {"major": {}, "minor": {}}, "supplemental_fonts": [{"theme_font_role": "major", "script_code": "Hant", "family": "微軟正黑體"}]},
    ]
    result = classify_theme_usage(
        profiles,
        master_theme_topology=[{"source_id": "M001", "target_id": "T001", "basis": "measured", "source_scope": "slide_master"}],
        slide_theme_topology=[],
    )
    assert result[0]["usage_state"] == "referenced"
    assert result[0]["authority_state"] == "active_professor_style"
    assert result[0]["supporting_master_ids"] == ["M001"]
    assert result[1]["usage_state"] == "unreferenced"
    assert result[1]["authority_state"] == "reference_only"
    assert result[1]["supplemental_fonts"][0]["authority_state"] == "reference_only"


def test_unknown_theme_topology_target_fails_closed():
    from thesis_deck_system.phase3_checkpoint2 import classify_theme_usage

    with pytest.raises(Checkpoint2PolicyViolation):
        classify_theme_usage(
            [{"theme_profile_id": "T001", "palette": {}, "font_scheme": {"major": {}, "minor": {}}, "supplemental_fonts": []}],
            master_theme_topology=[{"source_id": "M001", "target_id": "T999", "basis": "measured", "source_scope": "slide_master"}],
            slide_theme_topology=[],
        )


def test_active_theme_lookup_excludes_unreferenced_palette_and_supplemental_fonts():
    from thesis_deck_system.phase3_checkpoint2 import resolve_active_descriptor_theme_color

    descriptors = {
        "P3-TEMPLATE-PRIMARY-1": {
            "T001": {"usage_state": "referenced", "palette": {"accent1": "AAAAAA"}},
            "T002": {"usage_state": "unreferenced", "palette": {"accent1": "BBBBBB"}},
        }
    }
    assert resolve_active_descriptor_theme_color("accent1", profile_id="P3-TEMPLATE-PRIMARY-1", theme_profile_id="T001", descriptor_theme_profiles=descriptors) == "AAAAAA"
    assert resolve_active_descriptor_theme_color("accent1", profile_id="P3-TEMPLATE-PRIMARY-1", theme_profile_id="T002", descriptor_theme_profiles=descriptors) is None


def test_unknown_typography_cannot_satisfy_font_fidelity_owning_gate(tmp_path: Path):
    run = Checkpoint2Run.start(pre_open_passed=True, private_root=tmp_path / "raw")
    shell = {"alias_uri": ALIASES[0], "source_sha256": SHA, "profile_id": "P3-SHELL-001", "slide_size": {"width": 13.333, "height": 7.5, "basis": "measured"}, "master_count": 1, "layout_count": 1, "slide_count": 1, "measurement_basis": {key: "measured" for key in ("slide_size", "topology", "regions", "typography", "styles", "primitives", "placeholders")}, "layout_master_topology": [], "slide_layout_topology": [], "shell_regions": [], "safe_content_bounds": {"value": None, "basis": "not_observable_structurally", "source_scope": "not_observable_structurally", "evidence_ids": []}, "typography_roles": [], "style_roles": [], "shell_primitives": [], "placeholder_measurements": [], "theme_palette": []}
    body = {"alias_uri": ALIASES[1], "source_sha256": SHA, "profile_id": "P3-BODY-001", "slide_size": {"width": 13.333, "height": 7.5, "basis": "measured"}, "slide_count": 1, "candidate_families": [{"family": "other_insufficient_structural_evidence", "confidence": "insufficient_structural_evidence", "evidence_basis": ["insufficient"]}], "body_measurements": [{"slide_id": "SL001", "measurement_basis": "measured", "objects": [], "connectors": [], "groups": [], "panels": [], "metrics": {key: {"value": None, "basis": "not_observable_structurally", "evidence_state": "unavailable", "supporting_object_ids": []} for key in ("text_area_ratio", "figure_area_ratio", "dominant_figure_ratio", "figure_text_ratio", "annotation_density", "whitespace_fraction", "comparison_symmetry", "matrix_rows", "matrix_columns", "panel_count", "caption_candidate_count", "callout_candidate_count", "photo_schematic_relation")}, "style_roles": [], "typography_observations": [{"observation_id": "T001", "role": "unknown", "role_confidence": "unknown", "family": "unknown", "theme_font_role": None, "script_role": "latin", "font_evidence_state": "unknown", "size_pt": 11.0, "weight": "regular", "style": "normal", "basis": "measured", "source_scope": "slide_body", "supporting_object_id": "O001"}]}]}
    registry = SchemaRegistry(SCHEMAS, include_phase3=True)
    run.set_descriptor_quality([shell, {**shell, "alias_uri": ALIASES[2], "profile_id": "P3-SHELL-003"}], body, registry)
    assert next(item for item in run.evidence.descriptor_quality_checks if item["check_id"] == "CP2-DQ-FONT-FIDELITY")["status"] == "fail"
